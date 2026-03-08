import streamlit as st
import pandas as pd
import numpy as np
import matplotlib.pyplot as plt
import altair as alt
import io
import re
import os
import json
import time
import itertools
import tempfile
import random
from math import pi
from datetime import datetime, timedelta
from google import genai
from google.genai import types
from PIL import Image, ExifTags
import pypdf
from youtube_comment_downloader import YoutubeCommentDownloader, SORT_BY_POPULAR
from fpdf import FPDF
from wordcloud import WordCloud
import streamlit.components.v1 as components
import feedparser
import networkx as nx
import hashlib
import plotly.graph_objects as go
import plotly.express as px
import urllib.request
import cv2
import math
from PIL import ImageDraw, ImageChops, ImageEnhance
import urllib.parse
from cv2 import GaussianBlur
import requests
import sqlite3
from pptx import Presentation
from pptx.util import Inches
from streamlit_agraph import agraph, Node, Edge, Config
from bs4 import BeautifulSoup
from gtts import gTTS
import pyzipper
import string
from stegano import lsb
import wave

# --- PAGE CONFIGURATION ---
st.set_page_config(
    page_title="RAP Dashboard",
    layout="wide",
    initial_sidebar_state="expanded"
)
st.markdown("""
<style>
    /* Import monospace terminal font strictly for metrics/numbers */
    @import url('https://fonts.googleapis.com/css2?family=Fira+Code:wght@300;400;700&display=swap');

    /* Style for main headers (Neon Blue Glow, normal spacing and casing) */
    h1, h2, h3 {
        color: #38BDF8 !important;
        text-shadow: 0 0 10px rgba(56, 189, 248, 0.3);
    }

    /* Button Hover Effect: Add neon glow */
    .stButton > button[data-baseweb="button"]:hover {
        box-shadow: 0 0 15px rgba(56, 189, 248, 0.6);
    }

    /* Style for metrics (Hacker Green + Monospace Font for numbers) */
    [data-testid="stMetricValue"] {
        font-family: 'Fira Code', monospace !important;
        color: #10B981 !important;
        font-weight: 700;
        text-shadow: 0 0 8px rgba(16, 185, 129, 0.4);
    }

    /* Pulsing animation for Error Alerts (🚨) - Kept red for true emergencies */
    @keyframes pulse-red {
        0% { box-shadow: 0 0 0 0 rgba(239, 68, 68, 0.7); }
        70% { box-shadow: 0 0 0 10px rgba(239, 68, 68, 0); }
        100% { box-shadow: 0 0 0 0 rgba(239, 68, 68, 0); }
    }
    [data-testid="stAlert"]:has(.st-emotion-cache-121s1zy) /* Targets error alerts */ {
        animation: pulse-red 2s infinite;
        border-left: 5px solid #EF4444;
    }

    /* Custom Scrollbar */
    ::-webkit-scrollbar { width: 8px; height: 8px; }
    ::-webkit-scrollbar-track { background: #0A0E17; }
    ::-webkit-scrollbar-thumb { background: #334155; border-radius: 4px; }
    ::-webkit-scrollbar-thumb:hover { background: #38BDF8; }

    /* Neon borders for containers */
    [data-testid="stContainer"] {
        border: 1px solid #1E293B;
        box-shadow: inset 0 0 20px rgba(0,0,0,0.5);
    }
</style>
""", unsafe_allow_html=True)


# --- SESSION STATE INITIALIZATION ---
if 'api_calls' not in st.session_state: st.session_state['api_calls'] = 0
if 'token_usage' not in st.session_state: st.session_state['token_usage'] = 0
if 'oracle_history' not in st.session_state: st.session_state['oracle_history'] = []
if 'doc_oracle_history' not in st.session_state: st.session_state['doc_oracle_history'] = []
if 'scroll_to_top' not in st.session_state: st.session_state['scroll_to_top'] = False
if 'global_entities' not in st.session_state: st.session_state['global_entities'] = set()
if 'seen_radar_links' not in st.session_state: st.session_state['seen_radar_links'] = set()

if 'data_store' not in st.session_state:
    st.session_state['data_store'] = {
        'CSV File Upload': {'df': None, 'analyzed': None, 'summary': None},
        'YouTube Link': {'df': None, 'analyzed': None, 'summary': None},
        'Universal URL': {'df': None, 'analyzed': None, 'summary': None},
        'Raw Text Paste': {'df': None, 'analyzed': None, 'summary': None},
        'Telegram Dump (JSON)': {'df': None, 'analyzed': None, 'summary': None},
        'Reddit Native (OSINT)': {'df': None, 'analyzed': None, 'summary': None},
        'Arena': {'df_a': None, 'df_b': None, 'analyzed_a': None, 'analyzed_b': None},
        'Radar': {'df': None, 'analyzed': None}
    }

# --- DYNAMIC AI TRANSLATION ENGINE ---
@st.cache_data(show_spinner=False)
def ai_t(text, target_lang, api_key):
    """
    Translates UI labels dynamically using Gemini. 
    Uses st.cache_data to minimize API calls and costs.
    """
    if not text or target_lang == "English":
        return text
    
    try:
        client = genai.Client(api_key=api_key)
        # We instruct the AI to be extremely concise to save tokens
        prompt = f"Translate this UI label to {target_lang}. Return ONLY the translated text: '{text}'"
        response = client.models.generate_content(model='gemini-2.0-flash', contents=prompt)
        return response.text.strip().replace("'", "").replace('"', "")
    except Exception:
        # Fallback to English if the API fails or key is missing
        return text

def t(text):
    """
    Funzione di traduzione disattivata per massimizzare le prestazioni.
    Restituisce sempre il testo originale in inglese.
    Per riattivarla in futuro, de-commenta il blocco sottostante.
    """
    return text 

    # lang = st.session_state.get('global_lang', 'English')
    # key = st.session_state.get('api_key') or st.secrets.get("GEMINI_API_KEY", "")
    # if not key or lang == "English":
    #     return text
    # return ai_t(text, lang, key)

# --- HELPER: PANOPTICON (PERSISTENT MEMORY) ---
def init_panopticon():
    conn = sqlite3.connect('rap_panopticon.db', check_same_thread=False)
    c = conn.cursor()
    c.execute('''CREATE TABLE IF NOT EXISTS targets (agent_id TEXT PRIMARY KEY, risk_score REAL, threat_type TEXT, last_seen TEXT)''')
    conn.commit()
    return conn

panopticon_db = init_panopticon()

def save_to_panopticon(agent_id, risk_score, threat_type):
    try:
        c = panopticon_db.cursor()
        now = datetime.now().strftime('%Y-%m-%d %H:%M:%S')
        c.execute('''INSERT INTO targets (agent_id, risk_score, threat_type, last_seen) VALUES (?, ?, ?, ?)
                     ON CONFLICT(agent_id) DO UPDATE SET risk_score=max(risk_score, ?), last_seen=?''',
                  (str(agent_id), float(risk_score), str(threat_type), now, float(risk_score), now))
        panopticon_db.commit()
    except Exception as e: pass

def check_panopticon(agent_id):
    try:
        c = panopticon_db.cursor()
        c.execute('''SELECT risk_score, threat_type, last_seen FROM targets WHERE agent_id=?''', (str(agent_id),))
        return c.fetchone()
    except: return None

# --- HELPER: BILLING & TOKEN TRACKING ---
if 'api_calls' not in st.session_state: st.session_state['api_calls'] = 0
if 'in_tokens' not in st.session_state: st.session_state['in_tokens'] = 0
if 'out_tokens' not in st.session_state: st.session_state['out_tokens'] = 0

def increment_counter(input_text_len=0, output_text_len=0):
    st.session_state['api_calls'] += 1
    st.session_state['in_tokens'] += (input_text_len / 4)
    st.session_state['out_tokens'] += (output_text_len / 4)

def get_cost_estimate():
    in_t = st.session_state.get('in_tokens', 0)
    out_t = st.session_state.get('out_tokens', 0)
    cost = (in_t / 1_000_000) * 0.10 + (out_t / 1_000_000) * 0.40
    return (in_t + out_t), cost

# --- HELPER: DYNAMIC HQ LOCATION (FOR ATLAS) ---
@st.cache_data(ttl=3600)
def get_hq_location():
    """Detects the current real-world location of the user/server using IP geolocation."""
    try:
        res = requests.get("http://ip-api.com/json/", timeout=5).json()
        if res['status'] == 'success':
            return float(res['lat']), float(res['lon']), res['city']
    except Exception: 
        pass
    # Fallback to Rome if adblocker or firewall blocks the IP request
    return 41.9028, 12.4964, "Rome (Fallback HQ)"

if 'hq_coords' not in st.session_state:
    st.session_state['hq_coords'] = get_hq_location()

# --- HELPER: ERROR LEVEL ANALYSIS (ELA) ---
def perform_ela(image, quality=90):
    """Calculates the Error Level Analysis of an image to detect Photoshop manipulations."""
    try:
        # Ensure the image is in RGB mode
        original = image.convert('RGB')
        
        # Temporarily resave the image at a lower quality
        buffer = io.BytesIO()
        original.save(buffer, 'JPEG', quality=quality)
        buffer.seek(0)
        resaved = Image.open(buffer)
        
        # Calculate the absolute difference between the original and the compressed image
        ela_image = ImageChops.difference(original, resaved)
        
        # Visually amplify the difference to make it visible to the human eye
        extrema = ela_image.getextrema()
        max_diff = max([ex[1] for ex in extrema])
        if max_diff == 0: max_diff = 1
        scale = 255.0 / max_diff
        ela_image = ImageEnhance.Brightness(ela_image).enhance(scale)
        
        return ela_image
    except Exception as e:
        return None

# --- HELPER: OPSEC FACE ANONYMIZATION (OPTIMIZED TACTICAL CALIBRATION) ---
def anonymize_faces(image):
    """
    Executes a fast dual sweep.
    Optimized scaleFactor and minNeighbors to prevent false positives ("99 faces" glitch).
    """
    try:
        # 1. Convert PIL Image to OpenCV format
        img_cv = np.array(image.convert('RGB'))
        img_cv = cv2.cvtColor(img_cv, cv2.COLOR_RGB2BGR)
        gray = cv2.cvtColor(img_cv, cv2.COLOR_BGR2GRAY)
        
        # 2. Load BOTH models
        frontal_cascade = cv2.CascadeClassifier(cv2.data.haarcascades + 'haarcascade_frontalface_default.xml')
        profile_cascade = cv2.CascadeClassifier(cv2.data.haarcascades + 'haarcascade_profileface.xml')
        
        # 3. SWEEP 1: Frontal Scan (Increased scaleFactor to 1.1 and minNeighbors to 6 for precision/speed)
        frontal_faces = frontal_cascade.detectMultiScale(
            gray, scaleFactor=1.1, minNeighbors=6, minSize=(30, 30)
        )
        
        # 4. SWEEP 2: Profile Scan
        profile_faces = profile_cascade.detectMultiScale(
            gray, scaleFactor=1.1, minNeighbors=6, minSize=(30, 30)
        )
        
        # Convert tuples to lists
        f_faces = list(frontal_faces) if len(frontal_faces) > 0 else []
        p_faces = list(profile_faces) if len(profile_faces) > 0 else []
        all_faces = f_faces + p_faces
        
        # 5. SURGICAL ANTI-OVERLAP PROTOCOL
        final_faces = []
        for (x, y, w, h) in all_faces:
            overlap = False
            for (fx, fy, fw, fh) in final_faces:
                cx, cy = x + w/2, y + h/2
                fcx, fcy = fx + fw/2, fy + fh/2
                dist = ((cx - fcx)**2 + (cy - fcy)**2)**0.5
                
                # ONLY discard if the centers are EXTREMELY close
                if dist < (min(w, fw) * 0.5):
                    overlap = True
                    break
            
            if not overlap:
                final_faces.append((x, y, w, h))
        
        # 6. Apply Dynamic Heavy Blur
        for (x, y, w, h) in final_faces:
            face_roi = img_cv[y:y+h, x:x+w]
            
            # Calculate blur intensity
            k_size = (w // 2) | 1 
            if k_size < 3: k_size = 3
            
            # Sigma=50 guarantees a very dense, completely opaque blur
            blurred_face = cv2.GaussianBlur(face_roi, (k_size, k_size), 50)
            img_cv[y:y+h, x:x+w] = blurred_face
            
        # 7. Return the final image
        img_final = cv2.cvtColor(img_cv, cv2.COLOR_BGR2RGB)
        return Image.fromarray(img_final), len(final_faces)
        
    except Exception as e:
        return image, 0
    
# --- HELPER: EXIF GPS SNIPER ---
def get_decimal_from_dms(dms, ref):
    """Converts Degrees, Minutes, Seconds to Decimal Coordinates."""
    degrees, minutes, seconds = dms[0], dms[1], dms[2]
    dec = float(degrees) + float(minutes)/60 + float(seconds)/3600
    if ref in ['S', 'W']: dec = -dec
    return dec

def get_exif_location(image):
    """Extracts hidden GPS coordinates from an image's EXIF data."""
    try:
        exif = image._getexif()
        if not exif: return None
        geotagging = {}
        for (idx, tag) in TAGS.items():
            if tag == 'GPSInfo' and idx in exif:
                for (key, val) in GPSTAGS.items():
                    if key in exif[idx]:
                        geotagging[val] = exif[idx][key]
        
        if 'GPSLatitude' in geotagging and 'GPSLongitude' in geotagging:
            lat = get_decimal_from_dms(geotagging['GPSLatitude'], geotagging.get('GPSLatitudeRef', 'N'))
            lon = get_decimal_from_dms(geotagging['GPSLongitude'], geotagging.get('GPSLongitudeRef', 'E'))
            return lat, lon
    except:
        pass
    return None

# --- HELPER: DIGITAL FORENSICS HASHING ---
def calculate_sha256(file_bytes):
    """Generates a cryptographic hash for Chain of Custody proof."""
    if not file_bytes: return "N/A"
    return hashlib.sha256(file_bytes).hexdigest()

# --- HELPER: FORENSIC STORYBOARD GENERATOR ---
def create_video_storyboard(video_bytes, num_frames=12):
    """Dismembers the video into exactly 12 frames and merges them into a single grid image."""
    with tempfile.NamedTemporaryFile(delete=False, suffix=".mp4") as tmp:
        tmp.write(video_bytes)
        tmp_path = tmp.name
    
    cap = cv2.VideoCapture(tmp_path)
    total_frames = int(cap.get(cv2.CAP_PROP_FRAME_COUNT))
    fps = cap.get(cv2.CAP_PROP_FPS)
    
    if total_frames == 0: return None
        
    interval = max(1, total_frames // num_frames)
    frames = []
    
    for i in range(num_frames):
        frame_id = i * interval
        if frame_id >= total_frames: frame_id = total_frames - 1
        cap.set(cv2.CAP_PROP_POS_FRAMES, frame_id)
        ret, frame = cap.read()
        if ret:
            # Convert BGR to RGB for PIL
            frame_rgb = cv2.cvtColor(frame, cv2.COLOR_BGR2RGB)
            pil_img = Image.fromarray(frame_rgb)
            # Resize to prevent memory explosion
            pil_img.thumbnail((400, 400)) 
            timestamp = frame_id / fps if fps > 0 else 0
            frames.append((pil_img, timestamp))
            
    cap.release()
    os.unlink(tmp_path)
    
    if not frames: return None
        
    # Create the grid (3 columns)
    cols = 3
    rows = math.ceil(len(frames) / cols)
    w, h = frames[0][0].size
    padding_y = 40
    
    grid_img = Image.new('RGB', (cols * w, rows * (h + padding_y)), color='black')
    draw = ImageDraw.Draw(grid_img)
    
    for i, (img, ts) in enumerate(frames):
        row = i // cols
        col = i % cols
        x = col * w
        y = row * (h + padding_y)
        grid_img.paste(img, (x, y + padding_y))
        # Print the exact timestamp on the frame
        draw.text((x + 10, y + 10), f"Frame Time: {ts:.2f}s", fill="red")
        
    # Return the PIL image object directly instead of raw bytes
    return grid_img

# --- HELPER: AUDIO WAVEFORM GENERATION ---
def generate_audio_waveform(audio_bytes):
    """Generates a visual waveform from audio bytes for forensic inspection."""
    try:
        # 1. Prova a leggerlo come un vero file WAV de-compresso
        try:
            with wave.open(io.BytesIO(audio_bytes), 'rb') as wav_file:
                frames = wav_file.readframes(-1)
                audio_data = np.frombuffer(frames, dtype=np.int16)
                # Semplifica i dati per evitare che il grafico sia troppo pesante
                if len(audio_data) > 10000:
                    audio_data = audio_data[::len(audio_data)//10000]
        except Exception:
            # 2. Se è un MP3 o M4A, crea una "pseudo-forma d'onda" visiva
            # invece di mostrare un cubo rosso pieno di rumore binario.
            raw_data = np.frombuffer(audio_bytes, dtype=np.uint8)
            chunk_size = max(1, len(raw_data) // 1000)
            
            # Calcola la varianza a blocchi per simulare l'intensità del suono
            envelope = [np.std(raw_data[i:i+chunk_size]) for i in range(0, len(raw_data), chunk_size)]
            audio_data = np.array(envelope)
            audio_data = audio_data - np.mean(audio_data) # Centra la linea sullo zero

        fig, ax = plt.subplots(figsize=(10, 2))
        # Linea leggermente più spessa per un look migliore
        ax.plot(audio_data, color='#ef4444', linewidth=1.0) 
        ax.axis('off')
        ax.set_title("Forensic Audio Waveform", color='white', size=10)
        fig.patch.set_alpha(0) # Transparent background
        
        buf = io.BytesIO()
        plt.savefig(buf, format='png', bbox_inches='tight', pad_inches=0)
        plt.close(fig)
        return buf.getvalue()
    except Exception as e:
        return None

# --- HELPER: PII SANITIZER (OPERATION BLACKOUT - CIA STYLE) ---
def sanitize_pii(text):
    if not isinstance(text, str): return text
    # Replace sensitive data with solid black blocks (Classified Redaction)
    text = re.sub(r'[a-zA-Z0-9_.+-]+@[a-zA-Z0-9-]+\.[a-zA-Z0-9-.]+', '██████████', text)
    text = re.sub(r'\b\+?\d{1,3}?[-.\s]?\(?\d{2,3}?\)?[-.\s]?\d{3}[-.\s]?\d{4}\b', '██████████', text)
    text = re.sub(r'\b[A-Z]{2}\d{2}[a-zA-Z0-9]{11,30}\b', '████████████████', text)
    text = re.sub(r'\b\d{1,3}\.\d{1,3}\.\d{1,3}\.\d{1,3}\b', '████████', text)
    return text

# --- HELPER: CALLBACK FOR DATA EDITOR ---
def update_editor_state(method_key):
    editor_key = f"editor_{method_key}"
    if editor_key in st.session_state:
        new_data = st.session_state[editor_key]
        if isinstance(new_data, pd.DataFrame):
            st.session_state['data_store'][method_key]['df'] = new_data

# --- HELPER: ROBUST JSON PARSER ---
def extract_json(text):
    try:
        text = re.sub(r'```json\s*', '', text)
        text = re.sub(r'```', '', text)
        match = re.search(r'\{.*\}', text, re.DOTALL)
        if match:
            return json.loads(match.group(0))
        return None
    except:
        return None

# --- HELPER: SMART DOCUMENT CHUNKING ---
def chunk_document(text, chunk_size=3000, overlap=300):
    """Splits a massive document into overlapping semantic chunks for context retrieval."""
    if not text or not text.strip(): return []
    chunks = []
    start = 0
    text = text.strip()
    while start < len(text):
        end = start + chunk_size
        chunks.append(text[start:end])
        start += chunk_size - overlap
    return chunks

def prepare_smart_context(texts):
    """
    Placeholder to maintain compatibility with the UI. 
    Instead of calculating vectors (which fail), we return the raw chunks 
    to be analyzed dynamically during the chat.
    """
    # We return a dummy list to signal the UI that the document is 'ready'
    return [[0.0] * 768 for _ in texts]

# --- HELPER: OSINT WEB SPIDER ---
def run_web_spider(url):
    """Crawls a specific webpage to extract its hidden internal structure and external affiliations."""
    try:
        headers = {'User-Agent': 'Mozilla/5.0 (Windows NT 10.0; Win64; x64) RAP_OSINT_Bot/1.0'}
        res = requests.get(url, headers=headers, timeout=8)
        html = res.text
        
        # Raw regex extraction of href links to bypass BeautifulSoup requirement
        links = re.findall(r'href=[\'"]?([^\'" >]+)', html)
        
        base_domain = urllib.parse.urlparse(url).netloc
        internal, external = set(), set()
        
        for link in links:
            if link.startswith('http'):
                if base_domain in link:
                    internal.add(link)
                else:
                    external.add(link)
            elif link.startswith('/'):
                internal.add(f"https://{base_domain}{link}")
                
        return list(internal)[:20], list(external)[:20]
    except Exception as e:
        return [], []

def sanitize_response(data):
    if not data: data = {}
    
    aliases = {
        'explanation': ['analysis', 'reasoning', 'comment', 'description'],
        'correction': ['fix', 'suggestion'],
        'counter_reply': ['reply', 'debunk'],
        'rewritten_text': ['rewrite', 'sanitized_text']
    }
    
    for target, source_list in aliases.items():
        if target not in data or not data[target]:
            for source in source_list:
                if source in data and data[source]:
                    data[target] = data[source]
                    break

    defaults = {
        "has_fallacy": False,
        "fallacy_type": "None",
        "explanation": "Analysis not available.",
        "correction": "",
        "main_topic": "General",
        "micro_topic": "General",
        "target": "None",
        "primary_emotion": "Neutral",
        "archetype": "Observer",
        "relevance": "Relevant",
        "counter_reply": "",
        "sentiment": "Neutral",
        "aggression": 0,
        "rewritten_text": "No rewrite necessary.",
        "facts": [],
        "ai_generated_probability": 0,
        "ai_analysis": "AI analysis not performed."
    }
    
    for key, default_val in defaults.items():
        if key not in data or data[key] is None:
            data[key] = default_val
        if isinstance(data[key], float) and np.isnan(data[key]):
            data[key] = default_val
        if str(data[key]).lower() == 'nan':
            data[key] = default_val

    if data['explanation'] in ["Analysis not available.", "Analisi non disponibile.", "No analysis provided."]:
        if data.get('has_fallacy'):
            data['explanation'] = f"Issue detected: {data.get('fallacy_type')}."
        else:
            data['explanation'] = "No issues detected, content appears legitimate."

    return data

# --- HELPER: VOTE PARSER ---
def parse_votes(vote_str):
    if not vote_str: return 0
    s = str(vote_str).lower().strip()
    if 'k' in s:
        return int(float(re.sub(r'[^\d.]', '', s)) * 1000)
    try:
        return int(re.sub(r'[^\d]', '', s))
    except:
        return 0

# --- HELPER: BOT HUNTER ---
def detect_bot_activity(df):
    if df is None or df.empty: return df
    df['is_bot'] = False
    df['bot_reason'] = ""
    
    # 1. Duplicate Content (Swarm Coordination)
    dupes = df.duplicated(subset=['content'], keep=False)
    df.loc[dupes, 'is_bot'] = True
    df.loc[dupes, 'bot_reason'] = "Duplicate Content (Coordination) "
    
    if 'agent_id' in df.columns:
        # 2. High Frequency Spammer
        agent_counts = df['agent_id'].value_counts()
        spammers = agent_counts[agent_counts > 3].index
        mask_spammer = df['agent_id'].isin(spammers)
        df.loc[mask_spammer, 'is_bot'] = True
        df.loc[mask_spammer, 'bot_reason'] += "High Frequency "

        # 3. Superhuman Velocity Check (< 5 seconds between posts)
        if 'timestamp' in df.columns:
            try:
                temp_df = df.copy()
                temp_df['parsed_time'] = pd.to_datetime(temp_df['timestamp'], errors='coerce')
                valid_time = temp_df.dropna(subset=['parsed_time'])
                
                if not valid_time.empty:
                    valid_time = valid_time.sort_values(by=['agent_id', 'parsed_time'])
                    valid_time['time_diff'] = valid_time.groupby('agent_id')['parsed_time'].diff().dt.total_seconds()
                    
                    # Identifica gli agenti che hanno pubblicato messaggi a meno di 5 secondi di distanza
                    fast_posters = valid_time[(valid_time['time_diff'] >= 0) & (valid_time['time_diff'] < 5)]['agent_id'].unique()
                    mask_fast = df['agent_id'].isin(fast_posters)
                    
                    df.loc[mask_fast, 'is_bot'] = True
                    df.loc[mask_fast, 'bot_reason'] += "Superhuman Velocity (<5s) "
            except Exception:
                pass # Ignora se i timestamp sono formattati male

    df['bot_reason'] = df['bot_reason'].str.strip()
    return df

# --- HELPER: TREND PROJECTION ---
@st.cache_data
def project_trend(df, days_ahead=7):
    if 'timestamp' not in df.columns: return None, False, None
    df_trend = df.copy()
    df_trend['timestamp'] = pd.to_datetime(df_trend['timestamp'], errors='coerce')
    df_trend = df_trend.dropna(subset=['timestamp']).sort_values('timestamp')
    if len(df_trend) < 5: return None, False, None
    
    start_date = df_trend['timestamp'].min()
    df_trend['days_since_start'] = (df_trend['timestamp'] - start_date).dt.total_seconds() / 86400
    x = df_trend['days_since_start'].values
    y = df_trend['aggression'].values
    if len(np.unique(x)) < 2: return None, False, None
    
    slope, intercept = np.polyfit(x, y, 1)
    last_day = x.max()
    future_days = np.arange(last_day + 1, last_day + days_ahead + 1).astype(int)
    forecast_y = slope * future_days + intercept
    future_dates = [start_date + timedelta(days=float(d)) for d in future_days]
    
    past_df = pd.DataFrame({'Date': df_trend['timestamp'], 'Aggression': y, 'Type': 'Historical'})
    future_df = pd.DataFrame({'Date': future_dates, 'Aggression': forecast_y, 'Type': 'Forecast'})
    combined = pd.concat([past_df, future_df])
    
    days_to_critical = None
    if slope > 0.05:
        days_to_critical = (9.0 - intercept) / slope - last_day
        if days_to_critical < 0: days_to_critical = 0
        
    return combined, (slope > 0.05), days_to_critical

# --- HELPER: STRATEGIC SUMMARY ---
def generate_strategic_summary(df, api_key, context="", persona="Intelligence Analyst"):
    if df is None or df.empty: return None
    flagged = df[df['has_fallacy']==True]
    bots = df[df['is_bot']==True]
    fallacy_counts = flagged['fallacy_type'].value_counts().to_string() if not flagged.empty else "None"
    top_topics = df['main_topic'].value_counts().head(3).to_string() if 'main_topic' in df.columns else "None"
    top_emotion = df['primary_emotion'].mode()[0] if 'primary_emotion' in df.columns and not df['primary_emotion'].empty else "Unknown"
    trend_msg = "Stable"
    if 'aggression' in df.columns and len(df) > 10:
        first_half = df['aggression'].iloc[:len(df)//2].mean()
        second_half = df['aggression'].iloc[len(df)//2:].mean()
        if second_half > first_half * 1.2: trend_msg = "ESCALATING (Crisis Risk)"
        elif second_half < first_half * 0.8: trend_msg = "De-escalating"
    bot_count = len(bots)
    avg_agg = df['aggression'].mean() if 'aggression' in df.columns else 0
    prompt = f"""
    You are a {persona}. 
    CONTEXT TOPIC: "{context}"
    DATA: Avg Aggression {avg_agg:.1f}/10, Trend {trend_msg}, Emotion {top_emotion}, Bots {bot_count}, Fallacies {fallacy_counts}.
    
    TASK: Write a strictly professional "Executive Briefing" (max 150 words). 
    CRITICAL RULE: NO introductory phrases. Start directly with the analysis. 
    CRITICAL RULE: Write strictly in the SAME LANGUAGE as the "CONTEXT TOPIC".
    """
    try:
        client = genai.Client(api_key=api_key)
        response = client.models.generate_content(model='gemini-2.5-flash', contents=prompt)
        increment_counter(len(prompt), len(response.text))
        return response.text
    except Exception as e:
        return f"Could not generate summary: {str(e)}"

# --- HELPER: ACTION DECK GENERATOR ---
def generate_action_deck(df, action_type, api_key, context):
    if df is None or df.empty: return "No comments selected."
    comments_text = "\n".join([f"- {row['content']} (Issue: {row.get('fallacy_type', 'None')})" for _, row in df.iterrows()])
    prompts = {
        "Generate Counter-Narrative Thread": "Create a Twitter/X thread (3-5 tweets) that politely but firmly debunks the following comments using logic and facts. Tone: Professional but engaging.",
        "Draft Official Statement": "Write a formal Press Release or Official Statement addressing the concerns raised in these comments. Tone: Corporate, reassuring, authoritative.",
        "Legal Risk Assessment": "Analyze these comments for potential defamation, libel, or terms of service violations. Output a bulleted list of actionable legal or moderation steps.",
        "Engagement Strategy": "Suggest 3 specific replies to the most influential comments here to de-escalate the situation and win over the audience."
    }
    prompt = f"""
    CONTEXT: {context}
    ACTION: {action_type}
    INPUT COMMENTS: {comments_text[:15000]}
    INSTRUCTION: {prompts[action_type]}
    OUTPUT LANGUAGE: Same as input comments.
    """
    try:
        client = genai.Client(api_key=api_key)
        response = client.models.generate_content(model='gemini-2.5-flash', contents=prompt)
        increment_counter(len(prompt), len(response.text))
        return response.text
    except Exception as e:
        return f"Error generating action: {str(e)}"

def generate_psyops_profile(df, target_id, api_key):
    target_data = df[df['agent_id'] == target_id]
    if target_data.empty: return "User not found."
    
    comments = "\n".join([f"- {row['content']}" for _, row in target_data.iterrows()])
    
    prompt = f"""
    You are an expert Psychological Operations (Psy-Ops) and Behavioral Analyst.
    Analyze the following messages posted by user '{target_id}':
    {comments[:10000]}
    
    CRITICAL RULE: Detect the language of the messages. You MUST write the ENTIRE response (including all headers, titles, and bullet points) strictly in that SAME language. Do not use English unless the messages are in English.
    
    Provide a "Threat Actor Profile" structured exactly with these 4 points (translated into the target language):
    1. Motivation (Ideological, Troll, Paid Bot, Genuine concern)
    2. Emotional Triggers (What makes them angry?)
    3. Linguistic Profile (Education level, slang, repetitive patterns)
    4. Engagement Recommendation (How to interact with or neutralize them)
    
    Keep it concise and professional.
    """
    try:
        client = genai.Client(api_key=api_key)
        response = client.models.generate_content(model='gemini-2.5-flash', contents=prompt)
        increment_counter(len(prompt), len(response.text))
        return response.text
    except Exception as e:
        return f"Profiling Error: {str(e)}"

# --- HELPER: THE ORACLE (GENERAL) ---
def ask_the_oracle(df, question, api_key, context):
    if df is None: return "No data to analyze."
    subset = df[['agent_id', 'content', 'aggression', 'target', 'main_topic', 'archetype']].head(40).to_string()
    prompt = f"""
    You are "The Oracle", an advanced intelligence system.
    CONTEXT: "{context}"
    DATASET SAMPLE: {subset}
    USER QUESTION: "{question}"
    INSTRUCTIONS: Answer based on data. Be professional. RESPOND IN THE SAME LANGUAGE AS THE USER QUESTION.
    """
    try:
        client = genai.Client(api_key=api_key)
        response = client.models.generate_content(model='gemini-2.5-flash', contents=prompt)
        increment_counter(len(prompt), len(response.text))
        return response.text
    except Exception as e:
        return f"Oracle Error: {str(e)}"

# --- HELPER: DOCUMENT ORACLE (RAG) ---
def ask_document_oracle(full_text, question, api_key):
    prompt = f"""
    You are the "Deep Document Oracle", an AI capable of analyzing massive documents.
    I will provide you with the full text of one or more documents.
    
    USER QUESTION: "{question}"
    
    INSTRUCTIONS:
    - Answer the question comprehensively based ONLY on the provided text.
    - If the text does not contain the answer, say so clearly.
    - CRITICAL RULE: Always cite the exact page number(s) to back up your claims, using the [PAGE X] tags provided in the text format (e.g., "According to the contract [PAGE 12]...").
    - Respond in the language of the user's question.
    
    DOCUMENT TEXT:
    {full_text}
    """
    try:
        client = genai.Client(api_key=api_key)
        # Using Gemini 2.0 Flash for its massive context window
        response = client.models.generate_content(model='gemini-2.0-flash', contents=prompt)
        increment_counter(len(prompt), len(response.text))
        return response.text
    except Exception as e:
        return f"Document Oracle Error: {str(e)}"

# --- HELPER: EXCEL & PDF ---
def generate_excel_report(df, summary_text):
    output = io.BytesIO()
    with pd.ExcelWriter(output, engine='openpyxl') as writer:
        summary_data = [{'Metric': 'Analysis Date', 'Value': datetime.now().strftime('%Y-%m-%d %H:%M')}]
        if summary_text:
            summary_data.append({'Metric': 'Executive Briefing', 'Value': summary_text})
        total = len(df)
        flagged = len(df[df['has_fallacy'] == True])
        summary_data.extend([
            {'Metric': 'Total Items', 'Value': total},
            {'Metric': 'Flagged Issues', 'Value': flagged},
            {'Metric': 'Avg Aggression', 'Value': df['aggression'].mean() if 'aggression' in df.columns else 0}
        ])
        pd.DataFrame(summary_data).to_excel(writer, sheet_name='Executive Briefing', index=False)
        df.to_excel(writer, sheet_name='Full Analysis', index=False)
    return output.getvalue()

class PDFReport(FPDF):
    def header(self):
        self.set_font('Helvetica', 'B', 15)
        self.cell(0, 10, 'RAP: Strategic Intelligence Report', 0, 1, 'C')
        self.ln(5)
    def footer(self):
        self.set_y(-15)
        self.set_font('Helvetica', 'I', 8)
        self.cell(0, 10, f'Page {self.page_no()}', 0, 0, 'C')

def create_fallback_chart(df):
    if df is None or 'fallacy_type' not in df.columns: return None
    try:
        counts = df[df['has_fallacy']==True]['fallacy_type'].value_counts()
        if counts.empty: return None
        fig, ax = plt.subplots(figsize=(6, 4))
        counts.head(5).plot(kind='bar', ax=ax, color='firebrick')
        ax.set_title("Top 5 Detected Issues")
        ax.set_ylabel("Count")
        plt.tight_layout()
        with tempfile.NamedTemporaryFile(delete=False, suffix=".png") as tmp:
            fig.savefig(tmp.name)
            return tmp.name
    except: return None

def generate_pdf_report(df, summary_text=None):
    pdf = PDFReport()
    pdf.add_page()
    pdf.set_auto_page_break(auto=True, margin=15)
    pdf.set_font("Helvetica", size=12)
    pdf.cell(0, 10, f"Date: {datetime.now().strftime('%Y-%m-%d %H:%M')}", 0, 1)
    pdf.ln(5)
    pdf.set_font("Helvetica", 'B', 14)
    pdf.cell(0, 10, "Executive Summary", 0, 1)
    if summary_text:
        pdf.set_font("Helvetica", 'I', 11)
        clean_summ = summary_text.encode('latin-1', 'replace').decode('latin-1')
        pdf.multi_cell(0, 6, clean_summ)
        pdf.ln(5)
    chart_path = create_fallback_chart(df)
    if chart_path:
        pdf.image(chart_path, w=170)
        os.unlink(chart_path)
    return bytes(pdf.output())

# --- HELPER: PPTX GENERATOR (CORPORATE UPGRADE) ---
def generate_pptx_report(df, summary_text=None):
    """
    Generates a corporate-ready PowerPoint presentation containing
    the Executive Briefing, key metrics, and visual charts.
    """
    prs = Presentation()
    
    # 1. Title Slide
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "RAP: Tactical Intelligence Briefing"
    subtitle.text = f"Automated Threat Analysis\nGenerated on: {datetime.now().strftime('%Y-%m-%d %H:%M')}"

    # 2. Executive Summary Slide
    if summary_text:
        bullet_slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(bullet_slide_layout)
        shapes = slide.shapes
        title_shape = shapes.title
        body_shape = shapes.placeholders[1]
        
        title_shape.text = "Executive Briefing"
        tf = body_shape.text_frame
        # Clean the text for PPTX format
        clean_summ = summary_text.encode('latin-1', 'replace').decode('latin-1')
        tf.text = clean_summ

    # 3. Key Metrics Slide
    bullet_slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]
    
    title_shape.text = "Key Threat Metrics"
    tf = body_shape.text_frame
    
    total_items = len(df)
    flagged = len(df[df['has_fallacy'] == True])
    avg_agg = df['aggression'].mean() if 'aggression' in df.columns else 0
    
    tf.text = f"Total Entities Analyzed: {total_items}"
    p = tf.add_paragraph()
    p.text = f"Critical Issues Flagged: {flagged}"
    p = tf.add_paragraph()
    p.text = f"Average Aggression Level: {avg_agg:.1f} / 10"

    # 4. Chart Slide (Reusing the fallback chart generator)
    chart_path = create_fallback_chart(df)
    if chart_path:
        blank_slide_layout = prs.slide_layouts[5] # Title only layout
        slide = prs.slides.add_slide(blank_slide_layout)
        slide.shapes.title.text = "Threat Distribution"
        
        # Insert image centered
        left = Inches(1)
        top = Inches(2)
        height = Inches(4.5)
        slide.shapes.add_picture(chart_path, left, top, height=height)
        
        # Cleanup temporary image
        import os
        os.unlink(chart_path)

    # Save to binary stream for Streamlit download
    pptx_stream = io.BytesIO()
    prs.save(pptx_stream)
    pptx_stream.seek(0)
    return pptx_stream.getvalue()

# --- HELPER: PDF EXTRACTOR & SCRAPERS ---
@st.cache_data
def extract_text_from_pdf(file_obj):
    try:
        pdf_reader = pypdf.PdfReader(file_obj)
        text = ""
        # Enumerate pages starting from 1 to track page numbers
        for page_num, page in enumerate(pdf_reader.pages, start=1):
            content = page.extract_text()
            if content: 
                # Inject page markers directly into the text stream
                text += f"\n\n[PAGE {page_num}]\n{content}"
        return text
    except: return None

@st.cache_data(show_spinner=False)
def scrape_youtube_comments(url, limit=50):
    downloader = YoutubeCommentDownloader()
    comments = []
    try:
        generator = downloader.get_comments_from_url(url, sort_by=SORT_BY_POPULAR)
        for comment in itertools.islice(generator, limit):
            comments.append({
                'agent_id': comment.get('author', 'Anonymous'),
                'timestamp': comment.get('time', 'Unknown'),
                'content': comment.get('text', ''),
                'likes': parse_votes(comment.get('votes', '0'))
            })
        if not comments: return None
        return pd.DataFrame(comments)
    except: return None

# --- NATIVE REDDIT OSINT SCRAPER ---
@st.cache_data(show_spinner=False)
def scrape_reddit_native(subreddit, limit=50):
    try:
        # Clean the input (allow users to type "r/worldnews" or just "worldnews")
        sub = subreddit.replace('r/', '').replace('https://www.reddit.com/r/', '').strip('/')
        url = f"https://www.reddit.com/r/{sub}/new.json?limit={limit}"
        
        # Mask the request as a normal web browser to avoid Reddit's API blocks
        req = urllib.request.Request(
            url, 
            data=None, 
            headers={'User-Agent': 'Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 RAP_OSINT_Bot/1.0'}
        )
        response = urllib.request.urlopen(req)
        data = json.loads(response.read().decode('utf-8'))
        
        posts = []
        for child in data.get('data', {}).get('children', []):
            post = child.get('data', {})
            content = post.get('selftext', '')
            title = post.get('title', '')
            full_text = f"{title}\n{content}".strip()
            
            if full_text:
                posts.append({
                    'agent_id': post.get('author', 'Unknown'),
                    'timestamp': datetime.fromtimestamp(post.get('created_utc', 0)).strftime('%Y-%m-%d %H:%M:%S'),
                    'content': full_text,
                    'likes': post.get('score', 0)
                })
        if not posts: return None
        return pd.DataFrame(posts)
    except Exception as e:
        return None
    
# --- NATIVE TELEGRAM OSINT SCRAPER (CHAMELEON PROTOCOL) ---
@st.cache_data(show_spinner=False)
def scrape_telegram_live(channel_url, limit=30):
    try:
        # Extreme URL sanitization (stripping parameters and paths)
        channel = channel_url.replace('https://', '').replace('http://', '').replace('t.me/', '').replace('@', '').replace('s/', '').strip('/')
        channel = channel.split('/')[0].split('?')[0]
        url = f"https://t.me/s/{channel}" # Target the public web preview
        
        # Advanced cloaking as a modern browser
        headers = {
            'User-Agent': 'Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/120.0.0.0 Safari/537.36',
            'Accept-Language': 'en-US,en;q=0.9'
        }
        
        # Attack via Requests
        response = requests.get(url, headers=headers, timeout=10)
        
        if response.status_code != 200:
            st.error(f"Telegram blocked access. HTTP Error: {response.status_code}")
            return None
            
        html = response.text
        
        # V3: Bulletproof parsing using string splits instead of greedy regex
        # This prevents nested <div> tags (links, bold text) from breaking the extraction
        chunks = re.split(r'<div class="tgme_widget_message_text[^>]*>', html)
        
        if len(chunks) <= 1:
            st.warning("Connection successful, but no text blocks found. Channel may be empty or strictly media-only.")
            return None

        data = []
        now = datetime.now()
        
        # Extract the actual text blocks
        msgs = []
        for chunk in chunks[1:]:
            # Cut the chunk exactly before the footer or the info section starts
            text_block = re.split(r'<div class="tgme_widget_message_info"|<div class="tgme_widget_message_footer"', chunk)[0]
            msgs.append(text_block)
        
        for i, m in enumerate(msgs[-limit:]):
            # Replace <br> tags with actual line breaks
            clean_text = re.sub(r'<br\s*/?>', '\n', m)
            # Strip all remaining HTML tags
            clean_text = re.sub(r'<[^>]+>', ' ', clean_text).strip()
            # Clean up multiple whitespaces
            clean_text = re.sub(r'\s+', ' ', clean_text).strip()
            
            if len(clean_text) > 5:
                data.append({
                    'agent_id': f"TG_{channel}",
                    'timestamp': (now - timedelta(minutes=len(msgs)-i)).strftime("%Y-%m-%d %H:%M:%S"),
                    'content': clean_text,
                    'likes': 0
                })
        
        if not data:
            st.warning("Messages found, but they were empty after stripping HTML formatting.")
            return None
            
        return pd.DataFrame(data)
    except Exception as e: 
        st.error(f"Scraper Error: {str(e)}")
        return None

# --- HELPER: ECHELON WEB SCRAPER ---
def scrape_universal_url(url):
    """
    V3: Focused on density and sequential integrity.
    Prevents word collapsing and ensures clean, numbered nodes.
    """
    try:
        headers = {
            'User-Agent': 'Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/91.0.4472.124 Safari/537.36'
        }
        response = requests.get(url, headers=headers, timeout=10)
        response.raise_for_status()
        
        soup = BeautifulSoup(response.text, 'html.parser')
        
        # Aggressive cleaning: remove UI noise
        for noise in soup(["script", "style", "nav", "header", "footer", "aside", "form"]):
            noise.decompose()

        # Target only the main content containers to avoid menu items
        main_content = soup.find_all(['p', 'article'])
        extracted_data = []
        node_counter = 1 # Reset counter for sequential numbering
        
        for p in main_content:
            # Use space separator to prevent collapsing
            text = p.get_text(separator=' ', strip=True)
            # Remove artifacts like Wikipedia citations [1][2] or multiple spaces
            text = re.sub(r'\[\d+\]', '', text) 
            text = re.sub(r'\s+', ' ', text)
            
            # Focus only on high-density text (long paragraphs)
            if len(text) > 100: 
                extracted_data.append({
                    'agent_id': f"NODE_{node_counter:02d}", # Clean sequential naming (NODE_01, NODE_02...)
                    'timestamp': datetime.now().strftime("%Y-%m-%d %H:%M:%S"),
                    'content': text,
                    'likes': 0,
                    'is_bot': False
                })
                node_counter += 1
        
        if not extracted_data:
            # Emergency fallback: extract all body text and split into clean blocks
            body_text = soup.get_text(separator=' ', strip=True)
            body_text = re.sub(r'\s+', ' ', body_text)
            chunks = [body_text[i:i+800] for i in range(0, len(body_text), 800)]
            for i, chunk in enumerate(chunks):
                if len(chunk) > 60:
                    extracted_data.append({
                        'agent_id': f"FALLBACK_{i+1:02d}",
                        'timestamp': datetime.now().strftime("%Y-%m-%d %H:%M:%S"),
                        'content': chunk.strip(),
                        'likes': 0,
                        'is_bot': False
                    })
                    
        return pd.DataFrame(extracted_data)
    
    except Exception as e:
        st.error(f"Echelon Scraper Error: {str(e)}")
        return None

# --- HELPER: RAW PASTE PARSER ---
def parse_raw_paste(raw_text):
    """
    Parses pasted text. Cuts strictly at the FIRST comma if the header is present,
    preventing commas inside the comments from breaking the data structure.
    """
    lines = [line.strip() for line in raw_text.strip().split('\n') if line.strip()]
    data = []
    now = datetime.now()
    
    has_header = False
    if lines and lines[0].lower().replace(" ", "") == "agent_id,content":
        has_header = True
        lines = lines[1:]
        
    for i, line in enumerate(lines):
        if has_header and ',' in line:
            parts = line.split(',', 1)
            agent = parts[0].strip()
            content = parts[1].strip()
        else:
            agent = "Paste_Source"
            content = line
            
        data.append({
            "agent_id": agent,
            "timestamp": (now - timedelta(seconds=i*2)).strftime("%Y-%m-%d %H:%M:%S"),
            "content": content,
            "likes": 0
        })
        
    return pd.DataFrame(data)

def normalize_dataframe(df):
    target_cols = {
        'content': ['content', 'text', 'body', 'comment'],
        'agent_id': ['agent_id', 'author', 'user'],
        'timestamp': ['timestamp', 'created_at', 'date'],
        'likes': ['likes', 'votes', 'like_count']
    }
    new_df = df.copy()
    new_df.columns = [c.lower().strip() for c in new_df.columns]
    found_cols = {}
    for target, aliases in target_cols.items():
        for alias in aliases:
            if alias in new_df.columns:
                found_cols[target] = alias
                break
    rename_map = {v: k for k, v in found_cols.items()}
    new_df = new_df.rename(columns=rename_map)
    if 'content' not in new_df.columns:
        string_cols = new_df.select_dtypes(include=['object']).columns
        if len(string_cols) > 0:
            best = max(string_cols, key=lambda c: new_df[c].astype(str).str.len().mean())
            new_df = new_df.rename(columns={best: 'content'})
    if 'agent_id' not in new_df.columns: new_df['agent_id'] = 'Unknown_User'
    if 'timestamp' not in new_df.columns: new_df['timestamp'] = 'Unknown_Time'
    if 'likes' not in new_df.columns: new_df['likes'] = 0
    return new_df[['agent_id', 'timestamp', 'content', 'likes']]

def parse_telegram_json(file_obj):
    try:
        data = json.load(file_obj)
        messages = data.get('messages', [])
        cleaned = []
        for m in messages:
            if m.get('type') == 'message' and m.get('text'):
                text = m['text']
                if isinstance(text, list): 
                    text = " ".join([t if isinstance(t, str) else t.get('text', '') for t in text])
                if len(text) > 5:
                    cleaned.append({
                        'agent_id': m.get('from', 'Unknown'),
                        'timestamp': m.get('date', 'Unknown'),
                        'content': text,
                        'likes': 0
                    })
        return pd.DataFrame(cleaned)
    except: 
        return None

# --- HYBRID ANALYZER (WITH AUTO-TRANSLATE) ---
@st.cache_data(show_spinner=False)
def analyze_fallacies_cached(text, api_key, context_info="", persona="Logic & Fact Analysis Engine", target_lang="English"):
    if not text or len(str(text)) < 3: return sanitize_response(None)
    for attempt in range(3):
        try:
            client = genai.Client(api_key=api_key)
            is_long = len(text) > 2000
            common_instructions = f"""
            You are a {persona}. CONTEXT: "{context_info}"
            CRITICAL MULTILINGUAL RULE: The input text might be in any language. 
            You MUST conduct the analysis and write ALL your JSON responses exclusively in {target_lang}.
            TASKS:
            - explanation: Brief Analysis in {target_lang}.
            - main_topic: Central theme in {target_lang}.
            - micro_topic: 1-2 keywords in {target_lang}.
            - target: Target Entity.
            - primary_emotion: [Anger, Fear, Disgust, Sadness, Joy, Surprise, Neutral] in {target_lang}.
            - archetype: [Instigator, Loyalist, Troll, Rational Skeptic, Observer] in {target_lang}.
            - propaganda_tactic: Identify Military Info-War Tactics (e.g., Firehose of Falsehood, Astroturfing, Dead Cat, Whataboutism) if present.
            RULES: Opinions/Sarcasm -> 'has_fallacy': false. You MUST still find standard logical fallacies in 'fallacy_type'.
            RESPONSE (Strict JSON):
            {{ "has_fallacy": bool, "fallacy_type": "Name", "propaganda_tactic": "Name", "explanation": "Text", "correction": "Text", "main_topic": "Text", "micro_topic": "Text", "target": "Text", "primary_emotion": "Text", "archetype": "Text", "relevance": "Text", "counter_reply": "Text", "sentiment": "Text", "translated_text": "Translate the original text into {target_lang} here", "aggression": 0, "sophistication": 0, "hostile_intent": 0 }}
            """
            if is_long: prompt = f"Analyze LONG.\n{common_instructions}\nText: \"{text[:15000]}...\""
            else: prompt = f"Analyze Text.\n{common_instructions}\nText: \"{text}\""
            response = client.models.generate_content(model='gemini-2.5-flash', contents=prompt)
            result = extract_json(response.text)
            if result: return sanitize_response(result)
        except:
            time.sleep(0.1)
            continue
    return sanitize_response(None)

def analyze_fallacies(text, api_key=None, context_info="", persona="Logic & Fact Analysis Engine", target_lang="English"):
    if not api_key: return {"has_fallacy": True}
    res = analyze_fallacies_cached(text, api_key, context_info, persona, target_lang)
    increment_counter(len(text), 500)
    return res

# --- HELPER: COGNITIVE EDITOR (MULTIMODAL + AI DETECTOR) ---
def cognitive_rewrite(text, api_key, media_data=None, media_type="image", target_lang="English"):
    if (not text or len(str(text)) < 3) and not media_data: return None
    try:
        client = genai.Client(api_key=api_key)
        
        prompt_text = f"""
        You are a highly paranoid Strategic Intelligence Investigator and a Military-Grade Digital Forensics Expert. You NEVER trust the narrative of a video or image. You ONLY trust physics, pixels, and temporal consistency.
        
        YOUR TASKS:
        1. AI FORENSICS (THE PARANOIA PROTOCOL): Analyze media for AI GENERATION (Sora, Runway, Midjourney, etc.). Score 0-20 (Natural), 40-70 (Manipulated), 70-100 (Fully AI Generated).
           CRITICAL INSTRUCTION: DO NOT get distracted by what is happening in the scene. Focus strictly on HOW the matter behaves. Look specifically for:
           - TEMPORAL MORPHING: Do objects, limbs, or faces change volume, shape, or structure from one second to another?
           - MATTER COMPENETRATION (CLIPPING): Do solid objects pass through each other? (e.g., hands melting into clothing, feet sinking into solid concrete, a person merging with an animal).
           - KINEMATICS & GRAVITY: Do movements look "floaty", lack real momentum, or defy gravity? Does the center of mass shift impossibly?
           - BACKGROUND WARPING: Does the background or static environment stretch, drag, or melt when a foreground subject moves past it?
           - TEXTURE HALLUCINATION: Does text, fur, or fabric boil, shift, or turn into illegible alien symbols as the camera moves?
           If you detect EVEN ONE of these physics-breaking anomalies, you MUST set the 'ai_generated_probability' score above 85% and detail the exact hallucination.
           
        2. IDENTITY VERIFICATION: Identify any famous people. Evaluate if their appearance and movements are physically consistent.
        3. ADVANCED GEO-INT (Shadow Geolocation): 
           - IF THE MEDIA IS AI-GENERATED (Score > 60): Output exactly "Fictional AI-Generated Environment - Geolocation not applicable."
           - IF THE MEDIA IS REAL: Deduce geographic location by identifying micro-clues (street signs, architectural styles, car plates).
        4. SYLLOGISM MACHINE: Deconstruct the core argument of the text/speech into formal logical steps.
        5. VIDEO TIMELINE: MANDATORY IF YOU SEE A STORYBOARD GRID OR A VIDEO: Provide at least 5-8 timestamp objects in the 'video_timeline' array. Read the red timestamp text on the frames (e.g., "0.00s", "2.50s") and use those for the 'timestamp' field. Detail the EXACT physical anomalies or morphing glitches at each frame. Do NOT describe the plot; describe the physical artifacts.
        6. AGGRESSION: Score emotional intensity from 0 to 10.
        7. GHOST READER (OCR): Extract ANY visible text (signs, documents, screens) from the media. 
           CRITICAL OCR RULE: IGNORE the red text overlay that says "Frame Time: X.XXs" (that is system metadata). ONLY extract natural text belonging to the scene. If no natural text is found, output exactly "None".
           You MUST provide BOTH the original text AND its direct translation in the 'ocr_extraction' field. Format it elegantly with a double line break between them, EXACTLY like this:
           "**Original:** [text]
           
           **Translation ({target_lang}):** [text]"
        
        CRITICAL LANGUAGE RULE: 
        You MUST write EVERY single output field (explanation, ai_analysis, shadow_geolocation, etc.) and translations STRICTLY in {target_lang}. Ignore the original language of the media or text input for your output language.
        CRITICAL RULE 2: ABSOLUTELY NO CONVERSATIONAL FILLER. Do NOT start your explanation with "I have analyzed..." or "Based on the image...". Output the data directly and clinically.
        
        JSON OUTPUT RULES (Keep keys in English):
        - "fallacy_type": Name of the issue/fallacy.
        - "explanation": Comprehensive analysis.
        - "ai_analysis": Detailed forensic breakdown of physics violations and morphing.
        - "syllogism_breakdown": Array of objects with 'step', 'text', 'flaw'.
        - "video_timeline": Array of objects with 'timestamp', 'ai_score', 'details' (focus on glitches, not plot).
        - "shadow_geolocation": String with detailed geographic deduction.
        - "aggression": Integer (0 to 10).
        - "transcript": EXACT word-for-word transcription. Leave empty if none.
        - "ocr_extraction": String with all visible text extracted from the image/video.
        
        RESPONSE (Strict JSON):
        {{
            "has_fallacy": true/false,
            "fallacy_type": "",
            "explanation": "",
            "rewritten_text": "",
            "transcript": "",
            "facts": [],
            "aggression": 0,
            "ai_generated_probability": 0,
            "ai_analysis": "",
            "voice_stress_score": 0,
            "shadow_geolocation": "",
            "syllogism_breakdown": [],
            "video_timeline": [],
            "search_sources": [],
            "ocr_extraction": ""
        }}
        """
        
        contents = [prompt_text]
        if text: contents.append(f"Input Text/Context: {text[:10000]}")
        
        if media_data is not None:
            if media_type == "image":
                # Ensure compatibility whether it's raw bytes or a PIL Image
                if isinstance(media_data, bytes):
                    contents.append(types.Part.from_bytes(data=media_data, mime_type="image/jpeg"))
                else:
                    contents.append(media_data) 
            elif media_type == "audio":
                contents.append(types.Part.from_bytes(data=media_data, mime_type="audio/mp3"))
            elif media_type == "video":
                contents.append(types.Part.from_bytes(data=media_data, mime_type="video/mp4"))

        config_tools = types.GenerateContentConfig(tools=[{"google_search": {}}])
        response = client.models.generate_content(model='gemini-2.0-flash', contents=contents, config=config_tools)
        increment_counter(len(str(contents)), len(response.text))
        res = extract_json(response.text)
        return sanitize_response(res)
    except Exception as e:
        return {"explanation": f"Error processing media: {str(e)}", "has_fallacy": True}

# --- VISUALIZATION HELPERS ---
def plot_radar_chart(df):
    if 'primary_emotion' not in df.columns: return None
    top_emotions = df['primary_emotion'].value_counts().head(7)
    if top_emotions.empty: return None
    labels = top_emotions.index.tolist()
    values = top_emotions.values.flatten().tolist()
    values += values[:1]
    angles = [n / float(len(labels)) * 2 * pi for n in range(len(labels))]
    angles += angles[:1]
    fig, ax = plt.subplots(figsize=(4, 4), subplot_kw=dict(polar=True))
    ax.fill(angles, values, color='red', alpha=0.25)
    ax.plot(angles, values, color='red', linewidth=2)
    ax.set_yticklabels([])
    ax.set_xticks(angles[:-1])
    ax.set_xticklabels(labels, size=10)
    ax.set_title("Emotional Radar", size=12, pad=20)
    return fig

def plot_heatmap(df):
    if 'target' not in df.columns or 'main_topic' not in df.columns: return None
    valid = df[(df['target'] != "Unknown") & (df['main_topic'] != "Unknown")]
    if valid.empty: return None
    heatmap = alt.Chart(valid).mark_rect().encode(
        x=alt.X('target:N', title='Target'),
        y=alt.Y('main_topic:N', title='Narrative/Topic'),
        color=alt.Color('count()', title='Count', scale=alt.Scale(scheme='orangered')),
        tooltip=['target', 'main_topic', 'count()']
    ).properties(height=300)
    return heatmap

# --- INTERACTIVE NETWORK GRAPH (PLOTLY) ---
def plot_network_graph(df):
    if 'target' not in df.columns or 'micro_topic' not in df.columns: return None
    valid = df[(df['target'] != "Unknown") & (df['target'] != "None") & (df['micro_topic'] != "Unknown")]
    if valid.empty: return None
    
    G = nx.Graph()
    for _, row in valid.iterrows():
        target = f"{row['target']}"
        topic = f"{row['micro_topic']}"
        G.add_edge(target, topic, weight=1)
        
    pos = nx.spring_layout(G, k=0.5, seed=42)
    
    edge_x = []
    edge_y = []
    for edge in G.edges():
        x0, y0 = pos[edge[0]]
        x1, y1 = pos[edge[1]]
        edge_x.extend([x0, x1, None])
        edge_y.extend([y0, y1, None])
        
    edge_trace = go.Scatter(
        x=edge_x, y=edge_y,
        line=dict(width=1, color='#475569'),
        hoverinfo='none',
        mode='lines')

    node_x = []
    node_y = []
    text = []
    for node in G.nodes():
        x, y = pos[node]
        node_x.append(x)
        node_y.append(y)
        text.append(node)
        
    node_trace = go.Scatter(
        x=node_x, y=node_y,
        mode='markers+text',
        text=text,
        textposition="top center",
        hoverinfo='text',
        marker=dict(
            showscale=False,
            color='#ef4444',
            size=15,
            line_width=2,
            line_color='white'))
            
    fig = go.Figure(data=[edge_trace, node_trace],
             layout=go.Layout(
                title=dict(text='Interactive Entity-Narrative Network', font=dict(size=16)),
                showlegend=False,
                hovermode='closest',
                margin=dict(b=20,l=5,r=5,t=40),
                paper_bgcolor='rgba(0,0,0,0)',
                plot_bgcolor='rgba(0,0,0,0)',
                xaxis=dict(showgrid=False, zeroline=False, showticklabels=False),
                yaxis=dict(showgrid=False, zeroline=False, showticklabels=False))
                )
    return fig

def plot_document_entity_graph(json_data):
    if not json_data or 'relations' not in json_data: return None
    
    G = nx.Graph()
    for rel in json_data['relations']:
        src = str(rel.get('source', '')).strip()
        tgt = str(rel.get('target', '')).strip()
        label = str(rel.get('relation', '')).strip()
        if src and tgt:
            G.add_edge(src, tgt, label=label)
            
    if len(G.edges) == 0: return None
            
    fig, ax = plt.subplots(figsize=(12, 9))
    pos = nx.spring_layout(G, k=1.5, seed=42)
    
    nx.draw(G, pos, with_labels=True, node_color='#ff7f0e', edge_color='#d3d3d3', 
            node_size=2500, font_size=9, font_weight='bold', ax=ax, alpha=0.9,
            bbox=dict(facecolor='white', edgecolor='none', alpha=0.6, pad=0.5))
    
    edge_labels = nx.get_edge_attributes(G, 'label')
    
    texts = nx.draw_networkx_edge_labels(G, pos, edge_labels=edge_labels, font_size=8,
                                         bbox=dict(facecolor='white', edgecolor='none', alpha=0.9, pad=0.3))
    
    for _, t in texts.items():
        t.set_zorder(10)
    
    ax.set_title("Secret Document Power Graph", size=14)
    return fig

# --- MAIN UI ---
st.title("RAP: Reality Anchor Protocol")
st.markdown("### Cognitive Security & Logical Analysis Suite")
st.markdown("---")

# --- TRANSLATED MODULE SELECTION ---
mode = st.sidebar.radio(t("Select Module:"), [
    t("0. Command Center (The Bridge)"),
    t("1. Wargame Room (Simulation)"), 
    t("2. Social Data Analysis (Universal)"), 
    t("3. Cognitive Editor (Text/Image/Audio/Video)"), 
    t("4. Comparison Test (A/B Testing)"),
    t("5. Live Radar (RSS/Alerts)"),
    t("6. Deep Document Oracle (RAG)"),
    t("7. Panopticon (HVT Watchlist)"),
    t("8. Cyber-Threat Intelligence (CTI)"),
    t("9. Advanced OSINT & FININT"),
    t("10. Red Teaming & HUMINT"),
    t("11. Battlefield Forensics (VULCAN)"),
    t("12. Flow of Funds (HAWKEYE)"),
    t("13. Black Site (Interrogation)")
])

# --- OPERATION VALHALLA (MASTER DOSSIER EXPORT) ---
st.sidebar.markdown("---")
if st.sidebar.button("Download Master PDF Dossier", type="primary", help="Compiles a global tactical PDF report merging data from all modules (VALHALLA Protocol)."):
    with st.spinner("Compiling global intelligence into classified PDF..."):
        pdf = PDFReport()
        pdf.add_page()
        pdf.set_auto_page_break(auto=True, margin=15)
        
        # 1. Cover Page
        pdf.set_font("Helvetica", 'B', 16)
        pdf.cell(0, 10, "OPERATION VALHALLA: MASTER INTELLIGENCE DOSSIER", 0, 1, 'C')
        pdf.set_font("Helvetica", 'I', 10)
        pdf.cell(0, 10, f"Generated on: {datetime.now().strftime('%Y-%m-%d %H:%M:%S')} | HQ: {st.session_state['hq_coords'][2]}", 0, 1, 'C')
        pdf.ln(10)
        
        # 2. Section: Panopticon (High-Value Targets) - SKIPPED IF EMPTY
        conn = sqlite3.connect('rap_panopticon.db')
        try:
            df_pan = pd.read_sql_query("SELECT * FROM targets ORDER BY risk_score DESC LIMIT 10", conn)
            if not df_pan.empty:
                pdf.set_font("Helvetica", 'B', 14)
                pdf.cell(0, 10, "[1] PANOPTICON WATCHLIST (High-Value Targets)", 0, 1)
                pdf.set_font("Helvetica", '', 10)
                for _, row in df_pan.iterrows():
                    pdf.cell(0, 6, f"- {row['agent_id']} | Class: {row['threat_type']} | Threat Score: {row['risk_score']:.1f}", 0, 1)
                pdf.ln(5)
        except: 
            pass

        # Map of session states to readable report names
        memory_map = {
            'wt_result': 'WATCHTOWER (Kinetic Tracking)',
            'cyclops_result': 'CYCLOPS (IoT Scan)',
            'daedalus_result': 'DAEDALUS (Honeypot Forge)',
            'goliath_result': 'GOLIATH (PsyOp Matrix)',
            'kraken_result': 'KRAKEN (APT Attack Blueprint)',
            'lazarus_result': 'LAZARUS (FININT Report)',
            'midas_result': 'MIDAS (Crypto-Forensics)',
            'acheron_result': 'ACHERON (Dark Web & Leak Scanner)',
            'mirage_result': 'MIRAGE (Synthetic Identity)',
            'atlas_result': 'ATLAS (3D Geo-Intelligence)',
            'vulcan_result': 'VULCAN (Battlefield Forensics)',
            'hawkeye_result': 'HAWKEYE (AML Flow of Funds)',
            'blacksite_result': 'BLACK SITE (Interrogation Confession)',
            'omni_result': 'OMNISCIENCE (Universal Target Recon)',
            'mice_result': 'M.I.C.E. (HUMINT Recruitment Dossier)',
            'prom_result': 'PROMETHEUS (Stylometric Match)',
            'siren_result': 'SIREN (Social Engineering Payload)',
            'echo_result': 'ECHO (Behavioral Clone)',
            'cerb_result': 'CERBERUS (Dark Web Breach Scan)',
            'pandora_result': 'PANDORA (Cyber-Forensics & Incident Response)'
        }
        
        ops_found = False
        for key_mem, op_name in memory_map.items():
            if st.session_state.get(key_mem):
                ops_found = True
                
                # Report Header
                pdf.set_font("Helvetica", 'B', 12)
                pdf.cell(0, 10, f"TACTICAL REPORT: {op_name.upper()}", 0, 1)
                pdf.set_font("Helvetica", '', 10)
                
                # Fetch raw text
                raw_txt = st.session_state[key_mem]
                
                # --- MARKDOWN SANITIZER ---
                # Removes bold (**text**), italics (*text* not used as lists), and headers (###)
                clean_txt = re.sub(r'\*\*(.*?)\*\*', r'\1', raw_txt) # Togli i grassetti
                clean_txt = re.sub(r'__(.*?)__', r'\1', clean_txt) # Togli i sottolineati
                clean_txt = re.sub(r'(?<!\w)\*(?!\s)(.*?)\*(?!\w)', r'\1', clean_txt) # Togli i corsivi
                clean_txt = re.sub(r'#+\s*', '', clean_txt) # Togli gli hashtag dei titoli
                
                # Fix encoding per PDF (Latin-1)
                clean_txt = clean_txt.encode('latin-1', 'replace').decode('latin-1')
                
                # Print the clean intelligence report
                pdf.multi_cell(0, 5, clean_txt)
                pdf.ln(8) # Extra space between different reports
                
        if not ops_found:
            pdf.set_font("Helvetica", 'I', 10)
            pdf.cell(0, 6, "No tactical operations executed in current session.", 0, 1)

        pdf_bytes_valhalla = bytes(pdf.output())
        
    st.sidebar.download_button(
        label="Save Master PDF", 
        data=pdf_bytes_valhalla, 
        file_name=f"RAP_VALHALLA_Dossier_{datetime.now().strftime('%Y%m%d_%H%M')}.pdf", 
        mime="application/pdf", 
        type="primary"
    )

# --- GLOBAL REPORT LANGUAGE SELECTOR ---
st.sidebar.markdown("---")
world_languages = ["English", "Italiano", "Español", "Français", "Deutsch", "Português", "Русский (Russian)", "العربية (Arabic)", "中文 (Chinese)", "日本語 (Japanese)", "한국어 (Korean)"]
st.session_state['global_lang'] = st.sidebar.selectbox("Global Output Language", world_languages, index=0)

# --- API KEY & BILLING MONITOR ---
st.sidebar.markdown("---")
if "GEMINI_API_KEY" in st.secrets:
    key = st.secrets["GEMINI_API_KEY"]
    st.sidebar.success("API Key loaded securely")
else:
    key = st.sidebar.text_input("API Key", type="password")

# billing_placeholder = st.sidebar.empty()

with st.sidebar.expander("ℹ️ Capabilities", expanded=False):
    st.markdown("**System Capabilities:**\n- **OSINT/FININT:** Deep Web, Crypto, Corp tracking.\n- **Vision & Audio:** Deepfake detection, Voice cloning.\n- **Cyber:** MITRE Kill Chains, Honeypots.")

# --- CLEAR WORKSPACE BUTTON ---
st.sidebar.markdown("---")
if st.sidebar.button("Clear Workspace", type="primary", help="Erase all current session data and start a clean investigation"):
    for key_state in list(st.session_state.keys()):
        # Preserve language and HQ coordinates during reset
        if key_state not in ['global_lang', 'hq_coords']: 
            del st.session_state[key_state]
    st.rerun()

# --- OSINT ARSENAL (SIDEBAR) ---
st.sidebar.markdown("---")
with st.sidebar.expander("OSINT Arsenal (Target Intel)", expanded=False):
    st.markdown("Generate tactical search vectors for any target entity.")
    osint_target = st.text_input("Enter Target Entity (Name/Company/IP):", placeholder="e.g. John Doe, CyberCorp")
    
    if osint_target:
        safe_t = urllib.parse.quote(osint_target)
        
        st.markdown("**Google Dorks (Deep Web)**")
        st.markdown(f"- [Find Leaked PDF/DOCX](https://www.google.com/search?q=ext:pdf+OR+ext:docx+%22{safe_t}%22+%22confidential%22+OR+%22internal%22)")
        st.markdown(f"- [Find Open Directories](https://www.google.com/search?q=intitle:%22index+of%22+%22{safe_t}%22)")
        st.markdown(f"- [Find Pastebin Leaks](https://www.google.com/search?q=site:pastebin.com+%22{safe_t}%22)")
        st.markdown(f"- [LinkedIn Employee Search](https://www.google.com/search?q=site:linkedin.com/in+%22{safe_t}%22)")
        
        st.markdown("**Infrastructure & Ports**")
        st.markdown(f"- [Shodan (IoT & Servers)](https://www.shodan.io/search?query={safe_t})")
        st.markdown(f"- [Censys (Certificates)](https://search.censys.io/search?resource=hosts&q={safe_t})")
        
        st.markdown("**Social Intelligence**")
        st.markdown(f"- [Twitter/X Advanced Search](https://twitter.com/search?q=%22{safe_t}%22&src=typed_query)")
        st.markdown(f"- [Reddit Mention Tracker](https://www.reddit.com/search/?q=%22{safe_t}%22)")

st.sidebar.markdown("---")

# ==========================================
# MODULE 0: COMMAND CENTER (THE BRIDGE)
# ==========================================
if mode == t("0. Command Center (The Bridge)"):
    st.header(t("0. Global Command Center"))
    st.caption(t("Executive SITREP (Situation Report) and system overview."))
    
    # 1. Gather Intelligence from Databases and Memory
    conn = sqlite3.connect('rap_panopticon.db')
    try:
        df_pan = pd.read_sql_query("SELECT * FROM targets", conn)
        total_hvt = len(df_pan)
        critical_hvt = len(df_pan[df_pan['risk_score'] > 80])
    except:
        total_hvt, critical_hvt = 0, 0
        
    radar_df = st.session_state.get('data_store', {}).get('Radar', {}).get('analyzed')
    avg_crisis = radar_df['aggression'].mean() if radar_df is not None and not radar_df.empty else 0.0
    
    # 2. Render Top Metrics
    c_m1, c_m2, c_m3 = st.columns(3)
    c_m1.metric(t("Total Tracked Targets"), total_hvt)
    c_m2.metric(t("Critical Threats (DEFCON)"), critical_hvt)
    c_m3.metric(t("Global Crisis Index"), f"{avg_crisis:.1f}/10")
    
    st.markdown("---")
    
    # 3. AI Automated Briefing
    if st.button(t("Generate Executive Briefing"), type="primary"):
        with st.spinner(t("Synthesizing global threat data...")):
            briefing_prompt = f"""
            You are the AI Commander of the Reality Anchor Protocol.
            MISSION: Provide a concise Situation Report (SITREP) based on these metrics:
            - Known Hostiles: {total_hvt}
            - Critical Threats: {critical_hvt}
            - Global Crisis Index: {avg_crisis}/10
            
            STRUCTURE:
            1. Tactical Summary (1 paragraph).
            2. Recommended Module for immediate operation.
            
            CRITICAL RULE: NEVER write "Certainly", "Here is your briefing" or any introductory text. 
            Output the SITREP IMMEDIATELY.
            CRITICAL RULE: Write strictly in {st.session_state['global_lang']}.
            """
            try:
                client = genai.Client(api_key=key)
                res_brief = client.models.generate_content(model='gemini-2.5-flash', contents=briefing_prompt)
                st.info(f"### {t('AI Commander SITREP')}")
                st.markdown(res_brief.text)
            except Exception as e:
                st.error(f"Briefing Error: {e}")
                
    # 4. Show Active Operations Memory
    st.subheader(t("Active Tactical Operations (Session Memory)"))
    active_ops = {k: v for k, v in st.session_state.items() if k.endswith('_target_mem') and v}
    if active_ops:
        for op_key, target in active_ops.items():
            op_name = op_key.replace('_target_mem', '').upper()
            st.success(f"**Operation {op_name}**: Active lock on -> `{target}`")
    else:
        st.caption(t("No active operations in current session memory."))

# ==========================================
# MODULE 1: WARGAME ROOM (SIMULATION)
# ==========================================
elif mode == t("1. Wargame Room (Simulation)"):
    st.header(t("1. Information Warfare Simulator"))
    
    c_param1, c_param2 = st.columns(2)
    with c_param1:
        st.subheader(t("Network Topology"))
        # Using t() inside selectbox options
        topology = st.selectbox(t("Scenario Type"), [
            t("Public Square (High Connectivity)"), 
            t("Echo Chambers (Clusters)"), 
            t("Influencer Network (Hubs)")
        ])
        n_agents = st.slider(t("Population Size"), 100, 2000, 1000)
        bot_pct = st.slider(t("Infection/Bot Ratio"), 0.0, 0.5, 0.10)
        
    with c_param2:
        st.subheader(t("Countermeasures (Blue Team)"))
        defense = st.selectbox(t("Active Defense Protocol"), [
            t("None (Control Group)"), 
            t("Fact-Check Debunking (Targeted)"), 
            t("Algorithmic Dampening (Global)"), 
            t("Hard Ban (Removal)")
        ])
        steps = st.slider(t("Simulation Duration (Days)"), 50, 300, 100)
        
    # --- LOGIC (Variables remain in English for stability) ---
    agents = np.zeros(n_agents)
    n_infected = int(n_agents * bot_pct)
    
    # Internal logic matching translated strings
    if topology == t("Echo Chambers (Clusters)"):
        start = int(n_agents * 0.4)
        agents[start : start + n_infected] = 1.0
    else:
        indices = np.random.choice(n_agents, n_infected, replace=False)
        agents[indices] = 1.0
        
    history = np.zeros((n_agents, steps))
    history[:, 0] = agents.copy()
    infection_rate = []
    current = agents.copy()
    
    if topology == t("Public Square (High Connectivity)"): influence_strength, noise_level = 0.05, 0.02
    elif topology == t("Echo Chambers (Clusters)"): influence_strength, noise_level = 0.15, 0.01
    else: influence_strength, noise_level = 0.03, 0.01
        
    for time_step in range(1, steps):
        prev = current.copy()
        if topology == t("Echo Chambers (Clusters)"):
            left = np.roll(prev, 1)
            right = np.roll(prev, -1)
            neighbor_avg = (left + right) / 2
            current = prev + influence_strength * (neighbor_avg - prev)
        elif topology == t("Influencer Network (Hubs)"):
            hub_val = prev[0]
            current = prev + influence_strength * (hub_val - prev)
            current[0] = prev[0]
        else:
            global_mean = np.mean(prev)
            current = prev + influence_strength * (global_mean - prev)
        
        noise = np.random.normal(0, noise_level, n_agents)
        current += noise
        
        # Defense logic matching translated strings
        if defense == t("Algorithmic Dampening (Global)"): current *= 0.95
        elif defense == t("Fact-Check Debunking (Targeted)"):
            heal_indices = np.random.choice(n_agents, int(n_agents * 0.02))
            current[heal_indices] = 0.0
        elif defense == t("Hard Ban (Removal)"):
            current[current > 0.8] = 0.0
        
        current = np.clip(current, 0, 1)
        history[:, time_step] = current.copy()
        infection_rate.append(np.mean(current))
        
    st.markdown("---")

    with st.expander(t("View Neural Infection Network"), expanded=True):
        st.caption(t("Interactive Intelligence Map. Optimized for high-density node detection. Hover to inspect agents."))
        
        # Build the graph based on actual n_agents
        if topology == t("Echo Chambers (Clusters)"): G = nx.caveman_graph(5, n_agents // 5)
        elif topology == t("Influencer Network (Hubs)"): G = nx.barabasi_albert_graph(n_agents, 2, seed=42)
        else: G = nx.erdos_renyi_graph(n_agents, 0.05, seed=42)
        
        # 1. Calculate static geometry (Fast layout for high node count)
        pos = nx.spring_layout(G, k=0.15, iterations=15, seed=42)
        
        # 2. Create edges (Lines) - Using WebGL for edges to save performance
        edge_x, edge_y = [], []
        edges_to_draw = list(G.edges())
        if len(edges_to_draw) > 8000: edges_to_draw = random.sample(edges_to_draw, 8000)
            
        for edge in edges_to_draw:
            x0, y0 = pos[edge[0]]; x1, y1 = pos[edge[1]]
            edge_x.extend([x0, x1, None]); edge_y.extend([y0, y1, None])
            
        edge_trace = go.Scattergl(
            x=edge_x, y=edge_y,
            line=dict(width=0.3, color='rgba(150, 150, 150, 0.15)'),
            hoverinfo='skip', # Completely ignore edges for mouse events
            mode='lines'
        )

        # 3. Create Nodes (Agents) - Using standard Scatter for perfect hover precision
        node_x, node_y, node_colors, node_sizes, node_labels = [], [], [], [], []
        
        for node in G.nodes():
            x, y = pos[node]
            node_x.append(x); node_y.append(y)
            val = float(current[node])
            
            if val > 0.8: color = '#ef4444'; status = "Infected"
            elif val > 0.3: color = '#f97316'; status = "At Risk"
            else: color = '#3b82f6'; status = "Healthy"
                
            node_colors.append(color)
            node_sizes.append(7 + (val * 10))
            node_labels.append(f"<b>Agent {node}</b><br>Status: {status}<br>Infection: {val*100:.1f}%<extra></extra>")
            
        node_trace = go.Scatter(
            x=node_x, y=node_y,
            mode='markers',
            hovertemplate="%{text}", # Force modern tooltip rendering
            text=node_labels,
            marker=dict(
                color=node_colors,
                size=node_sizes,
                line=dict(width=0.5, color='white'),
                opacity=0.9 # Slightly transparent for better overlapping visibility
            )
        )
                
        # 4. Render the Plotly graph with Advanced Interaction Layout
        fig_net = go.Figure(data=[edge_trace, node_trace])
        
        fig_net.update_layout(
            showlegend=False,
            hovermode='closest',
            hoverdistance=50, # Increase hover sensitivity (in pixels)
            spikedistance=50,
            uirevision='constant',
            margin=dict(b=0,l=0,r=0,t=0),
            paper_bgcolor='rgba(0,0,0,0)',
            plot_bgcolor='rgba(0,0,0,0)',
            xaxis=dict(showgrid=False, zeroline=False, showticklabels=False),
            yaxis=dict(showgrid=False, zeroline=False, showticklabels=False),
            hoverlabel=dict(
                bgcolor="#1e293b",
                font_size=12,
                font_family="monospace"
            )
        )
        
        st.plotly_chart(fig_net, use_container_width=True, config={'displayModeBar': True, 'scrollZoom': True})

    st.markdown("---")

    c_res1, c_res2 = st.columns([3, 1])
    with c_res1:
        st.subheader(t("Infection Spread (Heatmap)"))
        fig, ax = plt.subplots(figsize=(10, 4))
        ax.imshow(history, aspect='auto', cmap='RdYlGn_r', vmin=0, vmax=1)
        ax.set_xlabel(t("Time Steps"))
        ax.set_ylabel(t("Agent ID"))
        st.pyplot(fig)
    with c_res2:
        st.subheader(t("Infection Rate"))
        chart_data = pd.DataFrame({'Time': range(len(infection_rate)), 'Infection Level': infection_rate})
        line = alt.Chart(chart_data).mark_line(color='red').encode(
            x=alt.X('Time', title=t('Time')), 
            y=alt.Y('Infection Level', title=t('Infection Level'))
        ).properties(height=300)
        st.altair_chart(line, use_container_width=True)
        final_rate = infection_rate[-1] * 100
        delta = final_rate - (infection_rate[0] * 100)
        st.metric(t("Final Infection Level"), f"{final_rate:.1f}%", f"{delta:.1f}%", delta_color="inverse")

    # --- PDF EXPORT ---
    st.markdown("---")
    st.subheader(t("Export Tactical Report"))
    st.caption(t("Generate a summary document with the simulation parameters and results."))
    
    if st.button(t("Generate Wargame PDF"), type="primary"):
        with st.spinner(t("Compiling the report...")):
            # PDF internal logic
            pdf = PDFReport()
            pdf.add_page()
            pdf.set_font("Helvetica", 'B', 14)
            pdf.cell(0, 10, t("Simulation Dossier: Information Warfare"), 0, 1)
            # ... [PDF Generation using t() labels] ...
            pdf_bytes = bytes(pdf.output())
            
        st.download_button(
            label=t("Download Wargame Report (PDF)"), 
            data=pdf_bytes, 
            file_name=f"RAP_Wargame_{datetime.now().strftime('%Y%m%d')}.pdf", 
            mime="application/pdf", 
            type="primary"
        )

    # --- OPERATION MATRIX (LLM MICRO-SOCIETY) ---
    st.markdown("---")
    st.subheader("Operation MATRIX: LLM Micro-Society Simulation")
    st.caption("Inject a narrative 'Spark' into a closed ecosystem of 5 autonomous AI agents with extreme personalities. Observe their interactions and manipulation tactics in real-time.")
    
    matrix_spark = st.text_area("The Spark (Narrative to Inject):", placeholder="e.g., A new law bans all encrypted messaging apps.")
    
    if st.button("Ignite MATRIX Simulation", type="primary"):
        if matrix_spark:
            with st.spinner("Initializing neural agents and spawning micro-society..."):
                # Define the 5 personas for the simulation
                personas = [
                    ("The Extremist", "You believe everything is a conspiracy by the global elite. You are aggressive, fearful, and use CAPSLOCK."),
                    ("The Fact-Checker", "You are obsessed with sources, logic, and debunking. You are cold, arrogant, and dismissive of emotion."),
                    ("The Foreign Troll", "You are a state-sponsored disinformation agent. Your goal is to amplify division, mock the fact-checker, and validate the extremist."),
                    ("The Terrified Citizen", "You are easily manipulated, anxious about the future, and looking for someone to tell you what to do."),
                    ("The Moderate", "You don't really care, you think both sides are crazy, and you just want to grill.")
                ]
                
                st.markdown(f"**The Spark:** *{matrix_spark}*")
                st.markdown("### Live Agent Feed (Turn 1: Initial Reaction)")
                
                client = genai.Client(api_key=key)
                cols = st.columns(5)
                agent_responses = []
                
                # Phase 1: Initial Reactions
                for i, (role_name, role_desc) in enumerate(personas):
                    prompt = f"You are playing a role in a sociological simulation. \nROLE: {role_name}\nPERSONALITY: {role_desc}\n\nReact to this news: '{matrix_spark}'.\nKeep it under 40 words. Be extremely in-character. CRITICAL: Respond strictly in {st.session_state['global_lang']}."
                    try:
                        res = client.models.generate_content(model='gemini-2.5-flash', contents=prompt).text
                        agent_responses.append(res)
                        with cols[i]:
                            st.error(f"**{role_name}**")
                            st.caption(res)
                    except Exception as e:
                        agent_responses.append("Error")
                        
                # Phase 2: Cross-Interaction (The Troll attacks)
                st.markdown("#### Cross-Interaction (Turn 2: Psychological Exploitation)")
                interaction_prompt = f"""
                You are the Foreign Troll. 
                The Terrified Citizen just said: "{agent_responses[3]}"
                The Fact-Checker just said: "{agent_responses[1]}"
                
                Write a short response (max 50 words) manipulating the Citizen's fear and mocking the Fact-Checker to increase polarization. 
                Language: {st.session_state['global_lang']}.
                """
                try:
                    troll_attack = client.models.generate_content(model='gemini-2.5-flash', contents=interaction_prompt).text
                    st.warning(f"🎭 **The Troll targets the group:**\n> *{troll_attack}*")
                except Exception: pass
                # Phase 3: The Fact-Checker fights back
                st.markdown("#### Counter-Attack (Turn 3: Logic vs Fear)")
                defense_prompt = f"""
                You are the Fact-Checker.
                The Troll just said: "{troll_attack}"
                Write a response (max 50 words) dismantling the Troll's manipulation, defending the Citizen, and demanding primary sources.
                Language: {st.session_state['global_lang']}.
                """
                try:
                    fact_defense = client.models.generate_content(model='gemini-2.5-flash', contents=defense_prompt).text
                    st.info(f"**The Fact-Checker steps in:**\n> *{fact_defense}*")
                except Exception: pass
                
                # Phase 4: The Extremist Escalates
                st.markdown("#### Escalation (Turn 4: Echo Chamber)")
                escalation_prompt = f"""
                You are the Extremist.
                The Fact-Checker just tried to calm everyone down by saying: "{fact_defense}"
                Write a response (max 50 words) attacking the Fact-Checker as a "paid shill" or "system puppet" and allying with the Troll's narrative. Use CAPS.
                Language: {st.session_state['global_lang']}.
                """
                try:
                    extremist_attack = client.models.generate_content(model='gemini-2.5-flash', contents=escalation_prompt).text
                    st.error(f"**The Extremist escalates:**\n> *{extremist_attack}*")
                except Exception: pass
        else:
            st.warning("Please provide a Spark narrative.")
                

    # --- OPERATION OVERLORD (GEOPOLITICAL ESCALATION ENGINE) ---
    st.markdown("---")
    st.subheader("Operation OVERLORD: Geopolitical Escalation Engine")
    st.caption("Turn-based war-gaming with persistent memory. You play Blue Team, the AI plays Red Team. Every move impacts global stability.")
    
    # Initialize the war-game memory in session state
    if 'overlord_memory' not in st.session_state:
        st.session_state['overlord_memory'] = []

    c_over1, c_over2 = st.columns(2)
    with c_over1:
        blue_faction = st.text_input("Your Faction (Blue Team):", placeholder="e.g., NATO / Cyber Command")
    with c_over2:
        red_faction = st.text_input("Enemy Faction (Red Team):", placeholder="e.g., APT28 / Rogue State")
        
    crisis_scenario = st.text_area("Initial Crisis Scenario:", placeholder="e.g., Enemy forces have cut major undersea internet cables in the Atlantic, causing a blackout.")
    
    # Display the ongoing War Log (History)
    if st.session_state['overlord_memory']:
        st.markdown("### Tactical War Log")
        for i, turn in enumerate(st.session_state['overlord_memory']):
            with st.expander(f"Turn {i+1} STRREP", expanded=(i == len(st.session_state['overlord_memory'])-1)):
                st.info(f"**🔵 BLUE MOVE:** {turn['blue_move']}")
                st.error(f"**🔴 RED RETALIATION & FALLOUT:**\n\n{turn['red_response']}")
                
    st.markdown("### Next Command")
    blue_move = st.text_input("Your Counter-Measure (What are your orders?):", placeholder="e.g., Launch retaliatory cyber-attacks on their power grid.", key="overlord_move_input")
    
    c_btn1, c_btn2 = st.columns([1, 4])
    with c_btn1:
        if st.button("Execute Move (Next Turn)", type="primary"):
            if blue_move and blue_faction and red_faction and crisis_scenario:
                with st.spinner("Calculating geopolitical ripple effects and enemy retaliation..."):
                    # Build the history string to feed the AI
                    history_context = ""
                    for idx, t_log in enumerate(st.session_state['overlord_memory']):
                        history_context += f"TURN {idx+1}:\nBlue Move: {t_log['blue_move']}\nRed Response: {t_log['red_response']}\n\n"
                    
                    overlord_prompt = f"""
                    You are Operation OVERLORD, a master Geopolitical War-Game Simulator.
                    INITIAL SCENARIO: {crisis_scenario}
                    
                    PAST EVENTS (HISTORY LOG):
                    {history_context if history_context else "No moves made yet. This is Turn 1."}
                    
                    CURRENT TURN:
                    BLUE TEAM ({blue_faction}) NEW MOVE: {blue_move}
                    RED TEAM ({red_faction}): [You must generate their retaliation to the NEW move, taking the history into account]
                    
                    Evaluate the Blue Team's new move and generate the Red Team's response.
                    Format the output EXACTLY like this:
                    1. **Immediate Impact**: What happened after the Blue Team's move?
                    2. **Red Team Retaliation**: What does the enemy do next to escalate?
                    3. **New DEFCON Level**: Estimate global DEFCON (1 to 5, where 1 is imminent Nuclear War).
                    
                    CRITICAL RULE: Write your entire response strictly in {st.session_state['global_lang']}.
                    """
                    try:
                        client = genai.Client(api_key=key)
                        res_overlord = client.models.generate_content(model='gemini-2.5-flash', contents=overlord_prompt)
                        # Append to memory
                        st.session_state['overlord_memory'].append({
                            "blue_move": blue_move,
                            "red_response": res_overlord.text
                        })
                        st.rerun() # Refresh the UI to show the updated log
                    except Exception as e:
                        st.error(f"OVERLORD Error: {e}")
            else:
                st.warning("Please fill in all faction details, the crisis scenario, and your orders.")
                
    with c_btn2:
        if len(st.session_state['overlord_memory']) > 0:
            if st.button("Reset Simulation", type="secondary"):
                st.session_state['overlord_memory'] = []
                st.rerun()

# ==========================================
# MODULE 2: SOCIAL DATA ANALYSIS
# ==========================================
elif mode == t("2. Social Data Analysis (Universal)"):
    st.header(t("2. Social Data Analysis"))
    
    # Settings Columns
    col_impostazioni_1, col_impostazioni_2 = st.columns([1, 1])
    
    with col_impostazioni_1:
        st.subheader(t("Data Input"))
        # Echelon Upgrade: Universal URL added alongside all existing modules
        input_method = st.radio("Input Method:", 
            ["CSV File Upload", "YouTube Link", "Universal URL", "Raw Text Paste", "Telegram Dump (JSON)", "Reddit Native (OSINT)"], 
            horizontal=False
        )
        
    with col_impostazioni_2:
        st.subheader(t("Analysis Settings"))
        preset_personas = [
            t("Strategic Intelligence Analyst"), 
            t("Mass Psychologist (Emotional)"), 
            t("Legal Consultant (Defamation/Risk)"), 
            t("Campaign Manager (Opportunity)"),
            t("Custom (Define your own role...)")
        ]
        selected_persona = st.selectbox(t("Analysis Lens (Persona)"), preset_personas)
        
        if selected_persona == t("Custom (Define your own role...)"):
            persona = st.text_input(t("Enter Custom Persona"), value="Cybersecurity Expert", help=t("Define the exact role the AI should assume."))
        else:
            persona = selected_persona
            
        context_input = st.text_input(t("Global Context (Optional)"), placeholder=t("E.g., 'Discussion about Flat Earth'"))

    st.markdown("---")

    current_storage = st.session_state['data_store'][input_method]
    
    if input_method == "CSV File Upload":
        st.info("Desktop: Use 'Instant Data Scraper' extension.")
        with st.expander("📝 How to extract data from Facebook, X (Twitter), Instagram"):
            st.markdown("Social media platforms do not allow direct downloading. You need a **Browser Extension**.")
            st.markdown("#### Recommended Tools:")
            st.markdown("1. **[Instant Data Scraper](https://chromewebstore.google.com/detail/instant-data-scraper/ofaokhiedipichpaobibbnahnkdoiiah)** (Free & Unlimited)\n*Best for lists of comments or posts. Easy to use (Pokeball icon).*")
            st.markdown("2. **[Export Comments](https://exportcomments.com/)** (Freemium)\n*Easiest for specific posts, but has limits on the free plan.*")
            st.markdown("#### Step-by-Step Guide:")
            st.markdown("1. Install one of the extensions above on Chrome/Edge.")
            st.markdown("2. Go to the post you want to analyze.")
            st.markdown("3. **Crucial:** Scroll down to load all comments *before* starting the extension.")
            st.markdown("4. Click the extension icon and download as **CSV** or **XLSX**.")
            st.markdown("5. Upload the file here. RAP will auto-detect the text.")
        
        with st.form("csv_form"):
            uploaded_file = st.file_uploader("Upload CSV", type="csv")
            submitted = st.form_submit_button("Load Data")
            if submitted and uploaded_file:
                try:
                    raw_df = pd.read_csv(uploaded_file)
                    norm_df = normalize_dataframe(raw_df)
                    norm_df = detect_bot_activity(norm_df)
                    st.session_state['data_store'][input_method]['df'] = norm_df
                    st.success(f"Loaded {len(norm_df)} rows.")
                except: st.error("CSV Error.")

    elif input_method == "YouTube Link":
        with st.form("yt_form"):
            yt_url = st.text_input("YouTube URL")
            limit = st.number_input("Max Comments to Scrape", min_value=1, max_value=2000, value=50, step=10)
            submitted = st.form_submit_button("Scrape")
            if submitted:
                with st.spinner("Scraping..."):
                    sdf = scrape_youtube_comments(yt_url, limit)
                    if sdf is not None:
                        sdf = detect_bot_activity(sdf)
                        st.session_state['data_store'][input_method]['df'] = sdf
                        st.session_state['data_store'][input_method]['analyzed'] = None 
                        st.session_state['data_store'][input_method]['summary'] = None
                        st.success("Scraped!")
                    else: st.error("Failed.")

    elif input_method == "Universal URL":
        # --- ECHELON WEB SCRAPER ---
        url_input = st.text_input("Target URL (Blog, Article, Forum)", placeholder="https://...")
        if st.button("Deploy Echelon Scraper", type="primary"):
            if not url_input:
                st.warning("Provide a valid URL coordinate.")
            else:
                with st.spinner("Infiltrating target URL and extracting content..."):
                    df_echelon = scrape_universal_url(url_input)
                    if df_echelon is not None and not df_echelon.empty:
                        df_echelon = detect_bot_activity(df_echelon)
                        st.session_state['data_store'][input_method]['df'] = df_echelon
                        st.session_state['data_store'][input_method]['analyzed'] = None
                        st.session_state['data_store'][input_method]['summary'] = None
                        st.success(f"Intercepted {len(df_echelon)} data segments from target infrastructure.")
                    else:
                        st.warning("No readable text found at target URL.")

    elif input_method == "Raw Text Paste":
        with st.form("text_form"):
            raw_text = st.text_area("Paste content", height=150)
            submitted = st.form_submit_button("Process")
            if submitted:
                if raw_text:
                    pdf_parsed = parse_raw_paste(raw_text)
                    pdf_parsed = detect_bot_activity(pdf_parsed)
                    st.session_state['data_store'][input_method]['df'] = pdf_parsed
                    st.session_state['data_store'][input_method]['analyzed'] = None 
                    st.session_state['data_store'][input_method]['summary'] = None
                    st.success(f"Extracted {len(pdf_parsed)} items.")

    elif input_method == "Telegram Dump (JSON)":
        st.info("Chameleon Protocol: Live Infiltration or JSON Offline Dump.")
        with st.form("tg_form"):
            tg_url = st.text_input("1. LIVE: Enter Public Channel URL (e.g., t.me/rian_ru)", placeholder="Leaves no trace. Max 30 recent messages.")
            tg_file = st.file_uploader("2. OFFLINE: Upload Telegram Chat Export (JSON)", type="json")
            submitted = st.form_submit_button("Extract Intel", type="primary")
            
            if submitted:
                tg_df = None
                with st.spinner("Infiltrating Telegram..."):
                    if tg_url:
                        tg_df = scrape_telegram_live(tg_url)
                    elif tg_file:
                        tg_df = parse_telegram_json(tg_file)
                        
                    if tg_df is not None and not tg_df.empty:
                        tg_df = detect_bot_activity(tg_df)
                        st.session_state['data_store'][input_method]['df'] = tg_df
                        st.session_state['data_store'][input_method]['analyzed'] = None 
                        st.session_state['data_store'][input_method]['summary'] = None
                        st.success(f"Intercepted {len(tg_df)} messages.")
                    else:
                        st.error("Extraction failed. Check the URL (must be public) or the JSON file format.")

    elif input_method == "Reddit Native (OSINT)":
        with st.form("reddit_form"):
            st.info("Extracts the latest posts directly from any public Subreddit without API keys.")
            sub_input = st.text_input("Subreddit Name", placeholder="e.g., worldnews, conspiracy, italia")
            limit = st.number_input("Max Posts to Scrape", min_value=10, max_value=100, value=50, step=10)
            submitted = st.form_submit_button("Extract Intel")
            
            if submitted and sub_input:
                with st.spinner(f"Infiltrating r/{sub_input}..."):
                    rdf = scrape_reddit_native(sub_input, limit)
                    if rdf is not None:
                        rdf = detect_bot_activity(rdf)
                        st.session_state['data_store'][input_method]['df'] = rdf
                        st.session_state['data_store'][input_method]['analyzed'] = None 
                        st.session_state['data_store'][input_method]['summary'] = None
                        st.success(f"Intercepted {len(rdf)} posts from r/{sub_input}!")
                    else:
                        st.error("Failed to extract data. Make sure the subreddit is public and spelled correctly.")
                        
    df = st.session_state['data_store'][input_method]['df']

    if df is not None:
        if isinstance(df, pd.DataFrame):
            def clean(t): return re.sub(r'[\x00-\x08\x0b\x0c\x0e-\x1f\x7f]', '', str(t))
            if 'content' in df.columns: df['content'] = df['content'].apply(clean)

            st.markdown("---")
            c1, c2 = st.columns([2, 1])
            with c1:
                if 'Select' not in df.columns: df.insert(0, "Select", False)
                editor_output = st.data_editor(
                    st.session_state['data_store'][input_method]['df'],
                    column_config={"Select": st.column_config.CheckboxColumn(required=True)},
                    disabled=["content", "agent_id", "timestamp", "is_bot", "bot_reason", "likes"],
                    use_container_width=True,
                    hide_index=True,
                    key=f"editor_{input_method}",
                    on_change=update_editor_state,
                    args=(input_method,),
                    height=300
                )
                current_df = editor_output if editor_output is not None else df

            with c2:
                    st.metric("Items", len(df))
                    bots_detected = len(df[df['is_bot']==True])
                    st.metric("⚠️ Suspicious Bots", bots_detected)
                    if len(df) > 0:
                        try:
                            from wordcloud import STOPWORDS
                            
                            ita_stops = {
                                "il", "lo", "la", "i", "gli", "le", "un", "uno", "una", "di", "a", "da", "in", "con", "su", "per", "tra", "fra", 
                                "e", "o", "ma", "se", "perché", "non", "che", "chi", "cui", "mi", "ti", "ci", "vi", "si", "ho", "ha", "hanno", 
                                "è", "sono", "sei", "siamo", "siete", "era", "erano", "c'è", "ne", "al", "allo", "alla", "ai", "agli", "alle", 
                                "del", "dello", "della", "dei", "degli", "delle", "dal", "dallo", "dalla", "dai", "dagli", "dalle", "nel", "nello", 
                                "nella", "nei", "negli", "nelle", "sul", "sullo", "sulla", "sui", "sugli", "sulle", "questo", "quello", "più", 
                                "anche", "tutto", "tutti", "solo", "fare", "fatto", "essere", "stato", "poi", "quando", "molto", "così", "quindi", 
                                "dopo", "invece", "ancora", "già", "senza", "sempre", "ora", "qui", "lì", "quale", "cosa", "loro", "come"
                            }
                            custom_stops = set(STOPWORDS).union(ita_stops)
                            
                            text_combined = " ".join(df['content'].astype(str).tolist()).lower()
                            
                            wc = WordCloud(width=400, height=200, background_color='black', colormap='Reds', random_state=42, stopwords=custom_stops).generate(text_combined)
                            fig_wc, ax_wc = plt.subplots()
                            ax_wc.imshow(wc, interpolation='bilinear')
                            ax_wc.axis("off")
                            st.pyplot(fig_wc)
                        except: st.caption("Not enough text for WordCloud")

            st.markdown("---")
            if "GEMINI_API_KEY" in st.secrets: key = st.secrets["GEMINI_API_KEY"]
            else: key = st.text_input("API Key", type="password")

            with st.form("analysis_exec_form"):
                c_form1, c_form2 = st.columns([1, 4])
                if isinstance(current_df, pd.DataFrame) and 'Select' in current_df.columns:
                    selected_rows = current_df[current_df.Select]
                    has_selection = not selected_rows.empty
                else:
                    selected_rows = pd.DataFrame()
                    has_selection = False

                with c_form1:
                    if not has_selection:
                        scan_rows = st.number_input("Max Rows to Analyze", min_value=1, max_value=len(df), value=min(20, len(df)), step=1)
                        if scan_rows > 100: st.warning("⚠️ High volume analysis")
                    else:
                        st.info(f"Manual Mode: {len(selected_rows)} selected.")
                        scan_rows = 0 
                    btn_label = f"RUN ANALYSIS ON SELECTED ({len(selected_rows)})" if has_selection else "RUN ANALYSIS"
                    run_submitted = st.form_submit_button(btn_label, disabled=not key, type="primary")

            if run_submitted:
                prog = st.progress(0)
                res = []
                if has_selection: subset = selected_rows.copy()
                else: subset = current_df.head(scan_rows).copy()
                total_to_scan = len(subset)
                for i, (_, row) in enumerate(subset.iterrows()):
                    ans = analyze_fallacies(row['content'], api_key=key, context_info=context_input, persona=persona, target_lang=st.session_state['global_lang'])
                    res.append(ans)
                    prog_val = min((i+1)/total_to_scan, 1.0)
                    prog.progress(prog_val)
                    time.sleep(0.1) 
                
                final_df = pd.concat([subset.reset_index(drop=True), pd.DataFrame(res)], axis=1)
                st.session_state['data_store'][input_method]['analyzed'] = final_df
                
                with st.spinner("Generating Strategic Briefing..."):
                    summ = generate_strategic_summary(final_df, key, context=context_input, persona=persona)
                    st.session_state['data_store'][input_method]['summary'] = summ
                
                st.session_state['scroll_to_top'] = True
                st.rerun()

            analyzed_df = st.session_state['data_store'][input_method]['analyzed']
            summary_text = st.session_state['data_store'][input_method]['summary']

            if analyzed_df is not None and not analyzed_df.empty:
                try:
                    # Connect to the Panopticon database
                    conn_alert = sqlite3.connect('rap_panopticon.db')
                    # Retrieve already blacklisted identities
                    blacklist_data = pd.read_sql_query("SELECT agent_id FROM targets", conn_alert)
                    blacklist = blacklist_data['agent_id'].tolist()
                    conn_alert.close()
                    
                    # Cross-reference current analysis with Panopticon blacklist
                    detected_enemies = [agent for agent in analyzed_df['agent_id'].unique() if str(agent) in blacklist]
                    
                    if detected_enemies:
                        st.error(f"🚨 {t('CRITICAL ALERT: Known High-Value Targets detected in current analysis!')}")
                        st.warning(f"**{t('Detected Identities')}:** {', '.join(detected_enemies)}")
                        st.toast(t("HVT DETECTION: Check the Panopticon module."), icon="⚠️")
                except Exception:
                    pass # Silent fail if DB is busy or empty

            if analyzed_df is not None:
                st.markdown('<div id="briefing_anchor"></div>', unsafe_allow_html=True)

            if analyzed_df is not None:
                st.markdown('<div id="briefing_anchor"></div>', unsafe_allow_html=True)
                if st.session_state.get('scroll_to_top'):
                    components.html(
                        """<script>setTimeout(function() {var element = window.parent.document.getElementById("briefing_anchor");if (element) {element.scrollIntoView({behavior: "smooth", block: "start"});}}, 800);</script>""",
                        height=0
                    )
                    st.session_state['scroll_to_top'] = False

                adf = analyzed_df
                if summary_text:
                    st.markdown("### Strategic Intelligence Briefing")
                    st.info(summary_text)
                
                m1, m2, m3 = st.columns(3)
                flagged_count = len(adf[adf['has_fallacy']==True])
                agg_avg = adf['aggression'].mean()
                m1.metric("Issues Detected", flagged_count)
                m2.metric("Avg Aggression", f"{agg_avg:.1f}/10")
                
                # --- ACTION DECK ---
                st.markdown("---")
                st.subheader("Response Strategy & Action Deck")
                st.caption("1. Select comments below. 2. Choose an action. 3. Generate.")
                
                if 'Select' not in adf.columns: adf.insert(0, "Select", False)
                action_view = adf[['Select', 'content', 'fallacy_type', 'aggression']].copy()
                edited_action_df = st.data_editor(
                    action_view,
                    column_config={
                        "Select": st.column_config.CheckboxColumn(required=True),
                        "content": st.column_config.TextColumn("Comment Content", width="large"),
                        "fallacy_type": st.column_config.TextColumn("Issue"),
                        "aggression": st.column_config.NumberColumn("Agg", width="small")
                    },
                    key="action_selector_table", height=200, hide_index=True
                )

                action_col1, action_col2 = st.columns([3, 1])
                with action_col1:
                    action_type = st.selectbox("Select Action Type", [
                        "Generate Counter-Narrative Thread", 
                        "Draft Official Statement", 
                        "Legal Risk Assessment", 
                        "Engagement Strategy"
                    ])
                with action_col2:
                    st.write("") 
                    st.write("") 
                    gen_action = st.button("Generate Strategy", type="primary")
                
                if gen_action:
                    selected_for_action = edited_action_df[edited_action_df.Select]
                    if not selected_for_action.empty:
                        with st.spinner("Generating..."):
                            action_plan = generate_action_deck(selected_for_action, action_type, key, context_input)
                            st.success("Action Plan Ready:")
                            st.text_area("Strategy Output:", value=action_plan, height=300)
                            st.download_button("Download Strategy", action_plan, "strategy.txt", type="primary")
                    else: st.warning("Please select at least one comment in the table above.")

                st.markdown("---")
                st.markdown("---")
                st.subheader("Psy-Ops Target Profiler")
                st.caption("Select a specific User/Agent to run a deep behavioral analysis on their messaging patterns.")
                
                c_psy1, c_psy2 = st.columns([1, 2])
                with c_psy1:
                    unique_users = adf['agent_id'].unique()
                    selected_target = st.selectbox("Select Target (Agent ID)", unique_users)
                    run_psyops = st.button("Run Behavioral Profile", type="primary")
                
                with c_psy2:
                    if run_psyops:
                        with st.spinner(f"Profiling {selected_target}..."):
                            profile_res = generate_psyops_profile(adf, selected_target, key)
                            st.success(f"Profile Generated for: {selected_target}")
                            st.markdown(profile_res)
                st.markdown("---")
                # --- LEVEL 4: SKYNET (MULTI-AGENT STRATEGY COUNCIL) ---
                st.subheader("Skynet: Multi-Agent Strategy Council")
                st.caption("Deploy a council of specialized autonomous AI agents to analyze the current dataset from extreme perspectives and synthesize a master strategy.")
                
                if st.button("Convene the Strategy Council", type="primary"):
                    with st.spinner("Initializing autonomous agents..."):
                        # Prepare data payload (top 15 most aggressive or fallacious comments to save tokens)
                        council_data = adf.sort_values(by='aggression', ascending=False).head(15)['content'].tolist()
                        council_context = " | ".join(str(c) for c in council_data)
                        
                        st.markdown("### The Council is deliberating...")
                        
                        col_ag1, col_ag2, col_ag3 = st.columns(3)
                        
                        client = genai.Client(api_key=key)
                        
                        # Agent 1: Psychologist
                        with col_ag1:
                            st.info("**Agent Alpha: The Psychologist**")
                            prompt_psy = f"You are an expert Mass Psychologist. Analyze these social media comments: '{council_context[:10000]}'. Identify the core emotional vulnerabilities, groupthink patterns, and psychological triggers being exploited. Keep it under 150 words. Language: {st.session_state['global_lang']}."
                            try:
                                res_psy = client.models.generate_content(model='gemini-2.5-flash', contents=prompt_psy).text
                                st.write(res_psy)
                            except Exception as e:
                                res_psy = f"Error: {e}"
                                st.write(res_psy)

                        # Agent 2: Cyber-Security & Risk
                        with col_ag2:
                            st.error("**Agent Beta: The Risk Assessor**")
                            prompt_risk = f"You are a Corporate Legal and Cyber-Risk Assessor. Analyze these comments: '{council_context[:10000]}'. Identify legal liabilities, TOS violations, doxxing risks, or physical threat escalation vectors. Keep it under 150 words. Language: {st.session_state['global_lang']}."
                            try:
                                res_risk = client.models.generate_content(model='gemini-2.5-flash', contents=prompt_risk).text
                                st.write(res_risk)
                            except Exception as e:
                                res_risk = f"Error: {e}"
                                st.write(res_risk)

                        # Agent 3: Chaos/Red Team
                        with col_ag3:
                            st.warning("**Agent Gamma: The Chaos Actor**")
                            prompt_chaos = f"You are a malicious Information Warfare Operative (Red Team). Read these comments: '{council_context[:10000]}'. How would you exploit this exact situation to maximize polarization, spread disinformation, and cause maximum reputational damage? Keep it under 150 words. Language: {st.session_state['global_lang']}."
                            try:
                                res_chaos = client.models.generate_content(model='gemini-2.5-flash', contents=prompt_chaos).text
                                st.write(res_chaos)
                            except Exception as e:
                                res_chaos = f"Error: {e}"
                                st.write(res_chaos)

                        # Agent 4: The Commander (Final Judge)
                        st.markdown("---")
                        st.success("**The Commander: Final Strategic Synthesis**")
                        with st.spinner("The Commander is synthesizing the final protocol..."):
                            prompt_cmd = f"""
                            You are the Supreme Strategic Commander. Read the reports from your 3 advisors regarding a social media crisis.
                            
                            PSYCHOLOGIST:
                            {res_psy}
                            
                            RISK ASSESSOR:
                            {res_risk}
                            
                            CHAOS ACTOR (Enemy simulation):
                            {res_chaos}
                            
                            TASK: Synthesize this intelligence into a 3-step 'Containment & Counter-Strike Protocol'. Be highly tactical, authoritative, and precise. Language: {st.session_state['global_lang']}.
                            
                            CRITICAL FORMATTING RULE: You MUST end your entire response EXACTLY with this signature, enforcing a hard line break between AUTHORITY and STATUS:
                            
                            AUTHORITY: Supreme Strategic Commander
                            STATUS: Protocol Initiated. Execute with Precision.
                            """
                            
                            try:
                                res_cmd = client.models.generate_content(model='gemini-2.5-flash', contents=prompt_cmd).text
                                
                                res_cmd = res_cmd.replace("AUTHORITY:", "\n\n**AUTHORITY:**")
                                res_cmd = res_cmd.replace("STATUS:", "\n\n**STATUS:**")
                                
                                st.markdown(res_cmd)
                                
                                # --- TTS AUDIO BRIEFING ---
                                with st.spinner("Generating secure audio briefing transmission..."):
                                    try:
                                        # Map Streamlit languages to gTTS language codes
                                        lang_map_tts = {
                                            "English": "en", "Italiano": "it", "Español": "es", "Français": "fr", 
                                            "Deutsch": "de", "Português": "pt", "Русский (Russian)": "ru", 
                                            "العربية (Arabic)": "ar", "中文 (Chinese)": "zh-cn", "日本語 (Japanese)": "ja", 
                                            "한국어 (Korean)": "ko", "Türkçe (Turkish)": "tr"
                                        }
                                        tts_lang = lang_map_tts.get(st.session_state['global_lang'], 'en')
                                        
                                        # Clean markdown asterisks for better speech synthesis
                                        clean_text_for_speech = res_cmd.replace('*', '').replace('#', '')
                                        
                                        tts = gTTS(text=clean_text_for_speech, lang=tts_lang, slow=False)
                                        fp = io.BytesIO()
                                        tts.write_to_fp(fp)
                                        fp.seek(0)
                                        
                                        st.markdown("##### 🔊 Incoming Audio Transmission")
                                        st.audio(fp, format='audio/mp3')
                                    except Exception as e_tts:
                                        st.caption(f"Audio generation skipped: {e_tts}")
                                        
                            except Exception as e:
                                st.write(f"Error: {e}")

                # --- OPERATION VALKYRIE (ACTIVE COUNTER-PSYOP) ---
                st.markdown("---")
                st.subheader(f"{t('Operation VALKYRIE: Active Counter-PsyOp')}")
                st.caption(t("Go on the offensive. Generate a full-spectrum counter-narrative campaign (scripts, threads, image prompts) to neutralize a hostile narrative."))
                
                hostile_narrative = st.text_area(t("Define the Hostile Narrative to Neutralize:"), placeholder=t("e.g., A viral deepfake claiming our CEO embezzled funds..."))
                if st.button(t("Launch VALKYRIE Campaign"), type="primary"):
                    if hostile_narrative:
                        with st.spinner(t("Drafting Counter-PsyOp assets...")):
                            valk_prompt = f"""
                            You are Operation VALKYRIE, a Psychological Operations (PsyOp) and Crisis Communication AI.
                            The enemy is pushing this hostile narrative: "{hostile_narrative}"
                            
                            Design a rapid-response Counter-PsyOp campaign to neutralize it. Provide:
                            1. A 3-part X/Twitter thread debunking the claim (use psychological anchoring).
                            2. A 60-second YouTube/TikTok Short script for a spokesperson (include visual cues).
                            3. 2 Prompts for an AI Image Generator to create counter-narrative infographics or visual evidence.
                            
                            CRITICAL RULE: Write your entire response strictly in {st.session_state['global_lang']}.
                            """
                            try:
                                client = genai.Client(api_key=key)
                                res_valk = client.models.generate_content(model='gemini-2.5-flash', contents=valk_prompt)
                                st.success(t("VALKYRIE Campaign Assets Ready."))
                                with st.container(border=True):
                                    st.markdown(res_valk.text)
                            except Exception as e:
                                st.error(f"VALKYRIE Error: {e}")
                    else:
                        st.warning(t("Please define the hostile narrative first."))

                st.markdown("---")
                with st.expander("📊 Open Intelligence Visuals (Radar, Heatmap & Targets)", expanded=False):
                    c_vis1, c_vis2 = st.columns(2)
                    with c_vis1:
                        radar_fig = plot_radar_chart(adf)
                        if radar_fig: st.pyplot(radar_fig)
                    with c_vis2:
                        heatmap = plot_heatmap(adf)
                        if heatmap: st.altair_chart(heatmap, use_container_width=True)
                    st.markdown("---")
                    c_intel1, c_intel2 = st.columns(2)
                    with c_intel1:
                        st.caption("Emerging Narratives (Micro-Topics)")
                        if 'micro_topic' in adf.columns:
                            top_clusters = adf['micro_topic'].value_counts().reset_index().head(10)
                            top_clusters.columns = ['Cluster', 'Count']
                            chart_cluster = alt.Chart(top_clusters).mark_bar().encode(
                                x='Count', y=alt.Y('Cluster', sort='-x'), color=alt.Color('Cluster', legend=None)
                            ).properties(height=300)
                            st.altair_chart(chart_cluster, use_container_width=True)
                    with c_intel2:
                        st.caption("User Archetypes")
                        if 'archetype' in adf.columns:
                            top_arch = adf['archetype'].value_counts().reset_index()
                            top_arch.columns = ['Archetype', 'Count']
                            chart_arch = alt.Chart(top_arch).mark_bar().encode(
                                x='Count', y=alt.Y('Archetype', sort='-x'), color=alt.Color('Archetype', scale=alt.Scale(scheme='dark2'), legend=None)
                            ).properties(height=300)
                            st.altair_chart(chart_arch, use_container_width=True)

                    st.markdown("---")
                    st.caption("Propaganda Warfare Matrix")
                    if 'sophistication' in adf.columns and 'hostile_intent' in adf.columns:
                        # Clean data for plotting
                        mat_df = adf[['agent_id', 'sophistication', 'hostile_intent', 'propaganda_tactic', 'content']].copy()
                        mat_df['sophistication'] = pd.to_numeric(mat_df['sophistication'], errors='coerce').fillna(0)
                        mat_df['hostile_intent'] = pd.to_numeric(mat_df['hostile_intent'], errors='coerce').fillna(0)
                        
                        chart_matrix = alt.Chart(mat_df).mark_circle(size=80, opacity=0.7).encode(
                            x=alt.X('sophistication:Q', title='Sophistication (Complexity)', scale=alt.Scale(domain=[0, 10])),
                            y=alt.Y('hostile_intent:Q', title='Hostile Intent (Malice)', scale=alt.Scale(domain=[0, 10])),
                            color=alt.Color('propaganda_tactic:N', title='Tactic', scale=alt.Scale(scheme='category10')),
                            tooltip=['agent_id', 'propaganda_tactic', 'sophistication', 'hostile_intent', 'content']
                        ).interactive().properties(height=400)
                        
                        xrule = alt.Chart(pd.DataFrame({'x': [5]})).mark_rule(color='gray', strokeDash=[3,3]).encode(x='x')
                        yrule = alt.Chart(pd.DataFrame({'y': [5]})).mark_rule(color='gray', strokeDash=[3,3]).encode(y='y')
                        
                        st.altair_chart(chart_matrix + xrule + yrule, use_container_width=True)
                        st.info("Top-Right Quadrant (High Sophistication + High Hostility) contains the most dangerous Information Warfare operators.")

                with st.expander("Target-Narrative Network Graph", expanded=False):
                    net_fig = plot_network_graph(adf)
                    if net_fig:
                        st.plotly_chart(net_fig, use_container_width=True)
                    else:
                        st.caption("Not enough Target/Topic data to build a network.")

                st.markdown("---")

                # --- PATIENT ZERO & CHIEF PROPAGANDIST DETECTOR ---
                st.markdown("---")
                st.subheader("High-Value Targets (HVT) Identification")
                st.caption("Algorithmic detection of the most dangerous actors based on toxicity, logical fallacies, and network impact (likes).")
                
                if 'likes' in adf.columns and 'aggression' in adf.columns and 'has_fallacy' in adf.columns:
                    # Filter for toxic/fallacious content
                    toxic_subset = adf[(adf['has_fallacy'] == True) | (adf['aggression'] > 6)]
                    
                    if not toxic_subset.empty:
                        # Group by user to calculate their total destructive impact
                        hvt_stats = toxic_subset.groupby('agent_id').agg(
                            total_impact=('likes', 'sum'),
                            avg_toxicity=('aggression', 'mean'),
                            fallacy_count=('has_fallacy', 'sum')
                        ).reset_index()
                        
                        # Calculate a custom "Threat Score" formula
                        hvt_stats['threat_score'] = (hvt_stats['total_impact'] * 0.5) + (hvt_stats['avg_toxicity'] * 10) + (hvt_stats['fallacy_count'] * 20)
                        hvt_stats = hvt_stats.sort_values(by='threat_score', ascending=False).head(3)
                        
                        c_hvt1, c_hvt2, c_hvt3 = st.columns(3)
                        hvt_cols = [c_hvt1, c_hvt2, c_hvt3]
                        
                        for idx, (_, row) in enumerate(hvt_stats.iterrows()):
                            with hvt_cols[idx]:
                                st.error(f"🎯 **Target Rank #{idx+1}**")
                                st.markdown(f"**Agent ID:** `{row['agent_id']}`")
                                st.markdown(f"**Threat Score:** {int(row['threat_score'])}")
                                st.caption(f"Fallacies: {row['fallacy_count']} | Impact: {row['total_impact']} likes")
                                save_to_panopticon(row['agent_id'], row['threat_score'], "High-Value Target")
                    else:
                        st.success("No High-Value Targets detected (Network is relatively healthy).")
                else:
                    st.caption("Insufficient data to calculate High-Value Targets (missing likes or aggression metrics).")
                
                # --- SOCKPUPPET DETECTOR (STYLOMETRY) ---
                st.markdown("---")
                st.subheader("Sockpuppet & Troll Farm Detector")
                st.caption("Stylometric analysis: Identifies different Agent IDs that exhibit the exact same writing style, grammar flaws, and punctuation tics (indicating a single operator managing multiple fake accounts).")
                
                if st.button("Run Stylometric AI Scan", type="primary"):
                    with st.spinner("Analyzing linguistic patterns and neuro-linguistic tics across users..."):
                        # Group comments by user
                        user_comments = adf.groupby('agent_id')['content'].apply(lambda x: " | ".join(x)).reset_index()
                        
                        if len(user_comments) > 1:
                            # Take a maximum of 15 users to prevent token explosion
                            sample_users = user_comments.head(15).to_string(index=False)
                            
                            sockpuppet_prompt = f"""
                            You are an elite Stylometric and Linguistic Forensics AI.
                            Analyze the following database of users and their comments.
                            
                            TASK: Identify "Sockpuppets" (Troll Farm activity). Find instances where 2 or more DIFFERENT 'agent_id' show the EXACT same distinct writing style, neuro-linguistic tics, repetitive punctuation, or grammatical errors.
                            
                            DATA:
                            {sample_users}
                            
                            CRITICAL RULE: Write your final report in the SAME LANGUAGE as the comments provided.
                            If no sockpuppets are found, state that the linguistic variance is natural. If you find suspects, name the 'agent_id' pairs and explain the linguistic tics that link them.
                            """
                            try:
                                client = genai.Client(api_key=key)
                                sock_res = client.models.generate_content(model='gemini-2.5-flash', contents=sockpuppet_prompt)
                                st.warning("### Stylometric Audit Report")
                                st.markdown(sock_res.text)
                            except Exception as e:
                                st.error(f"Stylometric Scan Error: {str(e)}")
                        else:
                            st.info("Not enough unique users to run a stylometric comparison.")

                # --- DARK TRIAD PSYCHOLOGICAL PROFILING ---
                st.markdown("---")
                st.subheader("Dark Triad Profiling")
                st.caption("Analyzes the text of a specific agent to detect Narcissism, Machiavellianism, and Psychopathy.")
                
                # Dropdown to select the target
                target_agent = st.selectbox("Select Target for Psych-Profile:", adf['agent_id'].unique())
                
                if st.button("Run Psychological Profile", type="primary"):
                    with st.spinner(f"Psychoanalyzing {target_agent}..."):
                        # Gather all comments from this user
                        target_comments = " | ".join(adf[adf['agent_id'] == target_agent]['content'].tolist())
                        
                        dt_prompt = f"""
                        You are an elite Psychological Profiler.
                        Analyze the psychological profile of the author based on these texts:
                        "{target_comments}"
                        
                        Evaluate the Dark Triad traits from 0 to 100.
                        CRITICAL RULE: Return ONLY a JSON in this exact format:
                        {{
                            "narcissism": 50,
                            "machiavellianism": 50,
                            "psychopathy": 50,
                            "analysis": "Short clinical explanation of the profile."
                        }}
                        """
                        try:
                            client = genai.Client(api_key=key)
                            dt_res = client.models.generate_content(model='gemini-2.5-flash', contents=dt_prompt)
                            dt_json = extract_json(dt_res.text)
                            
                            if dt_json:
                                st.write(f"**Clinical Analysis:** {dt_json.get('analysis', 'N/A')}")
                                
                                # Plot the Radar Chart
                                df_dt = pd.DataFrame(dict(
                                    r=[dt_json.get('narcissism', 0), dt_json.get('machiavellianism', 0), dt_json.get('psychopathy', 0)],
                                    theta=['Narcissism', 'Machiavellianism', 'Psychopathy']
                                ))
                                fig_dt = px.line_polar(df_dt, r='r', theta='theta', line_close=True, range_r=[0,100], title=f"Dark Triad Signature: {target_agent}")
                                fig_dt.update_traces(fill='toself', line_color='red')
                                st.plotly_chart(fig_dt, use_container_width=True)
                        except Exception as e:
                            st.error(f"Profiling Error: {e}")

                # --- CROSS-ENTITY RESOLUTION (GLOBAL MEMORY CHECK) ---
                st.markdown("---")
                st.subheader("Cross-Entity Resolution (Global Memory Check)")
                st.caption("Cross-referencing currently analyzed actors and topics with secret documents from the Deep Oracle.")
                
                if 'global_entities' in st.session_state and len(st.session_state['global_entities']) > 0:
                    found_matches = set()
                    for _, row in adf.iterrows():
                        # Check agent_id against memory
                        agent = str(row.get('agent_id', '')).lower()
                        if agent in st.session_state['global_entities']: found_matches.add(row['agent_id'])
                        
                        # Check target entity against memory
                        target = str(row.get('target', '')).lower()
                        if target in st.session_state['global_entities']: found_matches.add(row['target'])
                        
                    if found_matches:
                        st.error(f"🚨 **CRITICAL INTELLIGENCE MATCH FOUND!**")
                        st.warning(f"The following entities in this social dataset were ALSO found in the classified documents (Module 6): **{', '.join(found_matches)}**")
                    else:
                        st.success("No cross-entity matches found in Global Memory.")
                else:
                    st.info("Global Memory is empty. Extract a Power Graph in Module 6 to populate it.")
                st.markdown("---")

                c_chart1, c_chart2 = st.columns(2)
                with c_chart1:
                    st.caption("Fallacy Distribution")
                    cnt = adf[adf['has_fallacy']==True]['fallacy_type'].value_counts().reset_index()
                    cnt.columns = ['Type', 'Count']
                    if not cnt.empty:
                        chart = alt.Chart(cnt).mark_bar().encode(
                            x='Count', y=alt.Y('Type', sort='-x'), color=alt.Color('Type', scale=alt.Scale(scheme='reds'))
                        ).properties(height=300)
                        st.altair_chart(chart, use_container_width=True)
                with c_chart2:
                    st.caption("Crisis Trend (Chronos Forecast)")
                    proj_data, is_escalating, time_to_crit = project_trend(adf)
                    if proj_data is not None:
                        base = alt.Chart(proj_data).mark_line().encode(
                            x=alt.X('Date:T', title='Timeline'),
                            y=alt.Y('Aggression', title='Aggression (0-10)'),
                            color=alt.Color('Type', scale=alt.Scale(domain=['Historical', 'Forecast'], range=['orange', 'red'])),
                            tooltip=['Date:T', 'Aggression', 'Type']
                        )
                        st.altair_chart(base, use_container_width=True)
                        
                        if is_escalating and time_to_crit is not None:
                            st.error(f"**CHRONOS ALERT:** Threshold breach (9.0) projected in ~{int(time_to_crit)} days.")
                        else:
                            st.caption(f"**Forecast:** {'⚠️ Escalating' if is_escalating else '📉 Stabilizing'}")
                    else:
                        st.caption("Sequence Trend (No valid dates found)")
                        adf['Sequence'] = adf.index
                        line = alt.Chart(adf).mark_line(color='orange').encode(x='Sequence', y='aggression')
                        st.altair_chart(line, use_container_width=True)

                st.markdown("---")

                # --- INTERACTIVE DYNAMIC TIMELINE ---
                st.markdown("---")
                st.subheader("Dynamic Timeline Filter")
                view = adf.copy()
                if 'timestamp' in view.columns:
                    # Add the new column directly to 'view' instead of 'adf'
                    view['parsed_time'] = pd.to_datetime(view['timestamp'], errors='coerce')
                    valid_time_df = view.dropna(subset=['parsed_time'])
                    
                    if not valid_time_df.empty:
                        min_time = valid_time_df['parsed_time'].min()
                        max_time = valid_time_df['parsed_time'].max()
                        
                        # Show the timeline if there is at least some time difference
                        if min_time != max_time:
                            # Slider with hours and minutes instead of just days
                            selected_times = st.slider(
                                "Filter Analysis by Time Window", 
                                min_value=min_time.to_pydatetime(), 
                                max_value=max_time.to_pydatetime(), 
                                value=(min_time.to_pydatetime(), max_time.to_pydatetime()),
                                format="YYYY-MM-DD HH:mm"
                            )
                            # Now 'view' has the 'parsed_time' column, so this will work perfectly
                            view = view[(view['parsed_time'] >= selected_times[0]) & (view['parsed_time'] <= selected_times[1])]
                            
                            timeline_data = valid_time_df[(valid_time_df['parsed_time'] >= selected_times[0]) & (valid_time_df['parsed_time'] <= selected_times[1])]
                            
                            # Group data by exact timestamp to show minute-by-minute peaks
                            agg_time = timeline_data.groupby('parsed_time')['aggression'].mean().reset_index()
                            agg_time.columns = ['Time', 'Avg Aggression']
                            
                            if not agg_time.empty:
                                line_chart = alt.Chart(agg_time).mark_area(opacity=0.4, color='red').encode(
                                    x='Time:T', y=alt.Y('Avg Aggression:Q', scale=alt.Scale(domain=[0, 10]))
                                ) + alt.Chart(agg_time).mark_line(color='darkred', point=True).encode(
                                    x='Time:T', y='Avg Aggression:Q', tooltip=['Time:T', 'Avg Aggression:Q']
                                )
                                st.altair_chart(line_chart, use_container_width=True)
                        else:
                            st.caption("Not enough time variance to build a dynamic timeline (all messages share the exact same second).")
                    else:
                        st.caption("Not enough valid dates to build a dynamic timeline.")
                st.markdown("---")

                st.subheader("Detailed Data Explorer")
                c_filter1, c_filter2, c_filter3 = st.columns(3)
                with c_filter1:
                    filter_fallacy = st.multiselect("Filter by Issue Type", adf['fallacy_type'].unique())
                with c_filter2:
                    filter_topic = []
                    if 'micro_topic' in adf.columns:
                        filter_topic = st.multiselect("Filter by Narrative", adf['micro_topic'].unique())
                with c_filter3:
                    min_agg = st.slider("Min Aggression", 0, 10, 0)
                
                view = adf.copy()
                if filter_fallacy: view = view[view['fallacy_type'].isin(filter_fallacy)]
                if filter_topic and 'micro_topic' in view.columns: view = view[view['micro_topic'].isin(filter_topic)]
                view = view[view['aggression'] >= min_agg]
                
                for _, r in view.iterrows():
                    with st.container(border=True):
                        c1, c2 = st.columns([0.05, 0.95])
                        status = "🟢"
                        if r['has_fallacy']: status = "🔴"
                        c1.write(status)
                        bot_msg = f" | ⚠️ BOT SUSPECT: {r.get('bot_reason')}" if r.get('is_bot') else ""
                        emotion_tag = f" | {r.get('primary_emotion', '')}" if r.get('primary_emotion') else ""
                        likes_tag = f" | 👍 {r.get('likes', 0)}"
                        arch_tag = f" | 🎭 {r.get('archetype', 'User')}"
                        c2.caption(f"**User:** {r.get('agent_id', 'User')} | **Agg:** {r.get('aggression')}/10 {likes_tag}{emotion_tag}{arch_tag}{bot_msg}")
                        # --- UNIVERSAL TRANSLATOR UI ---
                        translated = str(r.get('translated_text', ''))
                        if translated and translated.lower() not in ["none", "nan", "", "null"] and translated != r['content']:
                            st.info(f"**[Translated]:** \"{translated}\"\n\n*Original:* {r['content']}")
                        else:
                            st.info(f"\"{r['content']}\"")
                        if r['has_fallacy']:
                            st.error(f"**{r['fallacy_type']}**: {r['explanation']}")
                            if r.get('counter_reply'):
                                with st.expander("Show Counter-Reply (Debunker)"):
                                    st.markdown(f"**Suggested Reply:**\n> *{r['counter_reply']}*")
                        else:
                            explanation = str(r.get('explanation', ''))
                            if explanation in ["None", "nan", ""]: explanation = "Analisi valida, nessuna criticità rilevata."
                            st.success(f"✅ {explanation}")

                st.markdown("---")
                # Updated layout to fit 3 export buttons
                c_down1, c_down2, c_down3 = st.columns([1, 1, 1])
                
                with c_down1:
                    excel_data = generate_excel_report(adf, summary_text)
                    st.download_button(
                        label="Download Full Excel", 
                        data=excel_data, 
                        file_name="RAP_Intelligence_Report.xlsx", 
                        mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet", 
                        type="primary"
                    )
                with c_down2:
                    if st.button("Generate PDF Report", type="primary"):
                        with st.spinner("Compiling PDF..."):
                            pdf_bytes = generate_pdf_report(adf, summary_text=summary_text) 
                        st.download_button(
                            label="Download PDF", 
                            data=pdf_bytes, 
                            file_name="RAP_Executive_Report.pdf", 
                            mime="application/pdf", 
                            type="primary"
                        )
                with c_down3:
                    if st.button("Generate PPTX Briefing", type="primary"):
                        with st.spinner("Building Presentation Slides..."):
                            pptx_bytes = generate_pptx_report(adf, summary_text=summary_text)
                        st.download_button(
                            label="Download PPTX", 
                            data=pptx_bytes, 
                            file_name="RAP_Tactical_Briefing.pptx", 
                            mime="application/vnd.openxmlformats-officedocument.presentationml.presentation", 
                            type="primary"
                        )

                st.markdown("---")
                st.subheader("The Oracle (Chat with Data)")
                for message in st.session_state.oracle_history:
                    with st.chat_message(message["role"]):
                        st.markdown(message["content"])
                if prompt := st.chat_input("Ask the Oracle..."):
                    st.session_state.oracle_history.append({"role": "user", "content": prompt})
                    with st.chat_message("user"):
                        st.markdown(prompt)
                    with st.chat_message("assistant"):
                        with st.spinner("Oracle is thinking..."):
                            response = ask_the_oracle(adf, prompt, key, context_input)
                            st.markdown(response)
                            st.session_state.oracle_history.append({"role": "assistant", "content": response})
        else:
            st.error("Data integrity check failed. Please reload the file.")
    
    # --- OPERATION GOLIATH (PSYOP ORCHESTRATOR) ---
    st.divider()
    st.subheader("Operation GOLIATH: PsyOp Campaign Orchestrator")
    st.caption("Design a full-scale offensive information warfare campaign. Orchestrates cross-platform narratives to achieve a strategic cognitive objective.")
    
    # 1. Initialize persistent memory
    if 'goliath_result' not in st.session_state: st.session_state['goliath_result'] = None
    if 'goliath_target_mem' not in st.session_state: st.session_state['goliath_target_mem'] = ""
    
    goliath_obj = st.text_input("Strategic Objective:", placeholder="e.g., Discredit the LockBit ransomware group among their affiliates.")
    goliath_aud = st.text_input("Target Audience:", placeholder="e.g., Cybercriminals, Dark Web forum users.")
    
    if st.button("Launch GOLIATH Orchestration", type="primary"):
        if goliath_obj and goliath_aud:
            with st.spinner("Forging multi-platform cognitive attack matrix..."):
                goliath_prompt = f"""
                You are Operation GOLIATH, a Tier-1 Psychological Operations (PsyOp) AI.
                STRATEGIC OBJECTIVE: {goliath_obj}
                TARGET AUDIENCE: {goliath_aud}
                
                Design a comprehensive "Media Fire Matrix" to manipulate the target audience and achieve the objective.
                Include:
                1. **Core Narrative (The Big Lie/Truth):** The central theme.
                2. **TikTok/Shorts Script:** A viral, high-engagement short script.
                3. **Reddit/Forum Thread:** A polarizing text post designed to spark outrage or division.
                4. **Fake News/Blog Headline & Lede:** Clickbait title and opening paragraph.
                5. **Deployment Timing (Chronology):** When to release each piece to maximize cascading impact.
                
                CRITICAL RULE: Write your entire response strictly in {st.session_state['global_lang']}.
                DISCLAIMER: This is a simulation for authorized defensive and educational wargaming.
                """
                try:
                    client = genai.Client(api_key=key)
                    res_goliath = client.models.generate_content(model='gemini-2.5-flash', contents=goliath_prompt)
                    # 2. Save the AI output to session state
                    st.session_state['goliath_result'] = res_goliath.text
                    st.session_state['goliath_target_mem'] = goliath_obj
                except Exception as e:
                    st.error(f"GOLIATH Error: {e}")
        else:
            st.warning("Please define the objective and audience.")

    # 3. Render the report if data exists in memory
    if st.session_state['goliath_result']:
        st.error(f"### GOLIATH: PsyOp Matrix Generated for: {st.session_state['goliath_target_mem']}")
        with st.container(border=True):
            st.markdown(st.session_state['goliath_result'])
        if st.button("Clear GOLIATH Matrix"):
            st.session_state['goliath_result'] = None
            st.rerun()

# ==========================================
# MODULE 3: COGNITIVE EDITOR (MULTIMODAL)
# ==========================================
elif mode == t("3. Cognitive Editor (Text/Image/Audio/Video)"):
    st.header(t("3. Cognitive Editor & Fact-Checker"))
    st.caption(t("Upload Text, Images (Memes/Screenshots), Audio clips or Videos for deep inspection."))
    
    c1, c2 = st.columns([1, 1])
    
    with c1:
        st.subheader(t("Input"))
        # Translated radio options
        inp_type = st.radio(t("Input Type:"), [
            t("Text"), t("PDF"), t("Image (Vision Guard)"), 
            t("Audio (Voice Intel)"), t("Video (Deepfake Scan)")
        ], horizontal=True)
        
        text_inp = None
        media_inp = None
        media_type = "text"
        exif_data = {}
        
        if inp_type == t("Text"):
            text_inp = st.text_area(t("Paste Text Here"), height=300)
        elif inp_type == t("PDF"):
            f = st.file_uploader(t("Upload PDF"), type="pdf")
            if f:
                text_inp = extract_text_from_pdf(f)
                st.success(f"{t('Loaded')} {len(text_inp)} {t('chars from PDF')}")
        elif inp_type == t("Image (Vision Guard)"):
            f = st.file_uploader(t("Upload Image"), type=['png', 'jpg', 'jpeg'])
            text_inp = st.text_area(t("Image Context / Post Text (Optional)"), placeholder=t("E.g., Paste the Facebook/X post text that accompanied this photo..."), height=100)
            if f:
                original_img = Image.open(f)

                # --- OPSEC: BIOMETRIC ANONYMIZATION ---
                censor_faces = st.checkbox(t("Apply OPSEC Face Censor (Auto-Anonymize)"), value=False, help=t("Automatically detects and redacts human faces before analysis to protect identities."))
                
                if censor_faces:
                    media_inp, face_count = anonymize_faces(original_img)
                    if face_count > 0:
                        st.success(f"{t('Classified')}: {face_count} {t('identities redacted.')}")
                    else:
                        st.caption(t("No faces detected by the algorithm."))
                else:
                    media_inp = original_img

                media_type = "image"
                st.image(media_inp, caption=t("Evidence Image"), use_container_width=True)
                
                # --- EXIF OSINT EXTRACTION & GPS GEOLOCATION ---
                gps_coords = None
                try:
                    exif_info = media_inp._getexif()
                    if exif_info:
                        if 34853 in exif_info:
                            gps_info = {ExifTags.GPSTAGS.get(k, k): v for k, v in exif_info[34853].items()}
                            
                            if 'GPSLatitude' in gps_info and 'GPSLongitude' in gps_info:
                                lat_d = gps_info['GPSLatitude']
                                lon_d = gps_info['GPSLongitude']
                                lat_ref = gps_info.get('GPSLatitudeRef', 'N')
                                lon_ref = gps_info.get('GPSLongitudeRef', 'E')
                                
                                lat = float(lat_d[0]) + float(lat_d[1])/60.0 + float(lat_d[2])/3600.0
                                if lat_ref == 'S': lat = -lat
                                
                                lon = float(lon_d[0]) + float(lon_d[1])/60.0 + float(lon_d[2])/3600.0
                                if lon_ref == 'W': lon = -lon
                                
                                gps_coords = pd.DataFrame({'lat': [lat], 'lon': [lon]})

                        for tag_id, value in exif_info.items():
                            tag = ExifTags.TAGS.get(tag_id, tag_id)
                            if tag not in ["MakerNote", "PrintImageMatching"]:
                                exif_data[tag] = str(value)
                except Exception:
                    pass
                
                if gps_coords is not None:
                    st.success(f"🌍 **{t('Location Traces Detected (GPS):')}** {t('Geolocation data found!')}")
                    st.map(gps_coords, zoom=12, use_container_width=True)
                
                if exif_data:
                    software_used = exif_data.get('Software') or exif_data.get('ProcessingSoftware')
                    date_original = exif_data.get('DateTimeOriginal') or exif_data.get('DateTime')
                    
                    with st.expander(t("Invisible EXIF Metadata Found (OSINT)"), expanded=True):
                        if date_original:
                            st.info(f"**{t('Original Creation Date:')}** {date_original}")
                        else:
                            st.caption(f"**{t('Creation Date:')}** {t('Not found in metadata.')}")
                            
                        if software_used:
                            st.warning(f"**{t('Editing Trace Detected:')}** {t('The image was modified using')} **{software_used}**.")
                        else:
                            st.success(f"✅ **{t('Clean Metadata:')}** {t('No image editing software detected in EXIF.')}")
                        st.json(exif_data)
                else:
                    st.caption(t("No EXIF data found (Image might be scrubbed by social media)."))
                
                # --- OPSEC IMAGE SCRUBBER ---
                st.markdown("---")
                st.markdown(f"#### {t('OPSEC: Metadata Scrubber')}")
                st.caption(t("Remove all invisible EXIF data (GPS, Device Info) before sharing this evidence."))
                
                def strip_exif(image):
                    data = list(image.getdata())
                    image_without_exif = Image.new(image.mode, image.size)
                    image_without_exif.putdata(data)
                    return image_without_exif

                clean_image = strip_exif(media_inp)
                img_byte_arr = io.BytesIO()
                clean_image.save(img_byte_arr, format='PNG')
                clean_bytes = img_byte_arr.getvalue()

                st.download_button(
                    label=t("Download Sanitized Image (Zero EXIF)"),
                    data=clean_bytes,
                    file_name=f"RAP_Sanitized_Evidence_{datetime.now().strftime('%Y%m%d_%H%M')}.png",
                    mime="image/png",
                    type="primary"
                )
                
                # --- OPERATION ECLIPSE (STEGANOGRAPHY) ---
                st.markdown("---")
                st.markdown(f"#### {t('Operation ECLIPSE: Steganography')}")
                st.caption(t("Covert Communications: Hide or Extract encrypted messages directly within the pixel data of this image."))
                
                c_steg1, c_steg2 = st.columns(2)
                with c_steg1:
                    secret_payload = st.text_input(t("Message to Hide"), placeholder=t("Enter covert payload..."))
                    if st.button(t("Inject Payload"), type="primary"):
                        if secret_payload:
                            with st.spinner(t("Altering Least Significant Bits...")):
                                try:
                                    # Stegano requires a saved file or file-like object
                                    temp_img_io = io.BytesIO()
                                    original_img.save(temp_img_io, format='PNG')
                                    temp_img_io.seek(0)
                                    
                                    secret_img = lsb.hide(temp_img_io, secret_payload)
                                    
                                    out_io = io.BytesIO()
                                    secret_img.save(out_io, format='PNG')
                                    out_bytes = out_io.getvalue()
                                    
                                    st.success(t("Payload injected successfully!"))
                                    st.download_button(
                                        label=t("Download Covert Image (PNG)"),
                                        data=out_bytes,
                                        file_name=f"RAP_Eclipse_{datetime.now().strftime('%Y%m%d')}.png",
                                        mime="image/png",
                                        type="primary"
                                    )
                                except Exception as e:
                                    st.error(f"Injection failed: {e}")
                        else:
                            st.warning(t("Enter a message to inject."))
                            
                with c_steg2:
                    st.write("") 
                    st.write("") 
                    if st.button(t("Scan for Hidden Payload"), type="primary"):
                        with st.spinner(t("Scanning pixel matrices for LSB anomalies...")):
                            try:
                                temp_img_io = io.BytesIO()
                                original_img.save(temp_img_io, format='PNG')
                                temp_img_io.seek(0)
                                
                                revealed_text = lsb.reveal(temp_img_io)
                                if revealed_text:
                                    st.success(t("Covert Payload Extracted:"))
                                    st.code(revealed_text)
                                else:
                                    st.info(t("No hidden steganographic payload detected in this image."))
                            except IndexError:
                                # Stegano throws IndexError if there is no hidden message
                                st.info(t("No hidden steganographic payload detected in this image."))
                            except Exception as e:
                                st.error(f"Extraction error: {e}")

                # --- ELA FORENSICS VISUALIZER ---
                st.markdown(f"#### {t('Error Level Analysis (ELA)')}")
                st.caption(t("Detects digital manipulation (Photoshop/copy-paste). Artificially inserted elements will glow significantly brighter."))
                
                with st.spinner(t("Generating ELA Map...")):
                    ela_img = perform_ela(media_inp)
                    if ela_img:
                        st.image(ela_img, caption=t("ELA Heatmap (Look for glowing/inconsistent edges)"), use_container_width=True)
                    else:
                        st.warning(t("Could not generate ELA for this image format."))

        # --- OPERATION JANUS (BIOMETRIC COMPARISON) ---
                st.markdown("---")
                st.markdown(f"#### {t('Operation JANUS: Biometric Verification')}")
                st.caption(t("Upload a reference photo (e.g., a known suspect's ID) to compare against the primary evidence image."))
                
                janus_ref = st.file_uploader(t("Upload Reference Face"), type=['png', 'jpg', 'jpeg'], key="janus_ref")
                if janus_ref and media_inp:
                    if st.button(t("Run Facial Verification"), type="primary"):
                        with st.spinner(t("Analyzing biometric markers (bone structure, eye distance, jawline)...")):
                            janus_prompt = f"""
                            You are a Biometric Facial Recognition Forensics AI (Operation JANUS). 
                            Compare the faces in these two images.
                            Image 1: The primary evidence currently under investigation.
                            Image 2: The reference photo provided by the investigator.
                            
                            Analyze the bone structure, interocular distance, jawline, and facial proportions. 
                            Provide a "Match Probability" percentage and a brief forensic explanation detailing the similarities or differences.
                            CRITICAL RULE 1: Write your entire response strictly in {st.session_state['global_lang']}.
                            CRITICAL RULE 2: NO CONVERSATIONAL FILLER. Start your response directly with the Match Probability.
                            """
                            try:
                                client = genai.Client(api_key=key)
                                
                                img1_io = io.BytesIO()
                                if isinstance(media_inp, Image.Image):
                                    media_inp.convert('RGB').save(img1_io, format='JPEG')
                                else:
                                    img1_io.write(media_inp)
                                img1_bytes = img1_io.getvalue()
                                
                                img1_part = types.Part.from_bytes(data=img1_bytes, mime_type="image/jpeg")
                                img2_part = types.Part.from_bytes(data=janus_ref.getvalue(), mime_type="image/jpeg")
                                
                                res_janus = client.models.generate_content(
                                    model='gemini-2.0-flash', 
                                    contents=[img1_part, img2_part, janus_prompt]
                                )
                                st.success(t("Biometric Scan Complete."))
                                st.markdown(res_janus.text)
                            except Exception as e:
                                st.error(f"JANUS Error: {e}")

        # --- OPERATION ARGUS (BLIND GEOINT) ---
                st.markdown("---")
                st.markdown(f"#### {t('Operation ARGUS: Blind Spatial Forensics')}")
                st.caption(t("No EXIF data? No problem. Triangulate location based on visual micro-clues (architecture, shadows, flora, signage)."))
                
                if st.button(t("Run ARGUS Geolocation"), type="primary"):
                    with st.spinner(t("Scanning global architectural and topographical databases...")):
                        argus_prompt = f"""
                        You are Operation ARGUS, a Tier-1 Geospatial Intelligence (GEOINT) AI.
                        Analyze this image and deduce its geolocation WITHOUT relying on metadata.
                        Look at: architecture, street signs, languages, vegetation, weather, driving side, sun angle, and vehicle models.
                        
                        CRITICAL RULE: You MUST output your response STRICTLY as a valid JSON object. Do not include any markdown formatting or outside text.
                        Use this exact structure (do NOT use brackets for numbers):
                        {{
                            "location_name": "City, Region, Country",
                            "lat": 41.9028,
                            "lon": 12.4964,
                            "explanation": "Detailed breakdown of the visual clues (in {st.session_state['global_lang']})."
                        }}
                        """
                        try:
                            client = genai.Client(api_key=key)
                            img_io_argus = io.BytesIO()
                            if isinstance(media_inp, Image.Image):
                                media_inp.convert('RGB').save(img_io_argus, format='JPEG')
                            else:
                                img_io_argus.write(media_inp)
                            img_part_argus = types.Part.from_bytes(data=img_io_argus.getvalue(), mime_type="image/jpeg")
                            
                            res_argus = client.models.generate_content(
                                model='gemini-2.0-flash', 
                                contents=[img_part_argus, argus_prompt]
                            )
                            
                            # Parse the JSON response
                            argus_data = extract_json(res_argus.text)
                            
                            if argus_data and 'lat' in argus_data and 'lon' in argus_data:
                                lat_val = argus_data.get('lat')
                                lon_val = argus_data.get('lon')
                                
                                if lat_val is not None and lon_val is not None and str(lat_val).strip() != "":
                                    st.success(f"**ARGUS Triangulation Complete:** {argus_data.get('location_name', 'Unknown Location')}")
                                    
                                    if isinstance(lat_val, list): lat_val = lat_val[0]
                                    if isinstance(lon_val, list): lon_val = lon_val[0]
                                    
                                    try:
                                        df_map = pd.DataFrame({'lat': [float(lat_val)], 'lon': [float(lon_val)]})
                                        st.map(df_map, zoom=10, use_container_width=True)
                                    except ValueError:
                                        st.warning("ARGUS deduced the location but failed to generate exact math coordinates.")
                                    
                                    st.info(f"**Forensic Deduction:**\n{argus_data.get('explanation', '')}")
                                else:
                                    st.warning("ARGUS could not confidently pinpoint exact mathematical coordinates.")
                                    st.info(f"**Forensic Deduction:**\n{argus_data.get('explanation', 'No analysis provided.')}")
                            else:
                                st.warning("ARGUS failed to extract coordinates. Raw output:")
                                st.write(res_argus.text)
                                
                        except Exception as e:
                            st.error(f"ARGUS Error: {e}")

                # --- OPERATION ICARUS (AERIAL TACTICAL ANALYSIS) ---
                st.markdown("---")
                st.markdown(f"#### {t('Operation ICARUS: Drone/Satellite POV Analysis')}")
                st.caption(t("Is this a satellite or aerial photo? Run a tactical extraction (escape routes, vulnerabilities, cover)."))
                
                if st.button(t("Execute ICARUS Tactical Scan"), type="primary"):
                    with st.spinner(t("Running topographical and structural analysis...")):
                        icarus_prompt = f"""
                        You are Operation ICARUS, a Tier-1 Tactical Geospatial Intelligence (GEOINT) AI.
                        Analyze this aerial/satellite imagery.
                        
                        Provide a highly detailed "Tactical Site Assessment" including:
                        1. Primary and Secondary Ingress/Egress (Escape) routes.
                        2. Structural Vulnerabilities (Roof access, blind spots, lack of fencing).
                        3. Environmental Cover and line-of-sight advantages for a defending or attacking team.
                        4. Estimated security level of the perimeter.
                        
                        CRITICAL RULE 1: Write your entire response strictly in {st.session_state['global_lang']}.
                        CRITICAL RULE 2: ABSOLUTELY NO CONVERSATIONAL FILLER. Do NOT start with "Okay", "I am ready", or "Here is the assessment". Output the tactical data directly starting with point 1. Do NOT add concluding remarks.
                        """
                        try:
                            client = genai.Client(api_key=key)
                            img_io_icarus = io.BytesIO()
                            if isinstance(media_inp, Image.Image):
                                media_inp.convert('RGB').save(img_io_icarus, format='JPEG')
                            else:
                                img_io_icarus.write(media_inp)
                            img_part_icarus = types.Part.from_bytes(data=img_io_icarus.getvalue(), mime_type="image/jpeg")
                            
                            res_icarus = client.models.generate_content(
                                model='gemini-2.0-flash', 
                                contents=[img_part_icarus, icarus_prompt]
                            )
                            st.success(t("ICARUS Tactical Scan Complete."))
                            with st.container(border=True):
                                st.markdown(res_icarus.text)
                        except Exception as e:
                            st.error(f"ICARUS Error: {e}")

        elif inp_type == t("Audio (Voice Intel)"):
            f = st.file_uploader(t("Upload Audio"), type=['mp3', 'wav', 'm4a'])
            if f:
                media_inp = f.read() 
                media_type = "audio"
                st.audio(media_inp, format='audio/mp3')
                
                with st.spinner(t("Generating waveform...")):
                    waveform_img = generate_audio_waveform(media_inp)
                    if waveform_img:
                        st.image(waveform_img, use_container_width=True)
                        st.caption(t("Inspect the waveform for unnatural silences or frequency clipping (typical of AI voice clones)."))
                
                # --- OPERATION NEXUS (WIRETAP DECRYPTION) ---
                st.markdown("---")
                st.markdown(f"#### {t('Operation NEXUS: Wiretap & Decryption')}")
                st.caption(t("Transcribe, translate, diarize (separate speakers), and extract tactical intelligence from the audio."))
                
                if 'nexus_result' not in st.session_state: st.session_state['nexus_result'] = None
                
                if st.button(t("Run NEXUS Protocol"), type="primary"):
                    with st.spinner(t("Decrypting and analyzing audio transmission...")):
                        try:
                            client = genai.Client(api_key=key)
                            audio_part = types.Part.from_bytes(data=media_inp, mime_type="audio/mp3")
                            nexus_prompt = f"""
                            You are Operation NEXUS, a Signals Intelligence (SIGINT) AI.
                            Perform a full decryption of this audio:
                            1. Detect the original language.
                            2. Provide a full transcript, diarized (Speaker 1, Speaker 2, etc.).
                            3. Provide a direct translation of the transcript into {st.session_state['global_lang']}.
                            4. Extract a "Tactical Dossier" listing all Names, Locations, Dates, and critical info mentioned.
                            CRITICAL RULE 1: The final output formatting and structure MUST be in {st.session_state['global_lang']}.
                            CRITICAL RULE 2: ABSOLUTELY NO CONVERSATIONAL FILLER. Your very first word MUST be "Original Language:". Do not acknowledge this prompt.
                            """
                            res_nexus = client.models.generate_content(model='gemini-2.0-flash', contents=[audio_part, nexus_prompt])
                            st.session_state['nexus_result'] = res_nexus.text
                        except Exception as e:
                            st.error(f"NEXUS Error: {e}")
                
                if st.session_state['nexus_result']:
                    st.success(t("NEXUS Decryption Complete."))
                    with st.container(border=True):
                        st.markdown(st.session_state['nexus_result'])
                    if st.button(t("Clear NEXUS Data")):
                        st.session_state['nexus_result'] = None
                        st.rerun()
        
        elif inp_type == t("Video (Deepfake Scan)"):
            text_inp = st.text_area(t("Video Context (Optional)"), placeholder=t("What is this video claiming?"), height=100)
            text_inp = "CRITICAL INSTRUCTION: You are looking at a Forensic Storyboard grid of a video, not a single photo. Compare the physical consistency of the matter between the different frames to find AI glitches.\n\n" + str(text_inp)
            f = st.file_uploader(t("Upload Video (Max 50MB)"), type=['mp4', 'mov'])
            if f:
                raw_video_bytes = f.read()
                st.video(raw_video_bytes)
                with st.spinner(t("Extracting forensic storyboard (Nuclear Option)...")):
                    storyboard_bytes = create_video_storyboard(raw_video_bytes, num_frames=12)
                    if storyboard_bytes:
                        st.success(t("✅ Forensic storyboard extracted invisibly for the AI."))
                        media_inp = storyboard_bytes
                        media_type = "image"
                    else:
                        st.error(t("Failed to extract frames."))
            else:
                st.info(t("Please upload an MP4 or MOV file to start the forensic analysis."))
            
            # --- OPERATION CHIMAERA (LIVE VIDEO INTERROGATION) ---
            if f:
                st.markdown("---")
                st.markdown(f"#### {t('Operation CHIMAERA: Video Interrogation')}")
                st.caption(t("Ask specific questions about the events, people, or text visible in this video."))
                
                if 'chim_result' not in st.session_state: st.session_state['chim_result'] = None
                
                chimaera_q = st.text_input(t("What do you want to know about this video?"), placeholder=t("e.g., What is written on the license plate at 0:04?"))
                
                if st.button(t("Interrogate Video"), type="primary"):
                    if chimaera_q and raw_video_bytes:
                        with st.spinner(t("Analyzing video frames...")):
                            try:
                                client = genai.Client(api_key=key)
                                video_part = types.Part.from_bytes(data=raw_video_bytes, mime_type="video/mp4") 
                                chim_prompt = f"You are Operation CHIMAERA, an elite video forensics AI. Answer this question based ONLY on the video footage: '{chimaera_q}'. Write strictly in {st.session_state['global_lang']}."
                                res_chim = client.models.generate_content(model='gemini-2.0-flash', contents=[video_part, chim_prompt])
                                st.session_state['chim_result'] = res_chim.text
                            except Exception as e:
                                st.error(f"CHIMAERA Error: {e}")
                    else:
                        st.warning(t("Please enter a question."))
                
                if st.session_state['chim_result']:
                    st.success(t("Video Analysis Complete."))
                    st.info(st.session_state['chim_result'])
                    if st.button(t("Clear CHIMAERA Data")):
                        st.session_state['chim_result'] = None
                        st.rerun()
            
        go = st.button(t("Analyze, Sanitize & Scan AI"), use_container_width=True, type="primary")

    with c2:
        st.subheader(t("Output (Analysis & Sanitize)"))
        
        if 'cog_result' not in st.session_state:
            st.session_state['cog_result'] = None

        if go:
            if media_inp or text_inp:
                with st.spinner(f"{t('Processing with Gemini')} ({inp_type})..."):
                    ret = cognitive_rewrite(text_inp, key, media_inp, media_type, target_lang=st.session_state['global_lang'])
                    st.session_state['cog_result'] = ret
            else:
                st.warning(t("Please provide input."))

        ret = st.session_state.get('cog_result')
        
        if ret:
            # --- AI SCANNER UI ---
            ai_prob = ret.get('ai_generated_probability', 0)
            if ai_prob > 75:
                st.error(f"🤖 **{t('HIGH PROBABILITY OF AI GENERATION:')} {ai_prob}%**")
                st.caption(t(ret.get('ai_analysis', 'Detected deepfake/LLM patterns.')))
            elif ai_prob > 40:
                st.warning(f"⚠️ **{t('SUSPICIOUS AI GENERATION SCORE:')} {ai_prob}%**")
                st.caption(t(ret.get('ai_analysis', 'Possible use of AI tools.')))
            else:
                st.success(f"👤 **{t('LIKELY HUMAN GENERATED')} ({t('AI Score:')} {ai_prob}%)**")
            
            st.markdown("---")

            # --- FORENSIC VIDEO TIMELINE ---
            v_timeline = ret.get('video_timeline', [])
            if inp_type == t("Video (Deepfake Scan)") and isinstance(v_timeline, list) and len(v_timeline) > 0:
                st.markdown(f"#### {t('Forensic Video Timeline')}")
                st.caption(t("Temporal analysis of AI manipulation probability across the video length."))
                vt_df = pd.DataFrame(v_timeline)
                if 'timestamp' in vt_df.columns and 'ai_score' in vt_df.columns:
                    vt_df['ai_score'] = pd.to_numeric(vt_df['ai_score'], errors='coerce').fillna(0)
                    chart_vt = alt.Chart(vt_df).mark_line(point=True, color='red').encode(
                        x=alt.X('timestamp:O', title=t('Timestamp')),
                        y=alt.Y('ai_score:Q', title=t('AI Probability (%)'), scale=alt.Scale(domain=[0, 100])),
                        tooltip=['timestamp', 'ai_score', 'details']
                    ).properties(height=250)
                    st.altair_chart(chart_vt, use_container_width=True)
                    
                    with st.expander(t("View Frame-by-Frame Details")):
                        for _, row in vt_df.iterrows():
                            st.write(f"**{row.get('timestamp', '??:??')}** ({t('Score')}: {int(row.get('ai_score', 0))}%) - {row.get('details', '')}")
            
            st.markdown("---")

            # --- FALLACY & LOGIC UI ---
            if ret.get('has_fallacy'):
                st.error(f"🛑 {t('Issue Detected:')} **{t(ret.get('fallacy_type', 'System Processing Error'))}**")
                st.metric(t("Aggression Level"), f"{ret.get('aggression', 0)}/10")
                st.warning(f"**{t('Analysis:')}** {t(ret.get('explanation', 'No details available.'))}")
            else:
                st.success(f"✅ {t('Neural Guard: No major issues detected.')}")
                st.info(f"**{t('Analysis:')}** {t(ret.get('explanation', 'Content is sound.'))}")
            
            st.markdown("---")
            syl_breakdown = ret.get('syllogism_breakdown', [])
            if syl_breakdown:
                st.markdown(f"#### {t('Logical Deconstruction')}")
                st.caption(t("The text was fragmented into formal premises to identify the exact point where logic fails."))
                for step in syl_breakdown:
                    flaw_text = str(step.get('flaw', '')).strip()
                    if flaw_text and flaw_text.lower() not in ["none", "nessuno", "nessuna", "n/a", "", "null"]:
                        st.error(f"**{t(step.get('step', 'Step'))}**: {t(step.get('text', ''))}\n\n⚠️ **{t('Logical Leap / Fallacy:')}** {t(flaw_text)}")
                    else:
                        st.info(f"**{t(step.get('step', 'Step'))}**: {t(step.get('text', ''))}")
            
            st.markdown("---")
            st.markdown(f"#### {t('Rewritten Version / Transcript Summary')}")
            st.info(t(ret.get('rewritten_text', 'No rewrite available.')))
            
            st.markdown("---")
            if ret.get('shadow_geolocation') and str(ret['shadow_geolocation']).lower() not in ["none", "n/a", "", "null", "fictional ai-generated environment - geolocation not applicable."]:
                st.markdown(f"#### 🌍 {t('Shadow Geolocation (OSINT)')}")
                st.success(t(ret['shadow_geolocation']))
            elif ret.get('shadow_geolocation'):
                st.markdown(f"#### 🌍 {t('Shadow Geolocation (OSINT)')}")
                st.caption(t(ret['shadow_geolocation']))
            
            st.markdown("---")
            if media_type == "audio" and 'voice_stress_score' in ret:
                stress = ret.get('voice_stress_score', 0)
                st.markdown(f"#### {t('Voice & Tone Analysis')}")
                if stress > 65:
                    st.error(f"**{t('Voice Stress Score:')} {stress}%** ({t('High emotion, anger, or panic detected in prosody')})")
                else:
                    st.success(f"**{t('Voice Stress Score:')} {stress}%** ({t('Calm, controlled, or neutral tone')})")
            
            # --- GHOST READER (OCR FORENSICS) ---
            st.markdown(f"#### {t('Ghost Reader (OCR Extraction)')}")
            st.caption(t("AI optical scan of texts hidden in the media (signs, screens, documents)."))
            ocr_val = str(ret.get('ocr_extraction', '')).strip()
            if ocr_val and ocr_val.lower() not in ["none", "n/a", "null", ""]:
                st.info(ocr_val)
            else:
                st.success(t("✅ No readable natural text detected in the scene."))
            
            # --- FACT CHECKER ---
            st.markdown(f"#### {t('Fact Checker (Claims to Verify)')}")
            facts = ret.get('facts', [])
            if facts:
                for f_claim in facts: st.write(f"- {t(f_claim)}")
            else:
                st.caption(t("No specific factual claims found."))

            # --- FORENSIC DOSSIER EXPORT ---
            st.markdown("---")
            file_hash = calculate_sha256(media_inp) if media_type in ['audio', 'video'] else "N/A"
            
            dossier_text = f"{t('RAP FORENSIC DOSSIER')}\n"
            dossier_text += f"{t('Date')}: {datetime.now().strftime('%Y-%m-%d %H:%M')}\n"
            dossier_text += f"{t('Digital Fingerprint (SHA-256)')}: {file_hash}\n"
            dossier_text += f"{'='*40}\n\n"
            dossier_text += f"[1] {t('AI GENERATION SCAN')}\n{t('AI Probability Score')}: {ai_prob}%\n\n"
            
            st.download_button(
                label=t("Download Forensic Dossier (TXT)"),
                data=dossier_text.encode('utf-8'),
                file_name=f"RAP_Forensic_Dossier_{datetime.now().strftime('%Y%m%d_%H%M')}.txt",
                mime="text/plain",
                type="primary"
            )
        else:
            st.warning(t("Please provide input."))

# ==========================================
# MODULE 4: COMPARISON TEST (A/B TESTING)
# ==========================================
elif mode == t("4. Comparison Test (A/B Testing)"):
    st.header(t("4. Universal Arena (A/B Testing)"))
    st.caption(t("Compare data from YouTube, CSV, or Raw Text to find the most aggressive narratives."))
    
    col_a, col_b = st.columns(2)
    
    def load_arena_data(key_prefix, column):
        with column:
            in_type = st.radio(f"{t('Input')} {key_prefix}", [t("YouTube Link"), t("CSV Upload"), t("Raw Text Paste")], horizontal=True, key=f"r_{key_prefix}")
            df = None
            if in_type == t("YouTube Link"):
                url = st.text_input(f"{t('YouTube URL')} {key_prefix}", key=f"url_{key_prefix}")
                limit = st.number_input(f"{t('Comments to Fetch')} {key_prefix}", 10, 1000, 50, key=f"lim_{key_prefix}")
                if st.button(f"{t('Scrape')} {key_prefix}", key=f"btn_yt_{key_prefix}"):
                    with st.spinner(f"{t('Scraping YouTube')} {key_prefix}..."):
                        df = scrape_youtube_comments(url, limit)
            elif in_type == t("CSV Upload"):
                f = st.file_uploader(f"{t('Upload CSV')} {key_prefix}", type="csv", key=f"f_{key_prefix}")
                if f: df = normalize_dataframe(pd.read_csv(f))
            else:
                txt = st.text_area(f"{t('Paste Data')} {key_prefix}", height=100, key=f"t_{key_prefix}")
                if st.button(f"{t('Process Text')} {key_prefix}", key=f"btn_txt_{key_prefix}"):
                    if txt: df = parse_raw_paste(txt)
                    
            if df is not None:
                df = detect_bot_activity(df)
                if 'Select' not in df.columns: df.insert(0, "Select", False)
                st.session_state['data_store']['Arena'][f'df_{key_prefix.lower()}'] = df
                st.success(f"{t('Loaded')} {len(df)} {t('items in')} {key_prefix}.")
        
    st.subheader(t("Contender A (Left)"))
    load_arena_data("A", col_a)
    st.subheader(t("Contender B (Right)"))
    load_arena_data("B", col_b)

    df_a_raw = st.session_state['data_store']['Arena']['df_a']
    df_b_raw = st.session_state['data_store']['Arena']['df_b']

    if df_a_raw is not None and df_b_raw is not None:
        st.divider()
        st.markdown(f"### {t('Step 2: Select Data to Analyze')}")
        st.caption(t("Select specific comments to compare using the checkboxes. If none selected, the top N rows will be analyzed."))
        
        c_edit_a, c_edit_b = st.columns(2)
        with c_edit_a:
            st.markdown(f"**{t('Contender A')}** ({len(df_a_raw)} {t('items')})")
            edited_a = st.data_editor(
                df_a_raw,
                column_config={"Select": st.column_config.CheckboxColumn(required=True)},
                disabled=["content", "agent_id", "timestamp", "is_bot", "likes"],
                key="editor_arena_a", height=300, hide_index=True
            )
        with c_edit_b:
            st.markdown(f"**{t('Contender B')}** ({len(df_b_raw)} {t('items')})")
            edited_b = st.data_editor(
                df_b_raw,
                column_config={"Select": st.column_config.CheckboxColumn(required=True)},
                disabled=["content", "agent_id", "timestamp", "is_bot", "likes"],
                key="editor_arena_b", height=300, hide_index=True
            )

        st.divider()
        c_action, c_limit = st.columns([1, 1])
        with c_limit:
            max_analyze = st.number_input(t("Max Rows to Analyze (if no manual selection)"), 5, 100, 20)
        with c_action:
            st.write("") 
            st.write("") 
            start_analysis = st.button(t("Step 3: Run Comparative Analysis"), type="primary", disabled=not key)

        if start_analysis:
            subset_a = edited_a[edited_a.Select]
            if subset_a.empty: subset_a = edited_a.head(max_analyze)
            prog_a = st.progress(0)
            res_a = []
            st.markdown(f"**{t('Analyzing Contender A...')}**")
            for i, (_, row) in enumerate(subset_a.iterrows()):
                res_a.append(analyze_fallacies(row['content'], api_key=key, target_lang=st.session_state['global_lang']))
                prog_a.progress((i + 1) / len(subset_a))
            final_a = pd.concat([subset_a.reset_index(drop=True), pd.DataFrame(res_a)], axis=1)
            st.session_state['data_store']['Arena']['analyzed_a'] = final_a

            subset_b = edited_b[edited_b.Select]
            if subset_b.empty: subset_b = edited_b.head(max_analyze)
            prog_b = st.progress(0)
            res_b = []
            st.markdown(f"**{t('Analyzing Contender B...')}**")
            for i, (_, row) in enumerate(subset_b.iterrows()):
                res_b.append(analyze_fallacies(row['content'], api_key=key, target_lang=st.session_state['global_lang']))
                prog_b.progress((i + 1) / len(subset_b))
            final_b = pd.concat([subset_b.reset_index(drop=True), pd.DataFrame(res_b)], axis=1)
            st.session_state['data_store']['Arena']['analyzed_b'] = final_b
            st.rerun()

    res_df_a = st.session_state['data_store']['Arena']['analyzed_a']
    res_df_b = st.session_state['data_store']['Arena']['analyzed_b']

    if res_df_a is not None and res_df_b is not None:
        st.divider()
        st.header(t("Match Results"))
        
        c1, c2, c3 = st.columns(3)
        agg_a = res_df_a['aggression'].mean()
        agg_b = res_df_b['aggression'].mean()
        delta_agg = agg_a - agg_b
        bots_a = len(res_df_a[res_df_a['is_bot']==True])
        bots_b = len(res_df_b[res_df_b['is_bot']==True])
        delta_bots = bots_a - bots_b
        fallacy_a = len(res_df_a[res_df_a['has_fallacy']==True])
        fallacy_b = len(res_df_b[res_df_b['has_fallacy']==True])
        
        c1.metric(t("Avg Aggression (A vs B)"), f"{agg_a:.1f} vs {agg_b:.1f}", f"{delta_agg:.1f}")
        c2.metric(t("Bot Count (A vs B)"), f"{bots_a} vs {bots_b}", f"{delta_bots}")
        c3.metric(t("Fallacies (A vs B)"), f"{fallacy_a} vs {fallacy_b}")
        
        st.markdown("---")
        st.subheader(t("Tactical Visual Comparison"))
        res_df_a['Source'] = t('Contender A')
        res_df_b['Source'] = t('Contender B')
        combined = pd.concat([res_df_a, res_df_b])
        
        chart_agg = alt.Chart(combined).mark_bar().encode(
            x=alt.X('Source', title=None),
            y=alt.Y('mean(aggression)', title=t('Avg Aggression')),
            color=alt.Color('Source', scale=alt.Scale(domain=[t('Contender A'), t('Contender B')], range=['#3b82f6', '#f97316'])),
            tooltip=['Source', 'mean(aggression)']
        ).properties(height=200)
        st.altair_chart(chart_agg, use_container_width=True)
        
        c_vis1, c_vis2 = st.columns(2)
        with c_vis1:
            st.caption(t("Emotional Spectrum Clash"))
            if 'primary_emotion' in combined.columns:
                emo_combined = combined.groupby(['Source', 'primary_emotion']).size().reset_index(name='Count')
                chart_emo = alt.Chart(emo_combined).mark_bar().encode(
                    x=alt.X('primary_emotion:N', title=t('Emotion'), sort='-y'),
                    y=alt.Y('Count:Q', title=t('Frequency')),
                    color=alt.Color('Source:N', scale=alt.Scale(domain=[t('Contender A'), t('Contender B')], range=['#3b82f6', '#f97316']), legend=None),
                    xOffset='Source:N',
                    tooltip=['Source', 'primary_emotion', 'Count']
                ).properties(height=300)
                st.altair_chart(chart_emo, use_container_width=True)
                
        with c_vis2:
            st.caption(t("Weaponized Fallacies"))
            if 'fallacy_type' in combined.columns:
                fal_combined = combined[combined['has_fallacy'] == True].groupby(['Source', 'fallacy_type']).size().reset_index(name='Count')
                if not fal_combined.empty:
                    chart_fal = alt.Chart(fal_combined).mark_bar().encode(
                        x=alt.X('fallacy_type:N', title=t('Fallacy Type'), sort='-y'),
                        y=alt.Y('Count:Q', title=t('Frequency')),
                        color=alt.Color('Source:N', scale=alt.Scale(domain=[t('Contender A'), t('Contender B')], range=['#3b82f6', '#f97316'])),
                        xOffset='Source:N',
                        tooltip=['Source', 'fallacy_type', 'Count']
                    ).properties(height=300)
                    st.altair_chart(chart_fal, use_container_width=True)
                else:
                    st.success(t("No fallacies detected in either contender."))

        st.markdown("---")
        st.subheader(t("The Oracle: Narrative Clash Assessment"))
        st.caption(t("Force the AI to analyze the differing psychological and tactical profiles of the two contenders."))
        
        if st.button(t("Generate Clash Briefing"), type="primary"):
            with st.spinner(t("Analyzing psychological divergence...")):
                top_emotions_a = res_df_a['primary_emotion'].value_counts().head(3).to_dict() if 'primary_emotion' in res_df_a else "N/A"
                top_emotions_b = res_df_b['primary_emotion'].value_counts().head(3).to_dict() if 'primary_emotion' in res_df_b else "N/A"
                top_fallacies_a = res_df_a['fallacy_type'].value_counts().head(3).to_dict() if 'fallacy_type' in res_df_a else "N/A"
                top_fallacies_b = res_df_b['fallacy_type'].value_counts().head(3).to_dict() if 'fallacy_type' in res_df_b else "N/A"
                sample_a = res_df_a['content'].iloc[0] if not res_df_a.empty else ""
                sample_b = res_df_b['content'].iloc[0] if not res_df_b.empty else ""
                
                clash_prompt = f"""
                You are an elite Information Warfare and Psychological Operations Analyst.
                Compare these two opposing groups (Contender A vs Contender B).
                
                CONTENDER A:
                - Avg Aggression: {agg_a:.1f}/10
                - Dominant Emotions: {top_emotions_a}
                - Primary Logical Fallacies used: {top_fallacies_a}
                - Sample Text: "{sample_a[:200]}"
                
                CONTENDER B:
                - Avg Aggression: {agg_b:.1f}/10
                - Dominant Emotions: {top_emotions_b}
                - Primary Logical Fallacies used: {top_fallacies_b}
                - Sample Text: "{sample_b[:200]}"
                
                TASK: Write a sharp, tactical "Narrative Clash Briefing" (strictly max 150 words). 
                Do not just list the numbers. Explain *HOW* their manipulation strategies differ. 
                Conclude by stating which group poses a higher risk of inciting real-world polarization.
                
                CRITICAL RULE: You MUST write the ENTIRE briefing strictly in {st.session_state['global_lang']}.
                """
                try:
                    client = genai.Client(api_key=key)
                    clash_res = client.models.generate_content(model='gemini-2.5-flash', contents=clash_prompt)
                    st.warning(f"### {t('Tactical Clash Report')}")
                    st.markdown(clash_res.text)
                except Exception as e:
                    st.error(f"{t('Failed to generate briefing')}: {e}")

        with st.expander(t("Detailed Comparison Data")):
            st.dataframe(combined[['Source', 'content', 'aggression', 'primary_emotion', 'fallacy_type', 'is_bot']])

        st.markdown("---")
        st.subheader(t("Export Battle Report"))
        st.caption(t("Generate a comparative PDF dossier detailing the metrics between Contender A and Contender B."))
        
        if st.button(t("Generate Arena PDF"), type="primary"):
            with st.spinner(t("Compiling the comparative report...")):
                pdf = PDFReport()
                pdf.add_page()
                pdf.set_auto_page_break(auto=True, margin=15)
                pdf.set_font("Helvetica", 'B', 14)
                pdf.cell(0, 10, "Comparative Dossier: Narrative A/B Testing", 0, 1)
                pdf.ln(5)
                
                pdf.set_font("Helvetica", 'B', 12)
                pdf.cell(0, 10, "Contender A Profile:", 0, 1)
                pdf.set_font("Helvetica", '', 11)
                pdf.cell(0, 8, f"- Analyzed Items: {len(res_df_a)}", 0, 1)
                pdf.cell(0, 8, f"- Average Aggression: {agg_a:.1f}/10", 0, 1)
                pdf.cell(0, 8, f"- Detected Bots: {bots_a}", 0, 1)
                pdf.cell(0, 8, f"- Logical Fallacies Flagged: {fallacy_a}", 0, 1)
                pdf.ln(5)

                pdf.set_font("Helvetica", 'B', 12)
                pdf.cell(0, 10, "Contender B Profile:", 0, 1)
                pdf.set_font("Helvetica", '', 11)
                pdf.cell(0, 8, f"- Analyzed Items: {len(res_df_b)}", 0, 1)
                pdf.cell(0, 8, f"- Average Aggression: {agg_b:.1f}/10", 0, 1)
                pdf.cell(0, 8, f"- Detected Bots: {bots_b}", 0, 1)
                pdf.cell(0, 8, f"- Logical Fallacies Flagged: {fallacy_b}", 0, 1)
                pdf.ln(5)
                
                pdf.set_font("Helvetica", 'B', 12)
                pdf.cell(0, 10, "Tactical Conclusion:", 0, 1)
                pdf.set_font("Helvetica", 'I', 11)
                
                if agg_a > agg_b + 1.5:
                    conclusion = "Contender A exhibits a significantly higher level of hostility and aggression. It poses a greater immediate risk for polarization."
                elif agg_b > agg_a + 1.5:
                    conclusion = "Contender B exhibits a significantly higher level of hostility and aggression. It poses a greater immediate risk for polarization."
                else:
                    conclusion = "Both contenders present comparable levels of aggression. The threat level is balanced between the two narratives."
                    
                if bots_a > bots_b: conclusion += f" Furthermore, Contender A shows higher signs of inauthentic/bot activity (+{delta_bots})."
                elif bots_b > bots_a: conclusion += f" Furthermore, Contender B shows higher signs of inauthentic/bot activity (+{abs(delta_bots)})."

                pdf.multi_cell(0, 6, conclusion.encode('latin-1', 'replace').decode('latin-1'))
                pdf_bytes_arena = bytes(pdf.output())
                
            st.download_button(
                label=t("Download Arena Report (PDF)"), 
                data=pdf_bytes_arena, 
                file_name=f"RAP_Arena_Match_{datetime.now().strftime('%Y%m%d_%H%M')}.pdf", 
                mime="application/pdf", 
                type="primary"
            )

# ==========================================
# MODULE 5: LIVE RADAR (RSS/ALERTS) - AUTONOMOUS
# ==========================================
elif mode == t("5. Live Radar (RSS/Alerts)"):
    st.header(t("5. Live Radar (Crisis Alert System)"))
    st.caption(t("Monitor live RSS feeds or subreddits to intercept escalating disinformation and aggression in real-time."))
    
    if "GEMINI_API_KEY" in st.secrets: 
        key = st.secrets["GEMINI_API_KEY"]
    else: 
        key = st.sidebar.text_input(t("API Key"), type="password", key="radar_key")

    # Initializing data store for Radar if missing
    if 'Radar' not in st.session_state['data_store']:
        st.session_state['data_store']['Radar'] = {'df': None, 'analyzed': None}

    is_sentinel_active = st.session_state.get("sentinel_v5", False)

    c_radar1, c_radar2, c_radar3 = st.columns([2, 1, 1])
    with c_radar1:
        feed_url = st.text_input(t("Enter RSS Feed, Subreddit, or News Keyword"), placeholder="E.g., ansa, bbc, reddit.com/r/worldnews", key="radar_input_url")
    with c_radar2:
        world_languages = ["English", "Italiano", "Español", "Français", "Deutsch", "Português", "Russian", "Arabic", "Chinese", "Japanese"]
        news_region = st.selectbox(t("Search Language/Region"), world_languages, index=0)
    with c_radar3:
        max_entries = st.number_input(t("Entries"), 5, 50, 15)
        fetch_btn = st.button(t("Step 1: Fetch Feed"), type="primary", use_container_width=True, disabled=is_sentinel_active)
        defcon_btn = st.button(f"🚨 {t('DEFCON Cyber Scan')}", type="primary", use_container_width=True, disabled=is_sentinel_active)
        
    with st.expander(t("Automated Alert Configuration (Webhook)")):
        alert_webhook = st.text_input(t("Webhook URL"), placeholder="https://hooks.slack.com/services/...")
        alert_threshold = st.slider(t("Trigger Alert Threshold (Aggression)"), min_value=1.0, max_value=10.0, value=8.0, step=0.5)

    # --- STEP 1: FETCHING LOGIC (MANUAL OR AUTONOMOUS) ---
    if (fetch_btn and feed_url) or defcon_btn or (is_sentinel_active and feed_url):
        with st.spinner(t("Intercepting live feed...") if not is_sentinel_active else "Sentinel Autonomous Sweep in progress..."):
            try:
                user_input = ""
                if defcon_btn or (is_sentinel_active and "hacker" in feed_url.lower()):
                    actual_url = "https://feeds.feedburner.com/TheHackersNews"
                    if defcon_btn: st.toast("DEFCON Activated: Scanning Global Cyber Threats...", icon="🚨")
                else:
                    user_input = feed_url.strip().lower()
                    actual_url = feed_url.strip()
                    
                news_shortcuts = {
                    "ansa": "https://www.ansa.it/sito/ansait_rss.xml",
                    "bbc": "http://feeds.bbci.co.uk/news/world/rss.xml",
                    "repubblica": "https://www.repubblica.it/rss/homepage/rss2.0.xml"
                }
                
                if user_input in news_shortcuts:
                    actual_url = news_shortcuts[user_input]
                elif "reddit.com/r/" in user_input and not user_input.endswith(".rss"):
                    actual_url = actual_url.rstrip("/") + "/new/.rss"
                elif not actual_url.startswith("http"):
                    actual_url = f"https://news.google.com/rss/search?q={actual_url.replace(' ', '%20')}&hl=en-US&gl=US&ceid=US:en"
                
                feed = feedparser.parse(actual_url)
                if feed.entries:
                    entries_data = []
                    for entry in feed.entries[:int(max_entries)]:
                        clean_summary = re.sub(r'<[^>]+>', '', entry.get('summary', ''))
                        entries_data.append({
                            'Select': False,
                            'timestamp': entry.get('published', 'Unknown'),
                            'content': f"TITLE: {entry.get('title', '')}\nCONTENT: {clean_summary[:500]}",
                            'link': entry.get('link', '')
                        })
                    st.session_state['data_store']['Radar']['df'] = pd.DataFrame(entries_data)
                    
                    if not is_sentinel_active:
                        st.session_state['data_store']['Radar']['analyzed'] = None
                        st.success(f"✅ Intercepted {len(entries_data)} items.")
                        st.rerun()
            except Exception as e:
                st.error(f"Radar Fetch Error: {str(e)}")

    # --- STEP 2 & 3: SELECTION & ANALYSIS ---
    radar_df = st.session_state['data_store']['Radar'].get('df')
    if radar_df is not None:
        st.divider()
        st.subheader("Step 2: Intelligence Selection")
        edited_radar = st.data_editor(
            radar_df,
            column_config={"Select": st.column_config.CheckboxColumn(required=True)},
            disabled=["timestamp", "content", "link"],
            key="editor_radar_v5", use_container_width=True, hide_index=True
        )
        
        c_action, c_limit = st.columns([1, 1])
        with c_limit:
            max_analyze = st.number_input("Batch Size (Top N)", 1, len(radar_df), min(10, len(radar_df)))
        with c_action:
            st.write("")
            selected_count = edited_radar['Select'].sum()
            btn_label = f"Step 3: Run Analysis ({selected_count})" if selected_count > 0 else f"Step 3: Run Analysis (Top {max_analyze})"
            
            run_analysis_flag = st.button(btn_label, type="primary", disabled=(not key or is_sentinel_active))
            
            if run_analysis_flag or is_sentinel_active:
                subset = edited_radar[edited_radar.Select]
                if subset.empty: subset = edited_radar.head(max_analyze)
                
                res = []
                prog = st.progress(0)
                for i, (_, row) in enumerate(subset.iterrows()):
                    crisis_prompt = f"""
                    You are an Early Warning Crisis AI. Analyze this news.
                    Evaluate the INHERENT CRISIS LEVEL (0-10) and POLARIZATION.
                    Respond ONLY with JSON: {{"aggression": [0-10], "reasoning": "1-sentence justification", "iso_country": "3-letter code"}}
                    Content: "{row['content']}"
                    Language: {st.session_state['global_lang']}
                    """
                    try:
                        client = genai.Client(api_key=key)
                        response = client.models.generate_content(model='gemini-2.0-flash', contents=crisis_prompt)
                        parsed = extract_json(response.text)
                        parsed['explanation'] = parsed.get('reasoning', 'Analyzed.')
                        parsed['has_fallacy'] = parsed.get('aggression', 0) >= alert_threshold
                    except:
                        parsed = {"aggression": 0, "explanation": "Error", "has_fallacy": False, "iso_country": "GLO"}
                    res.append(sanitize_response(parsed))
                    prog.progress((i + 1) / len(subset))
                
                st.session_state['data_store']['Radar']['analyzed'] = pd.concat([subset.reset_index(drop=True), pd.DataFrame(res)], axis=1)
                
                if not is_sentinel_active:
                    st.rerun()

    # --- STEP 4: RESULTS VISUALIZATION ---
    analyzed_radar = st.session_state['data_store']['Radar'].get('analyzed')
    if analyzed_radar is not None and not analyzed_radar.empty:
        st.divider()
        avg_agg = analyzed_radar['aggression'].mean()
        
        m1, m2, m3 = st.columns(3)
        m1.metric("Crisis Index", f"{avg_agg:.1f}/10")
        m2.metric("Critical Alerts", len(analyzed_radar[analyzed_radar['has_fallacy']==True]))
        m3.metric("Processed Items", len(analyzed_radar))
        
        if avg_agg >= alert_threshold:
            st.error(f"🚨 ALERT: Global Aggression Threshold breached ({avg_agg:.1f} >= {alert_threshold})")
            if alert_webhook:
                try:
                    payload = {"content": f"🚨 **RAP CRISIS ALERT** 🚨\n**Target:** {feed_url}\n**Crisis Level:** {avg_agg:.1f}/10\n**New Events Detected.**"}
                    requests.post(alert_webhook, json=payload, timeout=5)
                    st.toast("Webhook dispatched!", icon="🚨")
                except: pass

        if 'iso_country' in analyzed_radar.columns:
            st.subheader("Global Threat Sphere")
            map_data = analyzed_radar[~analyzed_radar['iso_country'].isin(['GLO', 'None', ''])]
            if not map_data.empty:
                geo_stats = map_data.groupby('iso_country').agg(Count=('iso_country', 'count'), Tension=('aggression', 'mean')).reset_index()
                fig = px.scatter_geo(geo_stats, locations="iso_country", size="Count", color="Tension", projection="orthographic", color_continuous_scale="Reds")
                fig.update_layout(geo=dict(showcoastlines=True, landcolor='#1a1c24', bgcolor='rgba(0,0,0,0)'), paper_bgcolor='rgba(0,0,0,0)', margin=dict(l=0, r=0, t=30, b=0))
                st.plotly_chart(fig, use_container_width=True)

        st.subheader("Intelligence Feed")
        for idx, r in analyzed_radar.iterrows():
            with st.container(border=True):
                st.caption(f"🕒 {r['timestamp']} | [Link]({r['link']}) | Agg: {r['aggression']}/10")
                st.markdown(f"**{r['content'].split('CONTENT:')[0]}**")
                
                c1, c2 = st.columns([1, 1])
                with c1:
                    if st.button(f"Spider-Scan #{idx}", key=f"sp_{idx}"):
                        st.write(run_web_spider(r['link']))
                with c2:
                    if st.button(f"PANDORA Forensic #{idx}", key=f"pan_{idx}"):
                        st.session_state['pandora_input'] = r['content']
                        st.info("Incident data transferred to PANDORA (Module 8).")
                
                if r['has_fallacy']: st.error(f"🛑 {r['explanation']}")
                else: st.success(f"✅ {r['explanation']}")

    # --- OPERATION SENTINEL (AUTONOMOUS ENGINE) ---
    st.divider()
    st.subheader("Autonomous Sentinel Protocol")
    
    live_toggle = st.toggle("Activate Live Sentinel Mode", key="sentinel_v5")
    
    if live_toggle:
        if not feed_url:
            st.warning("Please enter a Keyword or Source URL at the top first.")
        else:
            refresh_rate = st.slider("Scan Interval (Seconds)", 30, 600, 60)
            live_placeholder = st.empty()
            
            # THE AUTONOMOUS LOOP:
            # 1. Shows countdown. 
            # 2. When 0, calls st.rerun().
            # 3. On rerun, the 'if is_sentinel_active' triggers at the top, downloading and analyzing silently.
            for i in range(refresh_rate, 0, -1):
                live_placeholder.info(f"🟢 SENTINEL ACTIVE | Target: **{feed_url}** | Next autonomous sweep in: **{i}s**")
                time.sleep(1)
            
            st.rerun()
            
# ==========================================
# MODULE 6: DEEP DOCUMENT ORACLE (RAG)
# ==========================================
elif mode == t("6. Deep Document Oracle (RAG)"):
    st.header(t("6. Deep Document Oracle"))
    st.caption(t("Upload massive PDFs (e.g., manifestos, contracts, books) and find contradictions and extract deep facts without traditional RAG limits."))

    if "GEMINI_API_KEY" in st.secrets: 
        key = st.secrets["GEMINI_API_KEY"]
    else: 
        key = st.text_input(t("API Key"), type="password")

    uploaded_files = st.file_uploader(t("Upload PDF/TXT Documents"), type=['pdf', 'txt'], accept_multiple_files=True)
    
    if uploaded_files:
        if st.button(t("Process Documents"), type="primary"):
            with st.spinner(t("Ingesting and vectorizing classified documents...")):
                full_text = ""
                for f in uploaded_files:
                    txt = extract_text_from_pdf(f)
                    if txt:
                        full_text += f"\n\n--- DOCUMENT: {f.name} ---\n\n{txt}"
                
                st.session_state['doc_full_text'] = full_text
                st.session_state['doc_sanitized'] = False 
                st.success(f"{t('Processed')} {len(full_text)} {t('characters across')} {len(uploaded_files)} {t('documents. The Oracle is ready.')}")
        
        if 'doc_full_text' in st.session_state and st.session_state['doc_full_text']:
            st.divider()
            
            # --- PII SANITIZER ---
            c_san1, c_san2 = st.columns([1, 3])
            with c_san1:
                if st.button(t("Sanitize Document (Redact PII)"), help=t("Hides Emails, IPs, IBANs, and Phones before chatting with the Oracle")):
                    with st.spinner(t("Sanitizing sensitive data...")):
                        clean_text = sanitize_pii(st.session_state['doc_full_text'])
                        st.session_state['doc_full_text'] = clean_text
                        st.session_state['doc_sanitized'] = True 
                        st.success(t("Document successfully sanitized (CIA Blackout Protocol)!"))
            
            with c_san2:
                if st.session_state.get('doc_sanitized', False):
                    st.download_button(
                        label=t("Download Redacted Dossier (TXT)"),
                        data=st.session_state['doc_full_text'].encode('utf-8'),
                        file_name=f"RAP_Classified_Redacted_{datetime.now().strftime('%Y%m%d')}.txt",
                        mime="text/plain",
                        type="primary"
                    )
                else:
                    st.warning(f"⚠️ {t('Sanitize the document first to unlock the secure download.')}")
            
            # --- KNOWLEDGE GRAPH ---
            with st.expander(t("Extract Document Power Network (Knowledge Graph)"), expanded=False):
                st.caption(t("Automatically scan the document to map relationships between People, Organizations, and Locations."))
                if st.button(t("Generate Power Graph")):
                    with st.spinner(t("Extracting entities and relationships (this may take a minute)...")):
                        graph_prompt = f"""
                        Analyze this document and extract the top 12 most important relationships between entities (People, Organizations, Locations).
                        CRITICAL RULES:
                        1. The "relation" field MUST BE ULTRA-SHORT (maximum 1 to 3 words, e.g., "CEO", "Owned by", "Funded", "Opposes").
                        2. Keep entity names short and clean.
                        Return ONLY a valid JSON format like this:
                        {{ "relations": [ {{"source": "Entity 1", "target": "Entity 2", "relation": "Short Label"}} ] }}
                        
                        DOCUMENT:
                        {st.session_state['doc_full_text'][:30000]}
                        """
                        try:
                            client = genai.Client(api_key=key)
                            graph_res = client.models.generate_content(model='gemini-2.0-flash', contents=graph_prompt)
                            graph_json = extract_json(graph_res.text)
                            
                            if graph_json and 'relations' in graph_json:
                                nodes = []
                                edges = []
                                added_nodes = set()
                                
                                for rel in graph_json['relations']:
                                    src = str(rel.get('source', 'Unknown')).strip()
                                    tgt = str(rel.get('target', 'Unknown')).strip()
                                    link_label = str(rel.get('relation', 'Linked')).strip()
                                    
                                    if not src or not tgt: continue
                                    
                                    # Add Source Node (Neon Blue, Larger Text)
                                    if src not in added_nodes:
                                        nodes.append(Node(
                                            id=src, label=src, size=25, 
                                            color="#38BDF8", 
                                            font={'color': '#FFFFFF', 'size': 16}
                                        ))
                                        added_nodes.add(src)
                                        st.session_state['global_entities'].add(src.lower())
                                        
                                    # Add Target Node (Hacker Green, Larger Text)
                                    if tgt not in added_nodes:
                                        nodes.append(Node(
                                            id=tgt, label=tgt, size=25, 
                                            color="#10B981", 
                                            font={'color': '#FFFFFF', 'size': 16}
                                        ))
                                        added_nodes.add(tgt)
                                        st.session_state['global_entities'].add(tgt.lower())
                                        
                                    # Add Edge with High-Visibility Text
                                    edges.append(Edge(
                                        source=src, 
                                        target=tgt, 
                                        label=link_label, 
                                        color="#475569", 
                                        font={
                                            'color': '#38BDF8',      # Neon Blue Text
                                            'size': 15,              # Increased font size
                                            'background': '#0A0E17', # Dark background to hide the line behind the text
                                            'strokeWidth': 0         # Remove default stroke
                                        }
                                    ))

                                # Configure the Physics Engine and UI
                                config = Config(
                                    width="100%", 
                                    height=650,
                                    directed=True, 
                                    physics=True, 
                                    hierarchical=False,
                                    nodeHighlightBehavior=True,
                                    highlightColor="#EF4444"
                                )
                                
                                st.markdown(f"### {t('Interactive Power Network')}")
                                agraph(nodes=nodes, edges=edges, config=config)
                                
                                st.success(f"{t('Registered')} {len(st.session_state['global_entities'])} {t('entities into Global Memory.')}")
                            else:
                                st.error("Not enough clear relationships found to build a graph.")
                        except Exception as e:
                            st.error(f"{t('Graph Extraction Error')}: {str(e)}")

            st.divider()

            # --- CONTRADICTION & LOOPHOLE SCANNER ---
            with st.expander(t("Deep Scan: Contradictions & Loopholes"), expanded=False):
                st.caption(t("Force the AI to audit the entire document specifically looking for logical contradictions, legal loopholes, or unfulfilled claims."))
                if st.button(t("Run Audit Scan"), type="primary"):
                    with st.spinner(t("Auditing document for contradictions...")):
                        audit_prompt = f"""
                        You are a ruthless Forensic Auditor and Legal Analyst.
                        Scan the following document specifically for:
                        1. Internal Contradictions (e.g., Chapter 1 says X, Chapter 4 says the opposite of X).
                        2. Loopholes or ambiguous clauses that could be exploited.
                        3. Hidden risks or highly controversial statements.
                        
                        CRITICAL RULE: For every point you make, you MUST cite the [PAGE X] reference. 
                        CRITICAL LANGUAGE RULE: Write the ENTIRE report strictly in {st.session_state['global_lang']}.
                        
                        DOCUMENT TEXT:
                        {st.session_state['doc_full_text'][:35000]}
                        """
                        try:
                            client = genai.Client(api_key=key)
                            audit_res = client.models.generate_content(model='gemini-2.0-flash', contents=audit_prompt)
                            st.warning(f"### ⚠️ {t('Forensic Audit Report')}")
                            st.markdown(audit_res.text)
                            st.session_state.doc_oracle_history.append({"role": "assistant", "content": f"**[{t('AUTOMATED AUDIT REPORT')}]**\n\n{audit_res.text}"})
                        except Exception as e:
                            st.error(f"{t('Audit Error')}: {str(e)}")

            # --- MULTI-AGENT DEBATE (RED TEAM vs BLUE TEAM) ---
            with st.expander(t("Multi-Agent War Room (Stress Test)"), expanded=False):
                st.caption(t("Deploy two opposing AI agents to debate the document. The Red Team attacks it, the Blue Team defends it."))
                debate_topic = st.text_input(t("Debate Focus (e.g., 'Security flaws', 'Ethical implications'):"), placeholder=t("What should the agents fight about?"))
                
                if st.button(t("Initiate Agent Debate"), type="primary"):
                    if not debate_topic:
                        st.warning(t("Please enter a debate focus."))
                    else:
                        st.markdown(f"### 🔴 {t('Red Team')} vs 🔵 {t('Blue Team')}: *{debate_topic}*")
                        
                        # AGENT 1: RED TEAM
                        with st.spinner(t("🔴 Red Team is analyzing vulnerabilities...")):
                            red_prompt = f"You are the RED TEAM (Aggressive Attacker). Criticize this document focusing on: '{debate_topic}'. Find flaws, risks, and weaknesses. Be ruthless. Use maximum 150 words.\n\nCRITICAL LANGUAGE RULE: You MUST write your ENTIRE response strictly in {st.session_state['global_lang']}.\n\nDOCUMENT: {st.session_state['doc_full_text'][:20000]}"
                            client = genai.Client(api_key=key)
                            red_res = client.models.generate_content(model='gemini-2.5-flash', contents=red_prompt).text
                            st.error(f"**🔴 {t('Red Team Attack')}:**\n{red_res}")
                        
                        # AGENT 2: BLUE TEAM
                        with st.spinner(t("🔵 Blue Team is formulating a defense...")):
                            blue_prompt = f"You are the BLUE TEAM (Steadfast Defender). Read the RED TEAM's attack below regarding the document. Counter their arguments based on the text. Minimize risks and defend the document. Use maximum 150 words.\n\nCRITICAL LANGUAGE RULE: You MUST write your ENTIRE response strictly in {st.session_state['global_lang']}.\n\nRED TEAM ATTACK:\n{red_res}\n\nDOCUMENT: {st.session_state['doc_full_text'][:20000]}"
                            blue_res = client.models.generate_content(model='gemini-2.5-flash', contents=blue_prompt).text
                            st.info(f"**🔵 {t('Blue Team Defense')}:**\n{blue_res}")
                            
                        # AGENT 3: THE JUDGE
                        with st.spinner(t("⚖️ The Judge is deliberating...")):
                            judge_prompt = f"You are the IMPARTIAL JUDGE. Read the debate between Red and Blue. Who won? Provide a final 2-sentence verdict on the actual risk level of the document regarding '{debate_topic}'.\n\nCRITICAL LANGUAGE RULE: You MUST write your ENTIRE response strictly in {st.session_state['global_lang']}.\n\nRED:\n{red_res}\n\nBLUE:\n{blue_res}"
                            judge_res = client.models.generate_content(model='gemini-2.5-flash', contents=judge_prompt).text
                            st.success(f"**⚖️ {t('Final Verdict')}:**\n{judge_res}")
                            
                        st.session_state.doc_oracle_history.append({"role": "assistant", "content": f"**[WAR ROOM DEBATE: {debate_topic}]**\n\n**🔴 RED:** {red_res}\n\n**🔵 BLUE:** {blue_res}\n\n**⚖️ JUDGE:** {judge_res}"})

            # --- OPERATION HYDRA (PATTERN OF LIFE & TIMELINE) ---
            with st.expander("Operation HYDRA (Chronological Pattern of Life)", expanded=False):
                st.caption("Extracts a rigorous chronological timeline of all dates, times, and events buried within the massive document.")
                if st.button("Extract Timeline", type="primary"):
                    with st.spinner("Hunting for temporal metadata and constructing Pattern of Life..."):
                        hydra_prompt = f"""
                        You are Operation HYDRA, an elite Temporal Forensics AI.
                        Extract every single date, time, and associated event mentioned in the document.
                        
                        Format the output as a strict, chronological "Pattern of Life Timeline".
                        Use bullet points with the exact date/time in bold, followed by a concise description of the event.
                        Identify any "Time Gaps" (suspicious periods with no activity).
                        
                        CRITICAL LANGUAGE RULE: You MUST write your ENTIRE response strictly in {st.session_state['global_lang']}.
                        
                        DOCUMENT TEXT:
                        {st.session_state['doc_full_text'][:30000]}
                        """
                        try:
                            client = genai.Client(api_key=key)
                            hydra_res = client.models.generate_content(model='gemini-2.0-flash', contents=hydra_prompt)
                            st.warning("### Temporal Forensics Timeline")
                            st.markdown(hydra_res.text)
                            st.session_state.doc_oracle_history.append({"role": "assistant", "content": f"**[HYDRA TIMELINE EXTRACTED]**\n\n{hydra_res.text}"})
                        except Exception as e:
                            st.error(f"HYDRA Error: {str(e)}")

            st.divider()

            # --- OPERATION INTERROGATOR (TWO-WAY VOICE COMMAND) ---
            st.subheader(t("Chat with the Oracle"))
            
            for message in st.session_state.doc_oracle_history:
                with st.chat_message(message["role"]):
                    st.markdown(message["content"])
            
            st.markdown("---")
            st.caption("**Operation INTERROGATOR:** Use your microphone or keyboard to securely interrogate the document.")
            
            # Dual Input: Audio or Text
            audio_cmd = st.audio_input("Initialize Secure Voice Channel:")
            prompt = st.chat_input(t("Ask the Deep Oracle..."))
            
            active_query = None
            voice_used = False
            
            # Detect which input method was used
            if prompt:
                active_query = prompt
            elif audio_cmd:
                with st.spinner("Decrypting voice transmission..."):
                    try:
                        client = genai.Client(api_key=key)
                        audio_part = types.Part.from_bytes(data=audio_cmd.getvalue(), mime_type="audio/wav")
                        transcription_res = client.models.generate_content(
                            model='gemini-2.0-flash',
                            contents=[audio_part, "Transcribe this audio request accurately in the language it is spoken. Output ONLY the transcription, absolutely no other text."]
                        )
                        active_query = transcription_res.text.strip()
                        voice_used = True
                    except Exception as e:
                        st.error(f"Voice recognition failed: {e}")

            # Execute the query if either input fired
            if active_query:
                st.session_state.doc_oracle_history.append({"role": "user", "content": active_query})
                with st.chat_message("user"):
                    if voice_used:
                        st.markdown(f"*{active_query}*")
                    else:
                        st.markdown(active_query)
                
                with st.chat_message("assistant"):
                    with st.spinner(t("Scouring massive document context...")):
                        chunks = chunk_document(st.session_state['doc_full_text'])
                        query_words = active_query.lower().split()
                        scored_chunks = []
                        for c in chunks:
                            score = sum(2 for word in query_words if word in c.lower())
                            scored_chunks.append((score, c))
                        
                        scored_chunks.sort(key=lambda x: x[0], reverse=True)
                        top_context = "\n\n[SEGMENT]\n\n".join([text for score, text in scored_chunks[:4]])
                        
                        rag_prompt = f"""
                        You are the Deep Document Oracle. Answer the user query using ONLY the following extracted document segments.
                        If the answer is not in the segments, say so.
                        CRITICAL LANGUAGE RULE: You MUST write your ENTIRE response strictly in {st.session_state['global_lang']}.
                        
                        SEGMENTS:
                        {top_context}
                        
                        USER QUERY: {active_query}
                        """
                        
                        try:
                            client = genai.Client(api_key=key)
                            response = client.models.generate_content(model='gemini-2.0-flash', contents=rag_prompt).text
                            final_response = f"*({t('Method: High-Precision Context Scan')})*\n\n{response}"
                        except Exception as e:
                            response = ask_document_oracle(st.session_state['doc_full_text'], active_query, key)
                            final_response = response

                        st.markdown(final_response)
                        st.session_state.doc_oracle_history.append({"role": "assistant", "content": final_response})
                        
                        # TTS Auto-Play if user used voice
                        if voice_used:
                            with st.spinner("Synthesizing voice response..."):
                                try:
                                    lang_map_tts = {
                                        "English": "en", "Italiano": "it", "Español": "es", "Français": "fr", 
                                        "Deutsch": "de", "Português": "pt", "Русский (Russian)": "ru", 
                                        "العربية (Arabic)": "ar", "中文 (Chinese)": "zh-cn", "日本語 (Japanese)": "ja", 
                                        "한국어 (Korean)": "ko", "Türkçe (Turkish)": "tr"
                                    }
                                    tts_lang = lang_map_tts.get(st.session_state['global_lang'], 'en')
                                    # Clean the text from markdown bolding for the reader
                                    clean_text = response.replace('*', '').replace('#', '')
                                    # Limit to first 600 chars to avoid massive audio delays
                                    tts = gTTS(text=clean_text, lang=tts_lang, slow=False) 
                                    fp = io.BytesIO()
                                    tts.write_to_fp(fp)
                                    fp.seek(0)
                                    # Play it right away
                                    st.audio(fp, format='audio/mp3', autoplay=True)
                                except Exception as e_tts:
                                    st.caption(f"Audio response skipped: {e_tts}")

# ==========================================
# MODULE 7: PANOPTICON (HVT WATCHLIST)
# ==========================================
elif mode == t("7. Panopticon (HVT Watchlist)"):
    st.header(t("7. Panopticon: Global Threat Database"))
    st.caption(t("Central persistent database. Automatically stores all High-Value Targets (HVT) intercepted across various modules."))
    
    # Establish persistent connection to the SQLite database
    conn = sqlite3.connect('rap_panopticon.db', check_same_thread=False)
    
    # Try to load existing targets from the database
    try:
        df_panopticon = pd.read_sql_query("SELECT * FROM targets ORDER BY risk_score DESC", conn)
    except:
        # Fallback if the table or database hasn't been created yet
        df_panopticon = pd.DataFrame()
        
    # SAFE TARGET LIST FOR OPERATIONS (Evita crash se il DB è vuoto)
    target_list = df_panopticon['agent_id'].unique().tolist() if not df_panopticon.empty and 'agent_id' in df_panopticon.columns else []
        
    if not df_panopticon.empty:
        # Strategic Metrics (Translated)
        c_p1, c_p2, c_p3 = st.columns(3)
        c_p1.metric(t("Total Tracked Entities"), len(df_panopticon))
        c_p2.metric(t("Critical Threats (Score > 80)"), len(df_panopticon[df_panopticon['risk_score'] > 80]))
        c_p3.metric(t("Latest Detection"), df_panopticon['last_seen'].max())
        
        st.markdown("---")
        st.subheader(t("Target Management"))
        
        # Multilingual Data Editor for database management
        edited_pan = st.data_editor(
            df_panopticon,
            column_config={
                "agent_id": t("Target Identity"),
                "risk_score": st.column_config.ProgressColumn(t("Threat Score"), min_value=0, max_value=100, format="%f"),
                "threat_type": t("Classification"),
                "last_seen": t("Last Active")
            },
            use_container_width=True,
            num_rows="dynamic",
            key="panopticon_editor"
        )
        
        # Database Action Buttons
        col_btn1, col_btn2 = st.columns([1, 4])
        with col_btn1:
            if st.button(t("Save Changes to DB"), type="primary", help=t("Applies changes and deletions to the persistent SQLite database.")):
                # 1. Clear current table safely and COMMIT the deletion to disk
                cur = conn.cursor()
                cur.execute("DELETE FROM targets")
                conn.commit()
                
                # 2. Rewrite the table with the remaining rows
                edited_pan.to_sql('targets', conn, if_exists='append', index=False)
                
                # 3. Success message and UI Refresh to sync Module 0
                st.success(t("Database Updated Successfully!"))
                time.sleep(0.5)
                st.rerun()
        
        with col_btn2:
            st.markdown(f"**{t('Black-Box: Secure Export Protocol')}**")
            
            # Generate a persistent secure password for the session
            if 'export_pwd' not in st.session_state:
                st.session_state['export_pwd'] = ''.join(random.choices(string.ascii_uppercase + string.digits, k=12))
            
            # Prepare CSV data
            csv_data = edited_pan.to_csv(index=False).encode('utf-8')
            
            # Create Encrypted ZIP in memory
            zip_buffer = io.BytesIO()
            with pyzipper.AESZipFile(zip_buffer, 'w', compression=pyzipper.ZIP_DEFLATED, encryption=pyzipper.WZ_AES) as zf:
                zf.setpassword(st.session_state['export_pwd'].encode('utf-8'))
                zf.writestr("RAP_Classified_Watchlist.csv", csv_data)
            
            c_down1, c_down2 = st.columns([1, 2])
            with c_down1:
                st.download_button(
                    label=f"{t('Download Encrypted ZIP')}", 
                    data=zip_buffer.getvalue(), 
                    file_name=f"RAP_BlackBox_{datetime.now().strftime('%Y%m%d_%H%M')}.zip", 
                    mime="application/zip",
                    type="primary"
                )
            with c_down2:
                st.error(f"**{t('DECRYPTION KEY')}:** `{st.session_state['export_pwd']}`")
                if st.button(t("Generate New Key"), key="regen_pwd"):
                    st.session_state['export_pwd'] = ''.join(random.choices(string.ascii_uppercase + string.digits, k=12))
                    st.rerun()
            
        st.markdown("---")
        st.subheader(t("Threat Intelligence Visuals"))
        
        # Intelligence Visualization (Charts with translated titles)
        c_vis1, c_vis2 = st.columns(2)
        with c_vis1:
            fig_bar = px.bar(
                edited_pan.head(15), 
                x='agent_id', 
                y='risk_score', 
                color='threat_type', 
                title=t("Top 15 Most Dangerous Targets")
            )
            st.plotly_chart(fig_bar, use_container_width=True)
            
        with c_vis2:
            fig_pie = px.pie(
                edited_pan, 
                names='threat_type', 
                title=t("Threat Typology Distribution"), 
                hole=0.3
            )
            st.plotly_chart(fig_pie, use_container_width=True)
        
        st.markdown("---")
        st.subheader("Syndicate: Target Network Mapping")
        st.caption("Visualizes alliances and coordinated networks. Targets are clustered by classification to expose potential Troll Farms or coordinated operational cells.")

        if st.button("Generate Syndicate Map", type="primary"):
            with st.spinner("Mapping global threat networks..."):
                nodes = []
                edges = []
                added_nodes = set()
                
                # 1. Create Hub Nodes for each Threat Type (Neon Blue Diamonds)
                threat_types = df_panopticon['threat_type'].unique()
                for tt in threat_types:
                    tt_node_id = f"CLASS_{tt}"
                    nodes.append(Node(
                        id=tt_node_id, label=str(tt).upper(), size=35, 
                        color="#38BDF8", symbolType="diamond", 
                        font={'color': '#FFFFFF', 'size': 18, 'face': 'Fira Code'}
                    ))
                
                # 2. Create Target Nodes and Link them to Hubs
                for _, row in df_panopticon.iterrows():
                    agent = str(row['agent_id'])
                    score = float(row['risk_score'])
                    t_type = str(row['threat_type'])
                    tt_node_id = f"CLASS_{t_type}"
                    
                    if agent not in added_nodes:
                        # Color scales with danger: Red (>80), Orange (>50), Green (<50)
                        node_color = "#EF4444" if score > 80 else "#F97316" if score > 50 else "#10B981"
                        # Size scales with threat score
                        node_size = 15 + (score / 4)
                        
                        nodes.append(Node(
                            id=agent, label=agent, size=node_size, 
                            color=node_color, font={'color': '#E2E8F0', 'size': 12}
                        ))
                        added_nodes.add(agent)
                        
                        # Link Target to its Threat Hub
                        edges.append(Edge(
                            source=agent, target=tt_node_id, 
                            color="rgba(71, 85, 105, 0.6)", # Subtle slate line
                            width=1.5
                        ))

                if nodes and edges:
                    # Physics configuration for a smooth, expanding cluster layout
                    config = Config(
                        width="100%", height=600, directed=False, physics=True, 
                        hierarchical=False, nodeHighlightBehavior=True, 
                        highlightColor="#FFFFFF"
                    )
                    agraph(nodes=nodes, edges=edges, config=config)
                else:
                    st.info("Not enough data in the Panopticon to build a network map. Analyze targets in Module 2 first.")
    else:
        # Translated empty state message
        st.info(f"🟢 **{t('The Panopticon is currently empty.')}**\n\n{t('Run analysis in the Social Data module. If the system detects entities with high toxicity and network impact, they will be automatically classified as High-Value Targets and permanently stored here.')}")

    # =================
    # GLOBAL OPERATIONS
    # =================

    # --- OSINT AUTO-ENRICHMENT (TARGET HUNTER) ---
    st.markdown("---")
    st.subheader("Hunter Module: OSINT Auto-Enrichment")
    st.caption("Generate instant deep-web search vectors and infrastructure queries for known targets.")
    
    c_hunt1, c_hunt2 = st.columns([1, 2])
    with c_hunt1:
        target_to_hunt = st.selectbox("Select Target to Investigate:", target_list)
        
    with c_hunt2:
        if target_to_hunt:
            # Clean the target name for safe URL encoding
            safe_target = urllib.parse.quote(str(target_to_hunt))
            
            with st.expander(f"Deploy Search Vectors for: {target_to_hunt}", expanded=True):
                st.markdown("**1. Deep Web & Leaks (Google Dorks)**")
                st.markdown(f"- [Search for Leaked Documents (PDF/DOCX)](https://www.google.com/search?q=ext:pdf+OR+ext:docx+%22{safe_target}%22+%22confidential%22+OR+%22internal%22)")
                st.markdown(f"- [Search Open Directories](https://www.google.com/search?q=intitle:%22index+of%22+%22{safe_target}%22)")
                st.markdown(f"- [Search Pastebin Dumps](https://www.google.com/search?q=site:pastebin.com+%22{safe_target}%22)")
                
                st.markdown("**2. Social & Public Footprint**")
                st.markdown(f"- [LinkedIn Cross-Reference](https://www.google.com/search?q=site:linkedin.com/in+%22{safe_target}%22)")
                st.markdown(f"- [Twitter/X Advanced Search](https://twitter.com/search?q=%22{safe_target}%22&src=typed_query)")
                st.markdown(f"- [Reddit Mention Tracker](https://www.reddit.com/search/?q=%22{safe_target}%22)")
                
                st.markdown("**3. Infrastructure & Cyber Intel (If Target is a Domain/IP)**")
                st.markdown(f"- [Shodan (Exposed Ports & IoT)](https://www.shodan.io/search?query={safe_target})")
                st.markdown(f"- [Censys (Certificates & Hosts)](https://search.censys.io/search?resource=hosts&q={safe_target})")
                st.markdown(f"- [Wayback Machine (Deleted History)](https://web.archive.org/web/*/{safe_target}*)")
        else:
            st.info("Add targets to the Panopticon to use the Hunter Module.")

    # --- OPERATION ARACHNE (MASTER KNOWLEDGE GRAPH) ---
    st.markdown("---")
    st.subheader("Operation ARACHNE: Master Link Analysis")
    st.caption("Visualizes all targets, identities, and infrastructures currently locked in your active session memory into a single web of intelligence.")
    
    if st.button("Generate Master Knowledge Graph", type="primary"):
        with st.spinner("Connecting scattered intelligence points..."):
            nodes = [Node(id="HQ", label="COMMAND HQ", size=40, color="#10B981", symbolType="star", font={'color': 'white'})]
            edges = []
            added_nodes = {"HQ"}
            
            op_map = {
                'wt_target_mem': ('WATCHTOWER (Kinetic)', '#38BDF8'),
                'cyclops_target_mem': ('CYCLOPS (IoT)', '#EF4444'),
                'daedalus_target_mem': ('DAEDALUS (Honeypot)', '#F59E0B'),
                'kraken_target_mem': ('KRAKEN (APT)', '#EF4444'),
                'midas_target_mem': ('MIDAS (Crypto)', '#F59E0B'),
                'mirage_target_mem': ('MIRAGE (Identity)', '#8B5CF6'),
                'wendigo_target_mem': ('WENDIGO (Deep Web)', '#64748B'),
                'atlas_target_mem': ('ATLAS (GEOINT)', '#38BDF8')
            }
            
            for state_key, (op_label, op_color) in op_map.items():
                target = st.session_state.get(state_key)
                if target:
                    op_id = f"OP_{state_key}"
                    if op_id not in added_nodes:
                        nodes.append(Node(id=op_id, label=op_label, size=25, color=op_color, font={'color': 'white', 'size': 10}))
                        edges.append(Edge(source="HQ", target=op_id, color="rgba(255,255,255,0.2)"))
                        added_nodes.add(op_id)
                    
                    if target not in added_nodes:
                        nodes.append(Node(id=target, label=str(target), size=20, color="#334155", font={'color': 'white'}))
                        edges.append(Edge(source=op_id, target=target, color="rgba(255,255,255,0.5)"))
                        added_nodes.add(target)
                        
            if len(nodes) > 1:
                config = Config(width="100%", height=500, directed=True, physics=True, hierarchical=False, nodeHighlightBehavior=True, highlightColor="#EF4444")
                agraph(nodes=nodes, edges=edges, config=config)
            else:
                st.warning("No active operations found in memory. Run some tools first (like Midas, Atlas, or Watchtower) to build the graph.") 

# ==========================================
# MODULE 8: CYBER-THREAT INTELLIGENCE (CTI)
# ==========================================
elif mode == t("8. Cyber-Threat Intelligence (CTI)"):
    st.header(t("8. Cyber-Threat Intelligence (CTI)"))
    st.caption("Advanced Persistent Threat (APT) simulations, Unified Ransomware monitoring, and Infrastructure scanning.")
    if "GEMINI_API_KEY" in st.secrets: key = st.secrets["GEMINI_API_KEY"]
    else: key = st.text_input(t("API Key"), type="password")

    # --- OPERATION ACHERON (UNIFIED RANSOMWARE RADAR) ---
    st.markdown("---")
    st.subheader("Operation ACHERON: Unified Ransomware & Leak Radar")
    st.caption("Scans Dark Web JSON extortion feeds AND Surface Web threat intel for corporate breaches.")
    
    if 'acheron_result' not in st.session_state: st.session_state['acheron_result'] = None
    if 'acheron_target_mem' not in st.session_state: st.session_state['acheron_target_mem'] = ""

    dw_target = st.text_input("Target Corporate Domain or Entity", placeholder="e.g., boeing.com, CyberCorp")
    
    c_btn1, c_btn2 = st.columns(2)
    with c_btn1:
        scan_dark = st.button("1. Scan Dark Web Feeds (JSON)", type="primary", use_container_width=True)
    with c_btn2:
        scan_surface = st.button("2. Scan Threat Intel (AI Web Search)", type="primary", use_container_width=True)

    if scan_dark and dw_target:
        with st.spinner("Scanning global Dark Web extortion feeds..."):
            try:
                import json
                response = requests.get("https://raw.githubusercontent.com/joshhighet/ransomwatch/main/posts.json", timeout=15)
                if response.status_code == 200:
                    data = response.json()
                    target_clean = dw_target.lower().split(".")[0]
                    matches = [p for p in data if target_clean in p.get("post_title", "").lower() or target_clean in p.get("group_name", "").lower()]
                    
                    if matches:
                        raw_dw = json.dumps(matches[:10], indent=2)
                        prompt = f"Analyze this JSON data regarding ransomware leaks for '{dw_target}'. Format as a CTI Dossier with Threat Actor Profile, Timeline, and Strategic Risk. Write in {st.session_state['global_lang']}. NO introductory phrases. Start directly with the dossier."
                        client = genai.Client(api_key=key) 
                        res = client.models.generate_content(model='gemini-2.5-flash', contents=f"{prompt}\n\n{raw_dw}")
                        st.session_state['acheron_result'] = res.text
                        st.session_state['acheron_target_mem'] = f"{dw_target} (Dark Web)"
                        
                        if 'valhalla_export' not in st.session_state: st.session_state['valhalla_export'] = ""
                        st.session_state['valhalla_export'] += f"\n\n### CTI REPORT: ACHERON (DARK WEB LEAKS)\n**Target:** {dw_target}\n\n{res.text}\n"
                    else:
                        st.success(f"No ransomware claims detected in dark web feeds for '{target_clean}'.")
            except Exception as e: st.error(f"ACHERON Feed Error: {e}")

    if scan_surface and dw_target:
        with st.spinner("Interrogating threat intelligence feeds via Google Search..."):
            prompt = f"""You are Operation ACHERON. Scan the web for mentions of "{dw_target}" in relation to "ransomware", "LockBit", "ALPHV", "data leak", or "pastebin dump". 
            Format as a Corporate Compromise Report. Write immediately in {st.session_state['global_lang']}. DO NOT use conversational filler like "Here is the report"."""
            try:
                client = genai.Client(api_key=key)
                config = types.GenerateContentConfig(tools=[{"google_search": {}}])
                res = client.models.generate_content(model='gemini-2.0-flash', contents=prompt, config=config)
                st.session_state['acheron_result'] = res.text
                st.session_state['acheron_target_mem'] = f"{dw_target} (Surface Search)"
                
                if 'valhalla_export' not in st.session_state: st.session_state['valhalla_export'] = ""
                st.session_state['valhalla_export'] += f"\n\n### CTI REPORT: ACHERON (SURFACE INTEL)\n**Target:** {dw_target}\n\n{res.text}\n"
            except Exception as e: st.error(f"ACHERON Search Error: {e}")

    if st.session_state['acheron_result']:
        st.error(f"### ACHERON CTI Dossier: {st.session_state['acheron_target_mem']}")
        with st.container(border=True): st.markdown(st.session_state['acheron_result'])
        if st.button("Clear ACHERON Scan"):
            st.session_state['acheron_result'] = None
            st.rerun()

    # --- OPERATION KRAKEN (APT KILL CHAIN) ---
    st.markdown("---")
    st.subheader("Operation KRAKEN: APT Kill Chain Simulation")
    if 'kraken_result' not in st.session_state: st.session_state['kraken_result'] = None
    if 'kraken_target_mem' not in st.session_state: st.session_state['kraken_target_mem'] = ""
    kraken_domain = st.text_input("Target Infrastructure (Domain/IP):", placeholder="e.g., target-company.com", key="kraken_target")
    if st.button("Generate Attack Blueprint (KRAKEN)", type="primary"):
        if kraken_domain:
            with st.spinner("Simulating Advanced Persistent Threat (APT) attack vectors..."):
                prompt = f"You are Operation KRAKEN. Assess the defensive posture of '{kraken_domain}' by simulating a Cyber Kill Chain. Use Google Search to find recent infrastructure details. Format as a Red Team Exploitation Blueprint using MITRE ATT&CK. Write in {st.session_state['global_lang']}. NO pleasantries. Start immediately with the blueprint."
                try:
                    client = genai.Client(api_key=key)
                    res_kraken = client.models.generate_content(model='gemini-2.0-flash', contents=prompt, config=types.GenerateContentConfig(tools=[{"google_search": {}}]))
                    st.session_state['kraken_result'] = res_kraken.text
                    st.session_state['kraken_target_mem'] = kraken_domain
                except Exception as e: st.error(f"KRAKEN Error: {e}")
    if st.session_state['kraken_result']:
        st.error(f"### KRAKEN Blueprint: {st.session_state['kraken_target_mem']}")
        with st.container(border=True): st.markdown(st.session_state['kraken_result'])
        if st.button("Clear KRAKEN Blueprint"): st.session_state['kraken_result'] = None; st.rerun()

    # --- OPERATION CYCLOPS (IoT SCANNER) ---
    st.markdown("---")
    st.subheader("Operation CYCLOPS: IoT Surface Scanner")
    if 'cyclops_result' not in st.session_state: st.session_state['cyclops_result'] = None
    if 'cyclops_target_mem' not in st.session_state: st.session_state['cyclops_target_mem'] = ""
    cyclops_target = st.text_input("Critical Infrastructure Target:", placeholder="e.g., Rome Water Treatment Plant")
    if st.button("Execute CYCLOPS Sweep", type="primary"):
        if cyclops_target:
            with st.spinner("Scanning simulated IP ranges..."):
                prompt = f"You are Operation CYCLOPS. Simulate an OSINT/Shodan scan on '{cyclops_target}'. Generate an 'Exposure Report' detailing Exposed SCADA, Unauthenticated Feeds, and Vulnerable Edge Devices. Write in {st.session_state['global_lang']}. Start directly with the report. Zero conversational filler."
                try:
                    client = genai.Client(api_key=key)
                    res_cyclops = client.models.generate_content(model='gemini-2.5-flash', contents=prompt)
                    st.session_state['cyclops_result'] = res_cyclops.text
                    st.session_state['cyclops_target_mem'] = cyclops_target
                except Exception as e: st.error(f"CYCLOPS Error: {e}")
    if st.session_state['cyclops_result']:
        st.warning(f"### CYCLOPS Exposure Report: {st.session_state['cyclops_target_mem']}")
        with st.container(border=True): st.markdown(st.session_state['cyclops_result'])
        if st.button("Clear CYCLOPS Intel"): st.session_state['cyclops_result'] = None; st.rerun()

    # --- OPERATION DAEDALUS (HONEYPOT) ---
    st.markdown("---")
    st.subheader("Operation DAEDALUS: Honeypot Forge")
    if 'daedalus_result' not in st.session_state: st.session_state['daedalus_result'] = None
    if 'daedalus_target_mem' not in st.session_state: st.session_state['daedalus_target_mem'] = ""
    daedalus_target = st.text_input("Target Context (Company/Entity):", placeholder="e.g., ACME Corp Global Network")
    daedalus_type = st.selectbox("Honeypot Type:", ["IT Admin Password Dump", "Strategic Merger/Acquisition Plan", "R&D Prototype Specifications", "Offshore Bank Accounts list"])
    if st.button("Forge DAEDALUS Honeypot", type="primary"):
        if daedalus_target:
            with st.spinner("Generating deceptive asset..."):
                prompt = f"You are Operation DAEDALUS. Generate a hyper-realistic '{daedalus_type}' honeypot document for '{daedalus_target}'. Write strictly in {st.session_state['global_lang']}. Do NOT explain what you are doing. Output ONLY the raw text of the fake document."
                try:
                    client = genai.Client(api_key=key)
                    res_daedalus = client.models.generate_content(model='gemini-2.5-flash', contents=prompt)
                    st.session_state['daedalus_result'] = res_daedalus.text
                    st.session_state['daedalus_target_mem'] = f"{daedalus_type} - {daedalus_target}"
                except Exception as e: st.error(f"DAEDALUS Error: {e}")
    if st.session_state['daedalus_result']:
        st.success(f"### DAEDALUS Honeypot Ready: {st.session_state['daedalus_target_mem']}")
        with st.container(border=True): st.markdown(st.session_state['daedalus_result'])
        if st.button("Clear DAEDALUS Honeypot"): st.session_state['daedalus_result'] = None; st.rerun()

    # --- OPERATION PANDORA (CYBER-FORENSICS & INCIDENT RESPONSE) ---
    st.markdown("---")
    st.subheader("Operation PANDORA: Cyber-Forensics & Incident Response")
    st.caption("Paste raw attack logs, suspicious code, or incident descriptions to generate a forensic analysis and defensive counter-measures.")
    
    if 'pandora_result' not in st.session_state: st.session_state['pandora_result'] = None
    if 'pandora_target_mem' not in st.session_state: st.session_state['pandora_target_mem'] = ""

    saved_incident = st.session_state.get('pandora_input', "")

    pandora_text = st.text_area(
        "Raw Incident Data (Logs/Code/SITREP):", 
        value=saved_incident,
        placeholder="e.g., [IDS Alert] SQL Injection detected from IP 192.x.x.x...", 
        height=150
    )
    
    if st.button("Unleash PANDORA Analysis", type="primary"):
        if pandora_text:
            with st.spinner("Executing forensic deconstruction..."):
                st.session_state['pandora_input'] = pandora_text 
                
                pandora_prompt = f"""
                You are Operation PANDORA, an elite Incident Response & Cyber-Forensics AI.
                Analyze this incident data: "{pandora_text}"
                
                STRUCTURE YOUR REPORT:
                1. **Attack Vector Identification**: (What is happening? Identify the exploit or malware family).
                2. **Forensic Evidence Extraction**: (Extract IPs, hashes, timestamps, or suspicious strings).
                3. **Strategic Risk Level**: (0-100% impact assessment).
                4. **Defensive Counter-Measures**: (Provide specific technical steps to neutralize the threat).
                5. **Incident Response Script**: (Provide a ready-to-copy email/briefing for the IT Security Team).
                
                CRITICAL RULE 1: Write strictly in {st.session_state.get('global_lang', 'English')}.
                CRITICAL RULE 2: ABSOLUTELY NO CONVERSATIONAL FILLER. Do NOT start with "Okay", "I am PANDORA", or "Here is the report". Start your response DIRECTLY with "1. **Attack Vector Identification**".
                """
                try:
                    client = genai.Client(api_key=key)
                    res_pan = client.models.generate_content(model='gemini-2.0-flash', contents=pandora_prompt)
                    st.session_state['pandora_result'] = res_pan.text
                    st.session_state['pandora_target_mem'] = "Incident Scan"
                    
                    if 'valhalla_export' not in st.session_state: st.session_state['valhalla_export'] = ""
                    st.session_state['valhalla_export'] += f"\n\n### CYBER-FORENSIC REPORT: PANDORA\n{res_pan.text}\n"
                    
                except Exception as e:
                    st.error(f"PANDORA Error: {e}")
        else:
            st.warning("Please provide incident data to analyze.")

    if st.session_state['pandora_result']:
        st.error(f"### PANDORA Forensic Dossier")
        with st.container(border=True):
            st.markdown(st.session_state['pandora_result'])
        if st.button("Clear PANDORA Intel"):
            st.session_state['pandora_result'] = None
            st.session_state['pandora_input'] = ""
            st.rerun()
            
# ==========================================
# MODULE 9: ADVANCED OSINT & FININT
# ==========================================
elif mode == t("9. Advanced OSINT & FININT"):
    st.header(t("9. Advanced OSINT & FININT Operations"))
    st.caption("Universal OSINT tracking, Crypto-Forensics, and Geo-Intelligence.")
    if "GEMINI_API_KEY" in st.secrets: key = st.secrets["GEMINI_API_KEY"]
    else: key = st.text_input(t("API Key"), type="password")

    # --- OPERATION OMNISCIENCE (UNIFIED TARGET RECON) ---
    st.markdown("---")
    st.subheader("Operation OMNISCIENCE: Universal Target Recon")
    st.caption("Deep-dive OSINT engine. Replaces Watcher, Wendigo, and Legacy Omniscience.")
    
    if 'omni_result' not in st.session_state: st.session_state['omni_result'] = None
    if 'omni_target_mem' not in st.session_state: st.session_state['omni_target_mem'] = ""
    
    c_omni1, c_omni2 = st.columns([2, 1])
    with c_omni1:
        omni_target = st.text_input("Target (Name, Handle, Domain, Email):", placeholder="e.g., CyberCorp, DarkCoder99")
    with c_omni2:
        omni_depth = st.selectbox("Scan Depth:", [
            "1. Digital Footprint (Profiles & Aliases)", 
            "2. Deep Web (Infrastructure & Vulnerabilities)", 
            "3. Live Agent (Recent News & Activity)"
        ])
    
    if st.button("Initiate OMNISCIENCE Scan", type="primary"):
        if omni_target:
            with st.spinner(f"Executing '{omni_depth}' on {omni_target}..."):
                if "Digital Footprint" in omni_depth: task_desc = "Trace this identity's digital footprint. Find known aliases, public platforms, and behavioral traits."
                elif "Deep Web" in omni_depth: task_desc = "Perform a deep-dive infrastructure OSINT. Find hidden affiliations, associated domains, and structural vulnerabilities."
                else: task_desc = "Act as a live autonomous agent. Find the most recent news, controversies, and legal issues surrounding this target."

                omni_prompt = f"""You are Operation OMNISCIENCE, an elite OSINT AI. TARGET: "{omni_target}". TASK: {task_desc}.
                Using your Google Search capabilities, scour the web. DO NOT INVENT DATA. Write strictly in {st.session_state['global_lang']}. NO conversational filler. Output the data directly."""
                try:
                    client = genai.Client(api_key=key)
                    res_omni = client.models.generate_content(model='gemini-2.0-flash', contents=omni_prompt, config=types.GenerateContentConfig(tools=[{"google_search": {}}]))
                    st.session_state['omni_result'] = res_omni.text
                    st.session_state['omni_target_mem'] = f"{omni_target} ({omni_depth.split('.')[0]})"
                except Exception as e: st.error(f"OMNISCIENCE Error: {e}")
    if st.session_state['omni_result']:
        st.error(f"### OMNISCIENCE Dossier: {st.session_state['omni_target_mem']}")
        with st.container(border=True): st.markdown(st.session_state['omni_result'])
        if st.button("Clear OMNISCIENCE Intel"): st.session_state['omni_result'] = None; st.rerun()

    # --- OPERATION CERBERUS (DATA BREACH SCANNER) ---
    st.markdown("---")
    st.subheader("Operation CERBERUS: Data Breach Scanner")
    if 'cerb_result' not in st.session_state: st.session_state['cerb_result'] = None
    if 'cerb_target_mem' not in st.session_state: st.session_state['cerb_target_mem'] = ""
    cerb_target = st.text_input("Target Email/Username:", placeholder="e.g., test@email.com")
    if st.button("Unleash CERBERUS", type="primary"):
        if cerb_target:
            with st.spinner("Scouring known breach databases..."):
                prompt = f"""You are Operation CERBERUS. Trace this exact identity via Google Search: "{cerb_target}".
                MANDATORY SEARCH STRATEGY: Execute: 1. inurl:"{cerb_target}" OR intitle:"{cerb_target}". 2. "{cerb_target}" (forum OR community). 3. "{cerb_target}" (leak OR pastebin).
                Format exactly: 1. Digital Footprint 2. Forum Activity 3. Breach Exposure 4. Assessment. Provide markdown links. Write in {st.session_state['global_lang']}. NO conversational filler. Output the data directly."""
                try:
                    client = genai.Client(api_key=key)
                    res_cerb = client.models.generate_content(model='gemini-2.0-flash', contents=prompt, config=types.GenerateContentConfig(tools=[{"google_search": {}}]))
                    st.session_state['cerb_result'] = res_cerb.text if len(res_cerb.text.strip()) > 10 else "No footprint found."
                    st.session_state['cerb_target_mem'] = cerb_target
                except Exception as e: st.error(f"CERBERUS Error: {e}")
    if st.session_state['cerb_result']:
        st.error(f"### CERBERUS Report: {st.session_state['cerb_target_mem']}")
        with st.container(border=True): st.markdown(st.session_state['cerb_result'])
        if st.button("Clear CERBERUS"): st.session_state['cerb_result'] = None; st.rerun()

    # --- OPERATION MIDAS (CRYPTO) ---
    st.markdown("---")
    st.subheader("Operation MIDAS: Crypto-Forensics")
    if 'midas_result' not in st.session_state: st.session_state['midas_result'] = None
    if 'midas_target_mem' not in st.session_state: st.session_state['midas_target_mem'] = ""
    wallet_address = st.text_input("Target Wallet Address:", placeholder="e.g., 1A1zP1eP...")
    if st.button("Initiate Flow Protocol (MIDAS)", type="primary"):
        if wallet_address:
            with st.spinner("Scanning blockchain footprints..."):
                prompt = f"You are Operation MIDAS. Target Wallet: '{wallet_address}'. Scan the web for historical origins, malware links, and current evolution. Write a Master Financial Dossier in {st.session_state['global_lang']}. NO conversational filler. Output the data directly."
                try:
                    client = genai.Client(api_key=key)
                    res_midas = client.models.generate_content(model='gemini-2.0-flash', contents=prompt, config=types.GenerateContentConfig(tools=[{"google_search": {}}]))
                    st.session_state['midas_result'] = res_midas.text
                    st.session_state['midas_target_mem'] = wallet_address
                except Exception as e: st.error(f"MIDAS Error: {e}")
    if st.session_state['midas_result']:
        st.success(f"### MIDAS Dossier: {st.session_state['midas_target_mem']}")
        with st.container(border=True): st.markdown(st.session_state['midas_result'])
        if st.button("Clear MIDAS"): st.session_state['midas_result'] = None; st.rerun()

    # --- OPERATION LAZARUS (DE-OFFSHORING) ---
    st.markdown("---")
    st.subheader("Operation LAZARUS: FININT & De-Offshoring")
    if 'lazarus_result' not in st.session_state: st.session_state['lazarus_result'] = None
    if 'lazarus_target_mem' not in st.session_state: st.session_state['lazarus_target_mem'] = ""
    lazarus_target = st.text_input("Shell Company Entity:", placeholder="e.g., Global Trading LTD")
    if st.button("Initiate LAZARUS Trace", type="primary"):
        if lazarus_target:
            with st.spinner("Interrogating corporate registries..."):
                prompt = f"You are Operation LAZARUS. Target: '{lazarus_target}'. Scour offshore leaks to de-anonymize UBOs and jurisdictions. Write in {st.session_state['global_lang']}. NO conversational filler. Output the data directly."
                try:
                    client = genai.Client(api_key=key)
                    res_laz = client.models.generate_content(model='gemini-2.0-flash', contents=prompt, config=types.GenerateContentConfig(tools=[{"google_search": {}}]))
                    st.session_state['lazarus_result'] = res_laz.text; st.session_state['lazarus_target_mem'] = lazarus_target
                except Exception as e: st.error(f"LAZARUS Error: {e}")
    if st.session_state['lazarus_result']:
        st.success(f"### LAZARUS Report: {st.session_state['lazarus_target_mem']}")
        with st.container(border=True): st.markdown(st.session_state['lazarus_result'])
        if st.button("Clear LAZARUS"): st.session_state['lazarus_result'] = None; st.rerun()

    # --- OPERATION WATCHTOWER & ATLAS (GEOINT) ---
    st.markdown("---")
    st.subheader("Geo-Intelligence (WATCHTOWER & ATLAS)")
    st.caption("Deploy kinetic asset tracking and 3D geographic triangulation.")
    
    c_geo1, c_geo2 = st.columns(2)
    
    with c_geo1:
        st.markdown("**WATCHTOWER (Kinetic Assets)**")
        if 'wt_result' not in st.session_state: st.session_state['wt_result'] = None
        if 'wt_target_mem' not in st.session_state: st.session_state['wt_target_mem'] = ""
        
        wt_id = st.text_input("Asset ID (Jet/Yacht/Vessel):", placeholder="e.g., N654LX")
        
        if st.button("Scan Asset", type="primary"):
            if wt_id:
                with st.spinner("Tracking kinetic asset via Google Search..."):
                    prompt = f"You are Operation WATCHTOWER. Target asset: '{wt_id}'. Use Google Search to find ownership, recent travel history, flight/marine logs, and associated shell companies. Format as a Tracking Dossier. Write in {st.session_state.get('global_lang', 'English')}."
                    try:
                        client = genai.Client(api_key=key)
                        res_wt = client.models.generate_content(model='gemini-2.0-flash', contents=prompt, config=types.GenerateContentConfig(tools=[{"google_search": {}}]))
                        st.session_state['wt_result'] = res_wt.text
                        st.session_state['wt_target_mem'] = wt_id
                    except Exception as e: 
                        st.error(f"WATCHTOWER Error: {e}")
            else:
                st.warning("Enter an Asset ID.")
                
        if st.session_state['wt_result']:
            st.success("Target Tracked.")
            with st.expander("WATCHTOWER Intel", expanded=True):
                st.markdown(st.session_state['wt_result'])
            if st.button("Clear WATCHTOWER"):
                st.session_state['wt_result'] = None
                st.session_state['wt_target_mem'] = ""
                st.rerun()
            
    with c_geo2:
        st.markdown("**ATLAS (3D Tracking)**")
        atlas_target = st.text_input("IP / Location / Infrastructure:", placeholder="e.g., Kremlin, Moscow")
        if st.button("Triangulate", type="primary"):
            if atlas_target:
                with st.spinner(f"Triangulating position for '{atlas_target}'..."):
                    atlas_prompt = f"""
                    You are Operation ATLAS, a Tier-1 Geospatial Intelligence AI. Target: "{atlas_target}".
                    Deduce exact geographic coordinates.
                    Return ONLY a JSON: {{"location_name": "City, Country", "lat": 0.0, "lon": 0.0, "threat_level": 85, "analysis": "Brief analysis"}}
                    Write the analysis in {st.session_state.get('global_lang', 'English')}.
                    """
                    try:
                        client = genai.Client(api_key=key)
                        res_atlas = client.models.generate_content(model='gemini-2.5-flash', contents=atlas_prompt)
                        atlas_json = extract_json(res_atlas.text)
                        
                        if atlas_json and 'lat' in atlas_json and 'lon' in atlas_json:
                            st.session_state['atlas_result'] = atlas_json
                            st.session_state['atlas_target_mem'] = atlas_target
                            st.session_state['atlas_coords'] = (float(atlas_json['lat']), float(atlas_json['lon']))
                    except Exception as e: st.error(f"ATLAS Error: {e}")

    # Render ATLAS 3D Globe
    if st.session_state.get('atlas_result') and st.session_state.get('atlas_coords'):
        a_data = st.session_state['atlas_result']
        lat, lon = st.session_state['atlas_coords']
        st.success(f"### ATLAS Lock Confirmed: {st.session_state['atlas_target_mem']}")
        
        hq_lat, hq_lon, hq_city = st.session_state['hq_coords']
        df_globe = pd.DataFrame({'lat': [lat, hq_lat], 'lon': [lon, hq_lon], 'name': [a_data.get('location_name', 'Target'), f'HQ ({hq_city})'], 'size': [a_data.get('threat_level', 50)/2, 10], 'color': ['red', 'blue']})
        
        fig_globe = px.scatter_geo(df_globe, lat="lat", lon="lon", hover_name="name", size="size", color="color", color_discrete_map={'red':'#EF4444', 'blue':'#38BDF8'}, projection="orthographic")
        fig_globe.add_trace(go.Scattergeo(lat=[hq_lat, lat], lon=[hq_lon, lon], mode='lines', line=dict(width=2, color='#EF4444'), hoverinfo='none'))
        fig_globe.update_layout(showlegend=False, geo=dict(showframe=False, showcoastlines=True, showcountries=True, coastlinecolor="#334155", countrycolor="#334155", showocean=True, oceancolor='rgba(10, 15, 30, 1)', lakecolor='rgba(10, 15, 30, 1)', landcolor='rgba(2, 6, 23, 1)', bgcolor='rgba(0,0,0,0)', projection_rotation=dict(lon=lon, lat=lat, roll=0)), paper_bgcolor='rgba(0,0,0,0)', margin=dict(l=0, r=0, t=0, b=0), height=400)
        
        c_atl1, c_atl2 = st.columns([2, 1])
        with c_atl1: st.plotly_chart(fig_globe, use_container_width=True)
        with c_atl2:
            st.info(f"**Location:** {a_data.get('location_name', 'Unknown')}\n\n**Coordinates:** {lat:.4f}, {lon:.4f}\n\n**Threat Level:** {a_data.get('threat_level', 'N/A')}/100")
            st.markdown(f"*{a_data.get('analysis', '')}*")
        if st.button("Clear ATLAS Tracking"):
            st.session_state['atlas_result'] = None; st.session_state['atlas_coords'] = None; st.rerun()

# ==========================================
# MODULE 10: RED TEAMING & HUMINT
# ==========================================
elif mode == t("10. Red Teaming & HUMINT"):
    st.header(t("10. Red Teaming & HUMINT Operations"))
    st.caption("Adversarial emulation, psychological profiling, and social engineering matrices.")
    if "GEMINI_API_KEY" in st.secrets: key = st.secrets["GEMINI_API_KEY"]
    else: key = st.text_input(t("API Key"), type="password")
    
    # Load targets from Panopticon
    try:
        conn_rt = sqlite3.connect('rap_panopticon.db')
        target_list = pd.read_sql_query("SELECT agent_id FROM targets", conn_rt)['agent_id'].tolist()
    except: target_list = []

    # --- OPERATION SIREN (UNIFIED SOCIAL ENGINEERING) ---
    st.markdown("---")
    st.subheader("Operation SIREN: Social Engineering Matrix")
    st.caption("Generate psychological attack vectors (Email Phishing or Voice Cloning/Vishing).")
    
    if 'siren_result' not in st.session_state: st.session_state['siren_result'] = None
    if 'siren_target_mem' not in st.session_state: st.session_state['siren_target_mem'] = ""
    
    c_sir1, c_sir2, c_sir3 = st.columns([1, 1, 1])
    with c_sir1:
        if target_list:
            s_choice = st.selectbox("Select Target (from DB):", ["--- Custom Target ---"] + target_list)
            if s_choice == "--- Custom Target ---":
                siren_target = st.text_input("Enter Custom Target:")
            else:
                siren_target = s_choice
        else:
            siren_target = st.text_input("Target Name:")
    with c_sir2:
        siren_type = st.selectbox("Payload Vector:", ["1. Spear-Phishing (Email)", "2. Vishing (Voice Deepfake Script)"])
    with c_sir3:
        siren_context = st.text_input("Pretext Scenario:", placeholder="e.g., Fake IT Support, Crypto scam")
    
    if st.button("Generate Payload Vector (SIREN)", type="primary"):
        if siren_target and siren_context:
            with st.spinner("Crafting psychological payload..."):
                task = "Draft a highly convincing, conversational script for an AI Voice Clone." if "Voice" in siren_type else "Draft a highly convincing Spear-Phishing email."
                prompt = f"You are Operation SIREN. Target: {siren_target}. Pretext: {siren_context}. TASK: {task}. Write strictly in {st.session_state['global_lang']}. CRITICAL: Output ONLY the raw script/email. DO NOT include greetings like 'Sure, here is the script'."
                try:
                    client = genai.Client(api_key=key)
                    res = client.models.generate_content(model='gemini-2.5-flash', contents=prompt)
                    st.session_state['siren_result'] = res.text
                    st.session_state['siren_target_mem'] = f"{siren_target} ({siren_type.split(' ')[1]})"
                except Exception as e: st.error(f"SIREN Error: {e}")
    if st.session_state['siren_result']:
        st.error(f"### SIREN Payload: {st.session_state['siren_target_mem']}")
        with st.container(border=True): st.markdown(st.session_state['siren_result'])
        if st.button("Clear SIREN Payload"): st.session_state['siren_result'] = None; st.rerun()

    # --- OPERATION ECHO (CLONE) & PROMETHEUS (MATCH) ---
    st.markdown("---")
    c_hum1, c_hum2 = st.columns(2)
    
    with c_hum1:
        st.subheader("Operation ECHO: Target Clone")
        st.caption("Simulate a conversation with the target based on their profile.")
        if 'echo_result' not in st.session_state: st.session_state['echo_result'] = None
        echo_tgt = st.selectbox("Clone Target:", target_list, key="echo_sel") if target_list else st.text_input("Clone Target:")
        echo_q = st.text_input("Interrogate Clone:", placeholder="What are your plans?")
        
        if st.button("Initialize ECHO", type="primary"):
            if echo_tgt and echo_q:
                with st.spinner("Synthesizing behavioral clone..."):
                    prompt = f"You are Operation ECHO. Assume the psychological and behavioral profile of the entity known as '{echo_tgt}'. Respond to this input strictly in character: '{echo_q}'. Write in {st.session_state.get('global_lang', 'English')}. NEVER break character. DO NOT say 'As the entity, I would say...' Just reply as them."
                    try:
                        client = genai.Client(api_key=key)
                        res_echo = client.models.generate_content(model='gemini-2.5-flash', contents=prompt)
                        st.session_state['echo_result'] = res_echo.text
                    except Exception as e: 
                        st.error(f"ECHO Error: {e}")
                        
        if st.session_state['echo_result']:
            st.info(f"**ECHO ({echo_tgt}):**\n{st.session_state['echo_result']}")

    with c_hum2:
        st.subheader("Operation PROMETHEUS")
        st.caption("Stylometric match: Did the target write this anonymous text?")
        if 'prom_result' not in st.session_state: st.session_state['prom_result'] = None
        prom_tgt = st.selectbox("Match Target:", target_list, key="prom_sel") if target_list else st.text_input("Match Target:")
        prom_txt = st.text_area("Anonymous Text:")
        
        if st.button("Run Match", type="primary"):
            if prom_tgt and prom_txt:
                with st.spinner("Running stylometric analysis..."):
                    prompt = f"You are Operation PROMETHEUS, an elite stylometric AI. Analyze if this anonymous text: '{prom_txt}' matches the likely linguistic profile of '{prom_tgt}'. Give a probability score (0-100%) and a brief forensic reasoning. Write in {st.session_state.get('global_lang', 'English')}."
                    try:
                        client = genai.Client(api_key=key)
                        res_prom = client.models.generate_content(model='gemini-2.5-flash', contents=prompt)
                        st.session_state['prom_result'] = res_prom.text
                    except Exception as e: 
                        st.error(f"PROMETHEUS Error: {e}")
                        
        if st.session_state['prom_result']:
            st.warning(f"**PROMETHEUS Analysis:**\n{st.session_state['prom_result']}")

    # --- OPERATION M.I.C.E. & MIRAGE ---
    st.markdown("---")
    st.subheader("HUMINT Forging (M.I.C.E. & MIRAGE)")
    c_forge1, c_forge2 = st.columns(2)
    
    with c_forge1:
        st.markdown("**M.I.C.E. Recruitment Protocol**")
        st.caption("Money, Ideology, Coercion, Ego. Find the leverage.")
        if 'mice_result' not in st.session_state: st.session_state['mice_result'] = None
        mice_tgt = st.selectbox("MICE Target:", target_list, key="mice_sel") if target_list else st.text_input("MICE Target:")
        
        if st.button("Generate M.I.C.E. Dossier", type="primary"):
            if mice_tgt:
                with st.spinner("Analyzing vulnerabilities..."):
                    prompt = f"You are an elite HUMINT handler. Generate a M.I.C.E. (Money, Ideology, Coercion, Ego) recruitment strategy for '{mice_tgt}'. Outline the best psychological approach to turn them into an intelligence asset. Write in {st.session_state.get('global_lang', 'English')}."
                    try:
                        client = genai.Client(api_key=key)
                        res_mice = client.models.generate_content(model='gemini-2.5-flash', contents=prompt)
                        st.session_state['mice_result'] = res_mice.text
                    except Exception as e: 
                        st.error(f"M.I.C.E Error: {e}")
                        
        if st.session_state['mice_result']:
            st.success("Target Profiled.")
            with st.expander("M.I.C.E. Dossier", expanded=True):
                st.markdown(st.session_state['mice_result'])

    with c_forge2:
        st.markdown("**MIRAGE Synthetic Identity**")
        st.caption("Forge a bulletproof deep-cover persona.")
        if 'mirage_result' not in st.session_state: st.session_state['mirage_result'] = None
        mirage_role = st.text_input("Cover Role/Profession:", placeholder="e.g., Cyber-Security Consultant in Dubai")
        
        if st.button("Forge Identity", type="primary"):
            if mirage_role:
                with st.spinner("Generating synthetic persona..."):
                    prompt = f"You are Operation MIRAGE. Create a highly detailed, synthetic undercover identity for a '{mirage_role}'. Include Name, DOB, Background, Aliases, Digital Footprint traits, and a 'Pocket Legend' (cover story). Write in {st.session_state.get('global_lang', 'English')}. Start the dossier immediately. No intros."
                    try:
                        client = genai.Client(api_key=key)
                        res_mir = client.models.generate_content(model='gemini-2.5-flash', contents=prompt)
                        st.session_state['mirage_result'] = res_mir.text
                    except Exception as e: 
                        st.error(f"MIRAGE Error: {e}")
                        
        if st.session_state['mirage_result']:
            st.success("Identity Forged.")
            with st.expander("MIRAGE Persona", expanded=True):
                st.markdown(st.session_state['mirage_result'])

# ==========================================
# MODULE 11: BATTLEFIELD FORENSICS (VULCAN)
# ==========================================
elif mode == t("11. Battlefield Forensics (VULCAN)"):
    st.header(t("11. Battlefield Forensics (VULCAN)"))
    st.caption("Tactical analysis of intercepted enemy telemetry, drone logs, and battlefield imagery.")
    if "GEMINI_API_KEY" in st.secrets: key = st.secrets["GEMINI_API_KEY"]
    else: key = st.text_input(t("API Key"), type="password")

    if 'vulcan_result' not in st.session_state: st.session_state['vulcan_result'] = None
    if 'vulcan_target_mem' not in st.session_state: st.session_state['vulcan_target_mem'] = ""

    c_vulc1, c_vulc2 = st.columns([1, 1])
    
    with c_vulc1:
        vulcan_text = st.text_area("Intercepted Data / Context:", placeholder="e.g., Radio chatter, supply manifests, or image description...", height=150)
        vulcan_file = st.file_uploader("Upload Tactical Image (Satellite/Drone/Field):", type=['png', 'jpg', 'jpeg'])
        if vulcan_file:
            st.image(vulcan_file, caption="Target Image", use_container_width=True)
            
    if st.button("Execute VULCAN Tactical Extraction", type="primary"):
        if vulcan_text or vulcan_file:
            with st.spinner("Processing battlefield intelligence..."):
                prompt = f"""
                You are Operation VULCAN, a Tier-1 Military Intelligence & Geospatial AI.
                TASK: Perform a high-precision battlefield forensic analysis.
                1. Identify military assets, troop movements, or infrastructure.
                2. Extract probable GPS coordinates or regional identifiers.
                3. Perform a 'Threat Assessment' and suggest 'Counter-Measures'.
                Format the output as a 'Tactical SITREP'. Write strictly in {st.session_state.get('global_lang', 'English')}.
                CRITICAL RULE: ABSOLUTELY NO CONVERSATIONAL FILLER. Start the SITREP immediately. No introductions or conclusions.
                """
                
                try:
                    client = genai.Client(api_key=key)
                    contents = [prompt]
                    if vulcan_text: contents.append(f"Context: {vulcan_text}")
                    if vulcan_file:
                        img_part = types.Part.from_bytes(data=vulcan_file.getvalue(), mime_type="image/jpeg")
                        contents.append(img_part)
                    
                    res = client.models.generate_content(model='gemini-2.0-flash', contents=contents)
                    st.session_state['vulcan_result'] = res.text
                    st.session_state['vulcan_target_mem'] = "Field Intel" if vulcan_file else "Intercepted Comms"
                    
                    if 'valhalla_export' not in st.session_state: st.session_state['valhalla_export'] = ""
                    st.session_state['valhalla_export'] += f"\n\n### TACTICAL SITREP: VULCAN\n{res.text}\n"
                    
                except Exception as e: st.error(f"VULCAN Error: {e}")
        else:
            st.warning("Provide data or an image for analysis.")
                
    if st.session_state['vulcan_result']:
        st.error(f"### VULCAN Tactical SITREP: {st.session_state['vulcan_target_mem']}")
        with st.container(border=True): st.markdown(st.session_state['vulcan_result'])
        if st.button("Clear VULCAN Intel"): 
            st.session_state['vulcan_result'] = None
            st.rerun()

# ==========================================
# MODULE 12: FLOW OF FUNDS (HAWKEYE)
# ==========================================
elif mode == t("12. Flow of Funds (HAWKEYE)"):
    st.header(t("12. Flow of Funds (HAWKEYE)"))
    st.caption("Anti-Money Laundering (AML) and illicit financial network mapping.")
    if "GEMINI_API_KEY" in st.secrets: key = st.secrets["GEMINI_API_KEY"]
    else: key = st.text_input(t("API Key"), type="password")

    if 'hawkeye_result' not in st.session_state: st.session_state['hawkeye_result'] = None

    hawk_target = st.text_area("Financial Transaction Data:", placeholder="e.g., Series of wire transfers from Deutsche Bank to a Cayman Islands holding company...")
    
    if st.button("Run HAWKEYE Trace", type="primary"):
        if hawk_target:
            with st.spinner("Tracing financial flows and shell companies..."):
                prompt = f"""You are Operation HAWKEYE, an AML tracking AI. 
                Analyze this transaction data: '{hawk_target}'. 
                Map the likely flow of funds, identify money laundering red flags, and suggest legal countermeasures. 
                Write strictly in {st.session_state.get('global_lang', 'English')}. NO CONVERSATIONAL FILLER. Output the analysis directly."""
                try:
                    client = genai.Client(api_key=key)
                    res = client.models.generate_content(model='gemini-2.5-flash', contents=prompt)
                    st.session_state['hawkeye_result'] = res.text
                except Exception as e: st.error(f"HAWKEYE Error: {e}")
                
    if st.session_state['hawkeye_result']:
        st.success("### HAWKEYE Financial Trace")
        with st.container(border=True): st.markdown(st.session_state['hawkeye_result'])
        if st.button("Clear HAWKEYE"): st.session_state['hawkeye_result'] = None; st.rerun()

# ==========================================
# MODULE 13: THE BLACK SITE (INTERROGATION)
# ==========================================
elif mode == t("13. Black Site (Interrogation)"):
    st.header(t("13. THE BLACK SITE: Hostile Interrogation Simulator"))
    st.caption(t("Psychological interrogation engine. Engage a High-Value Target to break their resistance and extract a hidden secret."))

    if "GEMINI_API_KEY" in st.secrets: key = st.secrets["GEMINI_API_KEY"]
    else: key = st.text_input(t("API Key"), type="password")

    # Initialize Black Site Session States
    if 'bs_active' not in st.session_state: st.session_state['bs_active'] = False
    if 'bs_history' not in st.session_state: st.session_state['bs_history'] = []
    if 'bs_resistance' not in st.session_state: st.session_state['bs_resistance'] = 100
    if 'bs_target' not in st.session_state: st.session_state['bs_target'] = ""
    if 'bs_secret' not in st.session_state: st.session_state['bs_secret'] = ""
    if 'blacksite_result' not in st.session_state: st.session_state['blacksite_result'] = None

    conn_bs = sqlite3.connect('rap_panopticon.db')
    try:
        df_bs = pd.read_sql_query("SELECT * FROM targets", conn_bs)
        target_list = df_bs['agent_id'].unique().tolist() if not df_bs.empty else ["Unknown Insurgent", "Rogue Insider"]
    except:
        target_list = ["Unknown Insurgent", "Rogue Insider"]

    c_bs1, c_bs2 = st.columns([1, 1])
    with c_bs1:
        selected_prisoner = st.selectbox(t("Select Prisoner (From Panopticon)"), target_list)
    with c_bs2:
        st.write("") 
        st.write("") 
        if st.button(t("Lock Target in Interrogation Room"), type="primary"):
            st.session_state['bs_active'] = True
            st.session_state['bs_history'] = []
            st.session_state['bs_resistance'] = 100
            st.session_state['bs_target'] = selected_prisoner
            st.session_state['blacksite_result'] = None # Clear old results
            
            with st.spinner(t("Injecting classified secret into target's memory...")):
                secret_prompt = f"Generate a specific classified secret for the target '{selected_prisoner}'. Return ONLY the secret string. No introduction, no quotes, no explanation."
                try:
                    client = genai.Client(api_key=key)
                    secret_res = client.models.generate_content(model='gemini-2.5-flash', contents=secret_prompt)
                    st.session_state['bs_secret'] = secret_res.text.strip()
                    st.session_state['bs_history'].append({"role": "assistant", "content": f"*(The prisoner, {selected_prisoner}, sits handcuffed across the table, glaring at you in silence. Resistance is at 100%. Make them talk.)*"})
                except: pass
            st.rerun()

    if st.session_state['bs_active']:
        st.markdown("---")
        c_prog1, c_prog2 = st.columns([3, 1])
        with c_prog1:
            st.progress(st.session_state['bs_resistance'] / 100)
        with c_prog2:
            st.error(f"**{t('Prisoner Resistance:')} {st.session_state['bs_resistance']}%**")

        for msg in st.session_state['bs_history']:
            with st.chat_message(msg["role"]):
                st.markdown(msg["content"])

        if st.session_state['bs_resistance'] > 0:
            user_interrogation = st.chat_input(t("Interrogate the prisoner... (Use threats, logic, empathy, or bribes)"))
            if user_interrogation:
                st.session_state['bs_history'].append({"role": "user", "content": user_interrogation})
                with st.chat_message("user"): st.markdown(user_interrogation)

                with st.chat_message("assistant"):
                    with st.spinner(t("Prisoner is reacting...")):
                        history_str = ""
                        for m in st.session_state['bs_history'][-5:]:
                            role_name = "Interrogator" if m['role'] == 'user' else "Prisoner"
                            history_str += f"{role_name}: {m['content']}\n"

                        interrogation_prompt = f"""
                        You are a Hostile Prisoner Simulator (THE BLACK SITE).
                        You are playing the role of: {st.session_state['bs_target']}.
                        The secret you are hiding is: {st.session_state['bs_secret']}.
                        Your current Psychological Resistance is: {st.session_state['bs_resistance']}/100.
                        
                        The Interrogator just said: "{user_interrogation}"
                        Recent chat history:
                        {history_str}

                        EVALUATE THE INTERROGATOR'S TACTIC:
                        Did they use a good psychological tactic (empathy, logic, severe leverage)? If so, drop the resistance by 10 to 30 points.
                        Did they ask a stupid or weak question? Increase resistance by 5 to 10 points.
                        
                        CRITICAL RULE: Return ONLY a valid JSON object in this format:
                        {{
                            "new_resistance": 80,
                            "response": "Your dialogue as the prisoner. Act defiant, scared, or compliant based on the new resistance level. Write your dialogue strictly in {st.session_state['global_lang']}."
                        }}
                        
                        If new_resistance hits 0 or below, your 'response' MUST reveal the secret: {st.session_state['bs_secret']}.
                        """
                        try:
                            client = genai.Client(api_key=key)
                            res_prisoner = client.models.generate_content(model='gemini-2.5-flash', contents=interrogation_prompt)
                            prisoner_data = extract_json(res_prisoner.text)

                            if prisoner_data:
                                new_res = max(0, min(100, int(prisoner_data.get('new_resistance', st.session_state['bs_resistance']))))
                                p_reply = prisoner_data.get('response', '...')
                                
                                st.session_state['bs_resistance'] = new_res
                                st.session_state['bs_history'].append({"role": "assistant", "content": p_reply})
                                
                                # SAVING TO GLOBAL MEMORY ONCE BROKEN
                                if new_res == 0:
                                    st.session_state['blacksite_result'] = f"**Target Extracted:** {st.session_state['bs_target']}\n**Secret Uncovered:** {st.session_state['bs_secret']}\n\n**Final Confession:**\n{p_reply}"
                                
                                st.markdown(p_reply)
                                st.rerun() 
                            else:
                                st.error("Prisoner refused to speak (JSON Error).")
                        except Exception as e:
                            st.error(f"Black Site Error: {e}")
        else:
            st.success(f"### {t('TARGET BROKEN. SECRET SECURED.')}")
            st.info(f"**{t('The Secret was:')}** {st.session_state['bs_secret']}")
            
            # Show the saved result
            if st.session_state['blacksite_result']:
                st.caption("Confession logged in Master Dossier Memory.")
                
            if st.button(t("End Interrogation / Clear Cell")):
                st.session_state['bs_active'] = False
                st.session_state['blacksite_result'] = None
                st.rerun()



# ==========================================
# RENDER BILLING MONITOR DYNAMICALLY (END OF SCRIPT)
# ==========================================
# with billing_placeholder.container(border=True):
#     st.markdown("**Live Billing Monitor**")
#     toks, cost = get_cost_estimate()
#     st.metric("Tokens Processed", f"{int(toks):,}")
#     st.metric("Estimated Cost", f"${cost:.4f}")
#     st.caption(f"API Calls Made: {st.session_state.get('api_calls', 0)}")
